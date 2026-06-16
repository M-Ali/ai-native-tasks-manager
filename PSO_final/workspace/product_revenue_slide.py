"""
PSO Retail Fuels — Product Revenue & Volume Performance Slide
Sample data: PMG, HSD, R95, Lubes | National + Regional breakdown
Output: reports/PSO_Product_Revenue.pptx
"""
import sys, os
sys.path.insert(0, 'src')
os.environ['PYTHONIOENCODING'] = 'utf-8'

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml

NAVY  = RGBColor(0x1B,0x2A,0x4A); GOLD  = RGBColor(0xC9,0xA0,0x30)
WHITE = RGBColor(0xFF,0xFF,0xFF); LGREY = RGBColor(0xF2,0xF4,0xF8)
MGREY = RGBColor(0xD9,0xDC,0xE3); DBLUE = RGBColor(0x2E,0x5B,0x9A)
LBLUE = RGBColor(0xBD,0xD7,0xEE); GREEN = RGBColor(0x37,0x5B,0x25)
LGREEN= RGBColor(0xC6,0xEF,0xCE); RED   = RGBColor(0xC0,0x00,0x00)
LRED  = RGBColor(0xFF,0xC7,0xCE); LYELL = RGBColor(0xFF,0xEB,0x9C)
ORANGE= RGBColor(0xFF,0x82,0x00); LORANG= RGBColor(0xFF,0xE0,0xB2)

SW, SH = Inches(13.33), Inches(7.50)

PROD_COL = {
    "PMG":   DBLUE,
    "HSD":   GREEN,
    "R95":   RGBColor(0x8B,0x44,0x13),
    "Lubes": ORANGE,
}
PROD_LIGHT = {
    "PMG":   LBLUE,
    "HSD":   LGREEN,
    "R95":   LYELL,
    "Lubes": LORANG,
}

# ═══════════════════════════════════════════════════════════════════════════════
# DATA
# ═══════════════════════════════════════════════════════════════════════════════
PRODUCTS = ["PMG", "HSD", "R95", "Lubes"]
PROD_FULL = {
    "PMG":   "Motor Gasoline",
    "HSD":   "High Speed Diesel",
    "R95":   "Premium Euro-5",
    "Lubes": "Lubricants",
}

# National sample data (provided)
NAT = {
    "PMG":   {"rev_cy": 237.1, "rev_ly": 197.5, "rev_chg": 20.00, "vol_cy": 769.4},
    "HSD":   {"rev_cy": 160.8, "rev_ly": 133.4, "rev_chg": 20.50, "vol_cy": 505.1},
    "R95":   {"rev_cy": 10.2,  "rev_ly": 8.0,   "rev_chg": 27.10, "vol_cy": 32.1},
    "Lubes": {"rev_cy": 3.2,   "rev_ly": 2.8,   "rev_chg": 14.50, "vol_cy": 3.6},
}

# Derived national metrics
tot_rev_cy = sum(NAT[p]["rev_cy"] for p in PRODUCTS)
tot_rev_ly = sum(NAT[p]["rev_ly"] for p in PRODUCTS)
tot_vol_cy = sum(NAT[p]["vol_cy"] for p in PRODUCTS)
ovr_rev_chg = (tot_rev_cy - tot_rev_ly) / tot_rev_ly * 100

for p in PRODUCTS:
    d = NAT[p]
    d["rev_sh"] = d["rev_cy"] / tot_rev_cy * 100
    d["vol_sh"] = d["vol_cy"] / tot_vol_cy * 100
    d["rpl"]    = d["rev_cy"] * 1e9 / (d["vol_cy"] * 1e6)   # PKR per litre
    d["vol_ly"] = d["vol_cy"] / (1 + d["rev_chg"] / 100)    # estimated (assumes price flat)
    d["vol_chg"]= d["rev_chg"]                              # indicative

tot_vol_ly = sum(NAT[p]["vol_ly"] for p in PRODUCTS)
ovr_vol_chg = (tot_vol_cy - tot_vol_ly) / tot_vol_ly * 100

# Regional breakdown (indicative sample — proportional splits of national)
# South = Karachi+Sindh (biggest consumer market), North = KPK+N.Punjab, Central = Punjab
REG_SHARE = {
    #              PMG    HSD    R95    Lubes
    "South":   (0.42,  0.32,  0.36,  0.34),
    "Central": (0.28,  0.34,  0.26,  0.30),
    "North":   (0.30,  0.34,  0.38,  0.36),
}

REGIONS = ["National", "South", "Central", "North"]

def reg_vol(region, prod):
    if region == "National":
        return NAT[prod]["vol_cy"]
    idx = PRODUCTS.index(prod)
    return NAT[prod]["vol_cy"] * REG_SHARE[region][idx]

def reg_rev(region, prod):
    if region == "National":
        return NAT[prod]["rev_cy"]
    idx = PRODUCTS.index(prod)
    return NAT[prod]["rev_cy"] * REG_SHARE[region][idx]

def reg_tot_vol(region):
    return sum(reg_vol(region, p) for p in PRODUCTS)

def reg_tot_rev(region):
    return sum(reg_rev(region, p) for p in PRODUCTS)

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════
def rhex(c): return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def rect(sl,l,t,w,h,fill,line=None):
    s=sl.shapes.add_shape(1,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if line: s.line.color.rgb=line
    else:    s.line.fill.background()
    return s

def txt(sl,l,t,w,h,text,size=10,bold=False,fg=NAVY,
        align=PP_ALIGN.LEFT,italic=False,wrap=True):
    tb=sl.shapes.add_textbox(l,t,w,h)
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=str(text)
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=fg; r.font.italic=italic
    return tb

def cstyle(cell,text,size=8,bold=False,fg=NAVY,bg=None,
           align=PP_ALIGN.CENTER,italic=False):
    cell.text=str(text)
    tf=cell.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.runs[0] if p.runs else p.add_run()
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=fg; r.font.italic=italic
    if bg:
        pr=cell._tc.get_or_add_tcPr()
        sf=parse_xml(f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                     f'<a:srgbClr val="{rhex(bg)}"/></a:solidFill>')
        for ch in list(pr):
            if 'Fill' in ch.tag or 'fill' in ch.tag.lower(): pr.remove(ch)
        pr.append(sf)

def fp(v, d=1):
    if v is None: return "—"
    return f"+{v:.{d}f}%" if v >= 0 else f"{v:.{d}f}%"

def fv(v, d=1): return f"{v:,.{d}f}" if v is not None else "—"

def chg_bg(v):
    if v is None: return MGREY
    return LGREEN if v > 2 else (LRED if v < -2 else LYELL)

def bar_pct(sl, l, t, w, h, pct, fg_col, bg_col=LGREY):
    """Horizontal bar representing a percentage."""
    rect(sl, l, t, w, h, bg_col)
    rect(sl, l, t, w * min(pct, 1.0), h, fg_col)

def hdr(sl, title, subtitle=""):
    rect(sl, Inches(0), Inches(0), SW, Inches(0.75), NAVY)
    rect(sl, Inches(0), Inches(0.75), SW, Inches(0.04), GOLD)
    txt(sl, Inches(0.2),  Inches(0.06), Inches(10.5), Inches(0.40),
        title, size=14, bold=True, fg=GOLD)
    if subtitle:
        txt(sl, Inches(0.2), Inches(0.45), Inches(12.0), Inches(0.27),
            subtitle, size=8, fg=MGREY)
    txt(sl, Inches(10.9), Inches(0.25), Inches(2.2), Inches(0.24),
        "PSO Retail  |  10M FY26", size=7.5, fg=MGREY, align=PP_ALIGN.RIGHT)

def footer(sl, text):
    rect(sl, Inches(0), Inches(7.18), SW, Inches(0.32), NAVY)
    txt(sl, Inches(0.15), Inches(7.20), Inches(13), Inches(0.26),
        text, size=7.5, fg=MGREY, italic=True)

def sec_bar(sl, l, t, w, label, bg=DBLUE):
    rect(sl, l, t, w, Inches(0.26), bg)
    txt(sl, l+Inches(0.08), t+Inches(0.02), w-Inches(0.1), Inches(0.23),
        label, size=8.5, bold=True, fg=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
sl = prs.slides.add_slide(prs.slide_layouts[6])

hdr(sl, "PSO RETAIL FUELS — PRODUCT REVENUE & VOLUME PERFORMANCE",
    "All Products  |  National Overview + Regional Breakdown  |  10M FY26 vs 10M FY25  |  "
    "Revenue: PKR Billion  |  Volume: Million Litres  |  Regional figures are indicative")

# ── KPI STRIP ─────────────────────────────────────────────────────────────────
kpis = [
    ("Total Revenue CY",   f"PKR {fv(tot_rev_cy,1)} Bn",
     f"LY: PKR {fv(tot_rev_ly,1)} Bn", NAVY),
    ("Total Volume CY",    f"{fv(tot_vol_cy,1)} Mn L",
     f"LY: ~{fv(tot_vol_ly,0)} Mn L (est.)", DBLUE),
    ("Overall Rev Growth", fp(ovr_rev_chg),
     "All 4 products grew vs LY", GREEN),
    ("Fastest Growing",    "R95  +27.1%",
     "Premium segment surging — consumer upgrade trend", PROD_COL["R95"]),
]
kw = SW / len(kpis)
for ki, (lbl, val, sub, bg) in enumerate(kpis):
    rect(sl, ki*kw, Inches(0.80), kw, Inches(0.72), bg)
    txt(sl, ki*kw+Inches(0.04), Inches(0.81), kw-Inches(0.06), Inches(0.21),
        lbl, size=6.5, fg=MGREY, align=PP_ALIGN.CENTER)
    txt(sl, ki*kw+Inches(0.03), Inches(1.00), kw-Inches(0.05), Inches(0.32),
        val, size=12, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, ki*kw+Inches(0.03), Inches(1.30), kw-Inches(0.05), Inches(0.18),
        sub, size=6.5, fg=LGREY, align=PP_ALIGN.CENTER, italic=True)

# ── SECTION: NATIONAL PRODUCT TABLE (full width) ──────────────────────────────
sec_bar(sl, Inches(0.12), Inches(1.57), SW-Inches(0.24),
        "NATIONAL SUMMARY — REVENUE & VOLUME BY PRODUCT  |  "
        "Rev Chg% and Vol Chg% vs same period last year  |  "
        "Rev/Litre = average realisation per litre sold  |  Mix% = share of total")

# Table columns:
# Product | Rev CY | Rev LY | Rev Chg% | Rev Share | Vol CY | Vol LY(est) | Vol Chg%(est) | Vol Share | Rev/Litre
NCOLS = 10
tbl = sl.shapes.add_table(
    7, NCOLS,
    Inches(0.12), Inches(1.84),
    SW-Inches(0.24), Inches(2.56)
).table

CWS = [
    Inches(1.85),  # Product
    Inches(0.90),  # Rev CY
    Inches(0.90),  # Rev LY
    Inches(0.72),  # Rev Chg%
    Inches(0.65),  # Rev Share
    Inches(0.90),  # Vol CY
    Inches(0.90),  # Vol LY est
    Inches(0.72),  # Vol Chg% est
    Inches(0.65),  # Vol Share
    Inches(0.90),  # Rev/Litre
]
for ci, w in enumerate(CWS): tbl.columns[ci].width = w

# Row 0: group headers
for ci, (lbl, bg) in enumerate(zip(
    ["PRODUCT", "REVENUE (PKR Billion)", "", "", "",
     "VOLUME (Million Litres)", "", "", "", "REALISATION"],
    [NAVY, DBLUE, DBLUE, DBLUE, DBLUE,
     GREEN, GREEN, GREEN, GREEN, PROD_COL["R95"]]
)):
    cstyle(tbl.cell(0, ci), lbl, size=8, bold=True, fg=WHITE, bg=bg)
tbl.rows[0].height = Inches(0.22)

# Row 1: column names
COL_NAMES = [
    "Product", "CY (Bn)", "LY (Bn)", "Chg %", "Mix %",
    "CY (Mn L)", "LY est.", "Chg % est.", "Mix %", "PKR/Litre"
]
for ci, (h, bg) in enumerate(zip(COL_NAMES,
    [NAVY, DBLUE, DBLUE, DBLUE, DBLUE, GREEN, GREEN, GREEN, GREEN, PROD_COL["R95"]])):
    cstyle(tbl.cell(1, ci), h, size=7.5, bold=True, fg=WHITE, bg=bg,
           align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
tbl.rows[1].height = Inches(0.22)

# Data rows
ROW_H = Inches(0.42)
for ri, p in enumerate(PRODUCTS):
    d   = NAT[p]
    ri_t = ri + 2
    bg  = PROD_LIGHT[p]
    col = PROD_COL[p]

    cells = [
        (f"{p} — {PROD_FULL[p]}",    True,  PP_ALIGN.LEFT,   bg),
        (fv(d["rev_cy"]),             True,  PP_ALIGN.CENTER, bg),
        (fv(d["rev_ly"]),             False, PP_ALIGN.CENTER, bg),
        (fp(d["rev_chg"]),            True,  PP_ALIGN.CENTER, chg_bg(d["rev_chg"])),
        (f"{d['rev_sh']:.1f}%",       False, PP_ALIGN.CENTER, bg),
        (fv(d["vol_cy"], 1),          True,  PP_ALIGN.CENTER, bg),
        (f"~{fv(d['vol_ly'],0)}",     False, PP_ALIGN.CENTER, bg),
        (f"~{fp(d['vol_chg'])}",      False, PP_ALIGN.CENTER, chg_bg(d["vol_chg"])),
        (f"{d['vol_sh']:.1f}%",       False, PP_ALIGN.CENTER, bg),
        (f"{d['rpl']:,.0f}",          True,  PP_ALIGN.CENTER, bg),
    ]
    for ci, (v, bold, aln, cbg) in enumerate(cells):
        cstyle(tbl.cell(ri_t, ci), v, size=9, bold=bold, fg=NAVY, bg=cbg, align=aln)
    tbl.rows[ri_t].height = ROW_H

    # Color swatch on product name column
    rect(sl,
         Inches(0.12) + Inches(0.03),
         Inches(1.84) + Inches(0.22)*2 + ri*ROW_H + Inches(0.11),
         Inches(0.08), Inches(0.20), col)

# Total row
ri_tot = 6
tot_cells = [
    ("TOTAL", True, PP_ALIGN.LEFT, NAVY),
    (fv(tot_rev_cy), True, PP_ALIGN.CENTER, NAVY),
    (fv(tot_rev_ly), True, PP_ALIGN.CENTER, NAVY),
    (fp(ovr_rev_chg), True, PP_ALIGN.CENTER, NAVY),
    ("100%", True, PP_ALIGN.CENTER, NAVY),
    (fv(tot_vol_cy, 1), True, PP_ALIGN.CENTER, NAVY),
    (f"~{fv(tot_vol_ly,0)}", True, PP_ALIGN.CENTER, NAVY),
    (f"~{fp(ovr_vol_chg)}", True, PP_ALIGN.CENTER, NAVY),
    ("100%", True, PP_ALIGN.CENTER, NAVY),
    ("—", False, PP_ALIGN.CENTER, NAVY),
]
for ci, (v, bold, aln, bg) in enumerate(tot_cells):
    cstyle(tbl.cell(ri_tot, ci), v, size=9, bold=bold, fg=WHITE, bg=bg, align=aln)
tbl.rows[ri_tot].height = Inches(0.28)

# ── SECTION: REGIONAL BREAKDOWN (left 55%) + INSIGHTS (right 43%) ─────────────
REG_Y   = Inches(1.84) + Inches(0.22)*2 + 4*ROW_H + Inches(0.28) + Inches(0.10)
LEFT_W  = Inches(7.30)
RIGHT_L = Inches(7.52)
RIGHT_W = SW - RIGHT_L - Inches(0.12)

sec_bar(sl, Inches(0.12), REG_Y, LEFT_W,
        "REGIONAL BREAKDOWN — VOLUME (Mn L)  |  Indicative split of national volume")
sec_bar(sl, RIGHT_L, REG_Y, RIGHT_W,
        "KEY INSIGHTS", bg=PROD_COL["R95"])

# Regional table (left)
REG_COLS = 7  # Region | Total Vol | PMG | HSD | R95 | Lubes | Total Rev
reg_tbl_y = REG_Y + Inches(0.28)
rt = sl.shapes.add_table(
    len(REGIONS) + 2, REG_COLS,
    Inches(0.12), reg_tbl_y,
    LEFT_W, Inches(0.22) + Inches(0.22) + len(REGIONS)*Inches(0.36) + Inches(0.26)
).table

RT_CWS = [Inches(1.30), Inches(0.92), Inches(0.92), Inches(0.92),
          Inches(0.72), Inches(0.72), Inches(1.30)]
for ci, w in enumerate(RT_CWS): rt.columns[ci].width = w

# Row 0: group headers
for ci, (lbl, bg) in enumerate(zip(
    ["REGION", "TOTAL VOLUME", "PMG", "HSD", "R95", "Lubes", "TOTAL REVENUE"],
    [NAVY, LGREY, PROD_COL["PMG"], PROD_COL["HSD"], PROD_COL["R95"], PROD_COL["Lubes"], GREEN]
)):
    cstyle(rt.cell(0, ci), lbl, size=7.5, bold=True, fg=WHITE, bg=bg)
rt.rows[0].height = Inches(0.22)

# Row 1: sub-headers
for ci, (h, bg) in enumerate(zip(
    ["", "Mn Litres", "Mn Litres", "Mn Litres", "Mn Litres", "Mn Litres", "PKR Bn"],
    [NAVY, LGREY, PROD_LIGHT["PMG"], PROD_LIGHT["HSD"], PROD_LIGHT["R95"], PROD_LIGHT["Lubes"], LGREEN]
)):
    cstyle(rt.cell(1, ci), h, size=7, bold=False, fg=NAVY, bg=bg, italic=True)
rt.rows[1].height = Inches(0.22)

REG_BG = {
    "National": NAVY,
    "South":    DBLUE,
    "Central":  GREEN,
    "North":    ORANGE,
}
REG_ACCENT = {
    "National": MGREY,
    "South":    LBLUE,
    "Central":  LGREEN,
    "North":    LORANG,
}

for ri, region in enumerate(REGIONS):
    ri_t = ri + 2
    is_nat = region == "National"
    rbg  = REG_BG[region] if is_nat else REG_ACCENT[region]
    rfg  = WHITE if is_nat else NAVY
    rfg_b= WHITE if is_nat else PROD_COL.get("PMG", NAVY)

    tv  = reg_tot_vol(region)
    tr  = reg_tot_rev(region)
    vals = [
        (region, True,  PP_ALIGN.LEFT,   rbg),
        (fv(tv, 1), True, PP_ALIGN.CENTER, rbg),
        (fv(reg_vol(region, "PMG"), 1),   False, PP_ALIGN.CENTER, rbg),
        (fv(reg_vol(region, "HSD"), 1),   False, PP_ALIGN.CENTER, rbg),
        (fv(reg_vol(region, "R95"), 1),   False, PP_ALIGN.CENTER, rbg),
        (fv(reg_vol(region, "Lubes"), 1), False, PP_ALIGN.CENTER, rbg),
        (fv(tr, 1), True, PP_ALIGN.CENTER, rbg),
    ]
    for ci, (v, bold, aln, bg) in enumerate(vals):
        cstyle(rt.cell(ri_t, ci), v, size=9, bold=bold, fg=rfg, bg=bg, align=aln)
    rt.rows[ri_t].height = Inches(0.36)

    # Visual volume bars for PMG and HSD
    bar_x = Inches(0.12) + RT_CWS[0] + Inches(0.02)
    bar_w = RT_CWS[1] - Inches(0.04)
    bar_y = reg_tbl_y + Inches(0.22)*2 + ri*Inches(0.36) + Inches(0.26)
    if not is_nat:
        nat_tv = reg_tot_vol("National")
        bar_pct(sl, bar_x, bar_y, bar_w, Inches(0.08),
                tv / nat_tv if nat_tv else 0, REG_BG[region])

# Note row (indicative label)
ri_note = len(REGIONS) + 2
# Use rt.rows[ri_note] - but we need to add it. Actually the table has len(REGIONS)+2 = 6 rows.
# Last row index is 5 (0-based) = len(REGIONS)+2-1 = 5. Already added in table creation.

# ── INSIGHT BOXES (right side) ─────────────────────────────────────────────────
ins_y    = REG_Y + Inches(0.28)
ins_h_ea = Inches(0.68)
ins_gap  = Inches(0.06)

insights = [
    (DBLUE,
     "PMG & HSD DRIVE 97% OF REVENUE",
     (f"Motor Gasoline (PMG) is the largest product at PKR {fv(NAT['PMG']['rev_cy'])}Bn "
      f"({NAT['PMG']['rev_sh']:.0f}% of revenue, {NAT['PMG']['vol_sh']:.0f}% of volume). "
      f"HSD contributes {NAT['HSD']['rev_sh']:.0f}% — mostly commercial transport demand.")),
    (PROD_COL["R95"],
     "R95 IS THE GROWTH STORY (+27.1%)",
     (f"Premium Euro-5 grew the fastest at +27.1% revenue — consumers are upgrading to higher-grade fuel. "
      f"At PKR {NAT['R95']['rpl']:,.0f}/litre, it earns significantly more than standard PMG "
      f"(PKR {NAT['PMG']['rpl']:,.0f}/litre).")),
    (GREEN,
     "SOUTH LEADS VOLUME; NORTH LEADS HSD",
     (f"South (Karachi + Sindh) drives ~42% of PMG volume — high urban consumer density. "
      f"North contributes ~34–38% of HSD & R95 — longer transport routes and commercial "
      f"corridors mean higher diesel and premium fuel demand.")),
    (ORANGE,
     "LUBES: SMALL VOLUME, HIGH VALUE",
     (f"Lubricants are only {NAT['Lubes']['vol_sh']:.1f}% of total volume but earn "
      f"PKR {NAT['Lubes']['rpl']:,.0f}/litre — nearly 3× the PMG realization. "
      f"Growing at +14.5%, it is a high-margin segment worth expanding.")),
]

for ii, (accent, headline, body) in enumerate(insights):
    iy = ins_y + ii * (ins_h_ea + ins_gap)
    rect(sl, RIGHT_L, iy, RIGHT_W, ins_h_ea, LGREY)
    rect(sl, RIGHT_L, iy, Inches(0.06), ins_h_ea, accent)
    txt(sl, RIGHT_L+Inches(0.12), iy+Inches(0.06), RIGHT_W-Inches(0.16), Inches(0.20),
        headline, size=8, bold=True, fg=accent)
    txt(sl, RIGHT_L+Inches(0.12), iy+Inches(0.28), RIGHT_W-Inches(0.16), ins_h_ea-Inches(0.32),
        body, size=7.5, fg=NAVY, wrap=True)

footer(sl,
    "Revenue in PKR Billion  |  Volume in Million Litres  |  "
    "Vol LY and Vol Chg% are estimated (assumes price-flat scenario — actual may vary)  |  "
    "Regional split is indicative sample — update with actual regional data when available  |  "
    "Rev/Litre = Revenue ÷ Volume  |  10M FY26 = April 2025 to January 2026")

# ── SLIDE NOTES ───────────────────────────────────────────────────────────────
sl.notes_slide.notes_text_frame.text = (
    "SLIDE — PSO RETAIL PRODUCT REVENUE & VOLUME\n\n"
    "HOW TO PRESENT:\n"
    "Start with the KPI strip: PSO generated PKR 411Bn in revenue from 1,310Mn litres of fuel. "
    "All 4 products grew in revenue vs last year — a strong broad-based performance.\n\n"
    "PRODUCT TABLE:\n"
    "PMG is the backbone — 58% of revenue, 59% of volume. Growing at 20%.\n"
    "HSD is the commercial workhorse — 39% of revenue, 39% of volume. Also +20.5%.\n"
    "R95 (Premium Euro-5) is the star grower — +27.1%. Small absolute base (32Mn L) "
    "but highest revenue growth. Signals consumer preference shift to premium grades.\n"
    "Lubes: tiny volume but PKR 889/litre realisation — 3× PMG's. High-margin segment.\n\n"
    "REGIONAL:\n"
    "South (Karachi+Sindh) = biggest consumer market, PMG-heavy.\n"
    "North = longer routes, HSD-heavy (trucks, transport).\n"
    "Central (Punjab) = balanced mix.\n\n"
    "QUESTIONS:\n"
    "• Why is R95 growing so fast? — Euro-5 mandate pushed demand for higher-octane fuel. "
    "New car models increasingly require 95-octane.\n"
    "• Revenue grew 20% but volume — is that just price? — Vol LY is estimated here. "
    "If actual vol LY shows lower growth than revenue, it implies a price component.\n"
    "• Lubes only 3.6Mn L but PKR 3.2Bn? — Yes. Lubes are sold by weight/pack, "
    "not just litres, so the realisation per litre is not directly comparable to fuels.\n"
)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out = Path("reports/PSO_Product_Revenue.pptx")
out.parent.mkdir(exist_ok=True)
prs.save(str(out))
print(f"Saved: {out.resolve()}")
print(f"\nNATIONAL SUMMARY:")
print(f"{'Product':<12} {'RevCY':>8} {'RevLY':>8} {'RevChg':>8} {'VolCY':>10} {'RevSh':>7} {'VolSh':>7} {'RPL':>8}")
print("-"*70)
for p in PRODUCTS:
    d = NAT[p]
    print(f"{p:<12} {d['rev_cy']:>8.1f} {d['rev_ly']:>8.1f} {d['rev_chg']:>7.1f}% "
          f"{d['vol_cy']:>10.1f} {d['rev_sh']:>6.1f}% {d['vol_sh']:>6.1f}% {d['rpl']:>8.0f}")
print("-"*70)
print(f"{'TOTAL':<12} {tot_rev_cy:>8.1f} {tot_rev_ly:>8.1f} {ovr_rev_chg:>7.1f}% "
      f"{tot_vol_cy:>10.1f} {'100.0':>6}% {'100.0':>6}%")
