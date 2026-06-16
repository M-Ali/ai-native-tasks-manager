"""
PSO — Top 5 Cities Volume Focus Deck
1 cover + (1 city divider + 7 slides) × 5 = 41 slides
Objective: increase volumes. Notes = presenter talking points.
"""
import sys, os
sys.path.insert(0, 'src')
os.environ['PYTHONIOENCODING'] = 'utf-8'

import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.oxml import parse_xml
from pso import ingest

# ── palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B,0x2A,0x4A)
GOLD   = RGBColor(0xC9,0xA0,0x30)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LGREY  = RGBColor(0xF2,0xF4,0xF8)
MGREY  = RGBColor(0xD9,0xDC,0xE3)
DBLUE  = RGBColor(0x2E,0x5B,0x9A)
LBLUE  = RGBColor(0xBD,0xD7,0xEE)
GREEN  = RGBColor(0x37,0x5B,0x25)
LGREEN = RGBColor(0xC6,0xEF,0xCE)
RED    = RGBColor(0xC0,0x00,0x00)
LRED   = RGBColor(0xFF,0xC7,0xCE)
YELL   = RGBColor(0xFF,0xC0,0x00)
LYELL  = RGBColor(0xFF,0xEB,0x9C)
ORANGE = RGBColor(0xFF,0x82,0x00)
LORANG = RGBColor(0xFF,0xE0,0xB2)

SW, SH = Inches(13.33), Inches(7.50)

def rhex(c): return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

# ═══════════════════════════════════════════════════════════════════════════════
# DRAWING PRIMITIVES
# ═══════════════════════════════════════════════════════════════════════════════
def rect(sl, l,t,w,h, fill, line=None):
    s = sl.shapes.add_shape(1,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def txt(sl, l,t,w,h, text, size=10, bold=False, fg=NAVY, bg=None,
        align=PP_ALIGN.LEFT, italic=False):
    if bg: rect(sl,l,t,w,h,bg)
    tb = sl.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = str(text)
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=fg;  r.font.italic=italic
    return tb

def multi_txt(sl, l,t,w,h, lines, size=9, fg=NAVY, bg=None, spacing=None):
    """Write multiple bullet lines into one textbox."""
    if bg: rect(sl,l,t,w,h,bg)
    tb = sl.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = str(line)
        r.font.size=Pt(size); r.font.color.rgb=fg
        if spacing: p.space_before = spacing
    return tb

def cstyle(cell, text, size=8, bold=False, fg=NAVY, bg=None,
           align=PP_ALIGN.CENTER, italic=False):
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=fg;  r.font.italic=italic
    if bg:
        pr = cell._tc.get_or_add_tcPr()
        sf = parse_xml(f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                       f'<a:srgbClr val="{rhex(bg)}"/></a:solidFill>')
        for ch in list(pr):
            if 'Fill' in ch.tag or 'fill' in ch.tag.lower(): pr.remove(ch)
        pr.append(sf)

def set_notes(sl, text):
    sl.notes_slide.notes_text_frame.text = text

def hdr(sl, title, subtitle=""):
    rect(sl, Inches(0),Inches(0),SW,Inches(0.75),NAVY)
    rect(sl, Inches(0),Inches(0.75),SW,Inches(0.04),GOLD)
    txt(sl, Inches(0.2),Inches(0.06),Inches(10.5),Inches(0.40),
        title, size=15, bold=True, fg=GOLD)
    if subtitle:
        txt(sl, Inches(0.2),Inches(0.45),Inches(11.5),Inches(0.27),
            subtitle, size=8.5, fg=MGREY)
    txt(sl, Inches(10.9),Inches(0.25),Inches(2.2),Inches(0.24),
        "PSO Retail  |  10M FY26", size=7.5, fg=MGREY, align=PP_ALIGN.RIGHT)

def footer(sl, text):
    rect(sl, Inches(0),Inches(7.18),SW,Inches(0.32),NAVY)
    txt(sl,  Inches(0.15),Inches(7.20),Inches(13),Inches(0.26),
        text, size=7.5, fg=MGREY, italic=True)

def sec_bar(sl, l,t,w, label, bg=DBLUE):
    rect(sl, l,t,w,Inches(0.27),bg)
    txt(sl,  l+Inches(0.08),t+Inches(0.02),w-Inches(0.1),Inches(0.24),
        label, size=8.5, bold=True, fg=WHITE)

def insight(sl, l,t,w,h, title, body, tbg=GOLD, bbg=LYELL, tfg=NAVY, bfg=NAVY):
    rect(sl,l,t,w,Inches(0.30),tbg)
    txt(sl, l+Inches(0.08),t+Inches(0.03),w-Inches(0.12),Inches(0.25),
        title, size=9, bold=True, fg=tfg)
    rect(sl,l,t+Inches(0.30),w,h-Inches(0.30),bbg)
    txt(sl, l+Inches(0.10),t+Inches(0.34),w-Inches(0.16),h-Inches(0.40),
        body, size=8.5, fg=bfg)

def kpi(sl, l,t,w,h, label, value, sub="", bg=DBLUE, vfg=WHITE):
    rect(sl,l,t,w,h,bg)
    txt(sl, l+Inches(0.06),t+Inches(0.04),w-Inches(0.08),Inches(0.22),
        label, size=6.5, fg=LGREY, align=PP_ALIGN.CENTER)
    txt(sl, l+Inches(0.04),t+Inches(0.24),w-Inches(0.06),Inches(0.38),
        value, size=14, bold=True, fg=vfg, align=PP_ALIGN.CENTER)
    if sub:
        txt(sl, l+Inches(0.04),t+Inches(0.60),w-Inches(0.06),Inches(0.22),
            sub, size=7, fg=LGREY, align=PP_ALIGN.CENTER, italic=True)

def bar_chart(sl, l,t,w,h, cats, cy_v, ly_v, period):
    cd = ChartData()
    cd.categories = cats
    cd.add_series(f"CY {period}", cy_v)
    cd.add_series("LY", ly_v)
    cf = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,l,t,w,h,cd)
    ch = cf.chart
    ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout=False; ch.has_title=False
    try:
        ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb=NAVY
        ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb=MGREY
        ch.category_axis.tick_labels.font.size=Pt(8)
        ch.value_axis.tick_labels.font.size=Pt(7)
    except: pass

def chg_bg(v):
    if v is None: return MGREY
    return LGREEN if v>2 else (LRED if v<-2 else LYELL)

def fv(v,d=1): return f"{v:,.{d}f}" if v is not None else "—"
def fp(v,d=1):
    if v is None: return "—"
    return f"+{v:.{d}f}%" if v>=0 else f"{v:.{d}f}%"
def spct(cy,ly): return (cy-ly)/abs(ly)*100 if ly and ly!=0 else None

# ═══════════════════════════════════════════════════════════════════════════════
# DATA
# ═══════════════════════════════════════════════════════════════════════════════
print("Loading …")
df, _ = ingest.load("data/input/Working File Retail Fuels Data.xlsx")
period = df["_Period"].iloc[0]
retail = df[df["IsRetail"] & ~df["IsInternational"]].copy()

top10 = retail.groupby("CityNorm")["SalesLtr_CY"].sum().nlargest(10).index.tolist()
TOP5  = top10[:5]
retail_nat_vol = retail["SalesLtr_CY"].sum()/1e6

# 10-city per-station benchmark
t10_vols = {c: retail[retail["CityNorm"]==c]["SalesLtr_CY"].sum()/1e6 for c in top10}
t10_stns = {c: retail[retail["CityNorm"]==c]["Customer Number"].nunique() for c in top10}
t10_avg_per_stn = np.mean([t10_vols[c]/t10_stns[c] for c in top10])
t10_avg_city    = np.mean(list(t10_vols.values()))

DIESEL_P = ["HSD","LDO","SLUDGE-HSD","SLUDGE"]
PETROL_P = ["PMG","R95"]
LUBE_C   = ["DEO","PCMO","MCO","LOW GRADE","INDUSTRIAL GRADE","Greases","OTHERS"]
SEGS     = ["Diesel","Petrol","Lubricants"]

def compute(city_df):
    d = {}
    # totals
    d["tot_cy"]  = city_df["SalesLtr_CY"].sum()/1e6
    d["tot_ly"]  = city_df["SalesLtr_LY"].sum()/1e6
    d["tot_chg"] = spct(d["tot_cy"],d["tot_ly"])
    d["nat_sh"]  = d["tot_cy"]/retail_nat_vol*100

    # segments
    segs={}
    for seg in SEGS:
        s=city_df[city_df["FuelSegment"]==seg]
        cy=s["SalesLtr_CY"].sum()/1e6; ly=s["SalesLtr_LY"].sum()/1e6
        segs[seg]=dict(cy=cy,ly=ly,chg=spct(cy,ly),
                       sh=cy/d["tot_cy"]*100 if d["tot_cy"] else 0,
                       stns=s["Customer Number"].nunique())
    d["segs"]=segs

    # sub-products
    sub={}
    for p in DIESEL_P:
        s=city_df[(city_df["FuelSegment"]=="Diesel")&(city_df["ProductCategory"]==p)]
        cy=s["SalesLtr_CY"].sum()/1e6; ly=s["SalesLtr_LY"].sum()/1e6
        sub[p]=dict(grp="Diesel",cy=cy,ly=ly,chg=spct(cy,ly),stns=s["Customer Number"].nunique())
    for p in PETROL_P:
        s=city_df[(city_df["FuelSegment"]=="Petrol")&(city_df["ProductCategory"]==p)]
        cy=s["SalesLtr_CY"].sum()/1e6; ly=s["SalesLtr_LY"].sum()/1e6
        sub[p]=dict(grp="Petrol",cy=cy,ly=ly,chg=spct(cy,ly),stns=s["Customer Number"].nunique())
    for c in LUBE_C:
        s=city_df[(city_df["FuelSegment"]=="Lubricants")&(city_df["LubeCategory"]==c)]
        cy=s["SalesLtr_CY"].sum()/1e6; ly=s["SalesLtr_LY"].sum()/1e6
        sub[c]=dict(grp="Lubricants",cy=cy,ly=ly,chg=spct(cy,ly),stns=s["Customer Number"].nunique())
    d["sub"]=sub

    # stations
    stn=(city_df.groupby("Customer Number",as_index=False)
         .agg(name=("Name 1","first"),
              vol_cy=("SalesLtr_CY","sum"),
              vol_ly=("SalesLtr_LY","sum"))
         .sort_values("vol_cy",ascending=False).reset_index(drop=True))
    stn["ml_cy"]=stn["vol_cy"]/1e6; stn["ml_ly"]=stn["vol_ly"]/1e6
    stn["chg"]  =stn.apply(lambda r: spct(r["vol_cy"],r["vol_ly"]),axis=1)
    stn["sh"]   =stn["vol_cy"]/stn["vol_cy"].sum()*100
    stn["cumsh"]=stn["sh"].cumsum()
    stn["active"]=stn["vol_cy"]>0
    d["stn"]=stn

    n=len(stn)
    d["n_tot"]   =n
    d["n_inact"] =int((~stn["active"]).sum())
    d["n_active"]=n-d["n_inact"]
    d["n80"]     =int((stn["cumsh"]<80).sum())+1

    # product map
    pm={}
    for cnum,grp in city_df.groupby("Customer Number"):
        fp=set(grp[grp["FuelSegment"]!="Lubricants"]["ProductCategory"].dropna())
        lp=set(grp[grp["FuelSegment"]=="Lubricants"]["LubeCategory"].dropna())
        pm[cnum]=dict(fp=fp,lp=lp,
                      has_diesel=bool({"HSD","LDO"}&fp),
                      has_petrol=bool({"PMG","R95"}&fp),
                      has_lubes =bool(lp),
                      has_hsd="HSD" in fp, has_ldo="LDO" in fp,
                      has_pmg="PMG" in fp, has_r95="R95" in fp,
                      has_deo="DEO" in lp, has_pcmo="PCMO" in lp,
                      has_mco="MCO" in lp, has_lg="LOW GRADE" in lp)
    d["pm"]=pm
    pp=list(pm.values())
    d["n_all3"]   =sum(1 for p in pp if p["has_diesel"] and p["has_petrol"] and p["has_lubes"])
    d["n_d_p"]    =sum(1 for p in pp if p["has_diesel"] and p["has_petrol"] and not p["has_lubes"])
    d["n_d_only"] =sum(1 for p in pp if p["has_diesel"] and not p["has_petrol"] and not p["has_lubes"])
    d["n_p_only"] =sum(1 for p in pp if p["has_petrol"] and not p["has_diesel"] and not p["has_lubes"])
    d["n_lubes"]  =sum(1 for p in pp if p["has_lubes"])
    d["n_petrol"] =sum(1 for p in pp if p["has_petrol"])
    d["n_diesel"] =sum(1 for p in pp if p["has_diesel"])
    d["n_r95"]    =sum(1 for p in pp if p["has_r95"])
    d["n_pmg"]    =sum(1 for p in pp if p["has_pmg"])
    d["n_d_no_p"] =sum(1 for p in pp if p["has_diesel"] and not p["has_petrol"])
    d["n_fl"]     =sum(1 for p in pp if (p["has_diesel"] or p["has_petrol"]) and not p["has_lubes"])

    # volumes of gap segments
    dnp_custs=[c for c,p in pm.items() if p["has_diesel"] and not p["has_petrol"]]
    fl_custs =[c for c,p in pm.items() if (p["has_diesel"] or p["has_petrol"]) and not p["has_lubes"]]
    lub_custs=[c for c,p in pm.items() if p["has_lubes"]]

    d["dnp_vol"]=stn[stn["Customer Number"].isin(dnp_custs)]["ml_cy"].sum()
    d["fl_vol"] =stn[stn["Customer Number"].isin(fl_custs)]["ml_cy"].sum()

    lub_vol_cy=city_df[(city_df["FuelSegment"]=="Lubricants")&
                       (city_df["Customer Number"].isin(lub_custs))]["SalesLtr_CY"].sum()/1e6
    d["avg_lub_per_stn"]=lub_vol_cy/max(len(lub_custs),1)

    # benchmarks
    d["vol_per_stn"]  =d["tot_cy"]/n
    d["active_avg"]   =stn[stn["active"]]["ml_cy"].mean() if stn["active"].any() else 0
    d["top25_avg"]    =stn.head(max(1,n//4))["ml_cy"].mean()
    d["median_vol"]   =stn[stn["active"]]["ml_cy"].median() if stn["active"].any() else 0

    # recovery potentials
    inact=stn[~stn["active"]]
    d["inact_ly_vol"] =inact["ml_ly"].sum()
    d["rec_inact"]    =d["inact_ly_vol"]*0.65

    undp=stn[(stn["active"])&(stn["vol_ly"]>0)&(stn["vol_cy"]<stn["vol_ly"]*0.5)]
    d["n_undp"]       =len(undp)
    d["rec_undp"]     =(undp["ml_ly"]*0.75-undp["ml_cy"]).sum()
    d["undp_df"]      =undp

    d["rec_lub"]      =len(fl_custs)*d["avg_lub_per_stn"]
    d["rec_petrol"]   =d["dnp_vol"]*0.25
    d["tot_pot"]      =d["rec_inact"]+d["rec_undp"]+d["rec_lub"]+d["rec_petrol"]

    # product breadth vs volume
    pv={"1 product":[],"2 products":[],"3+ products":[]}
    for _,row in stn.iterrows():
        if not row["active"]: continue
        p=pm.get(row["Customer Number"],{})
        nc=(1 if p.get("has_diesel") else 0)+(1 if p.get("has_petrol") else 0)+(1 if p.get("has_lubes") else 0)
        k="1 product" if nc<=1 else "2 products" if nc==2 else "3+ products"
        pv[k].append(row["ml_cy"])
    d["pv"]=pv
    d["mul_avg"]=np.mean(pv.get("3+ products",[0]))
    d["one_avg"]=np.mean(pv.get("1 product",[0.001]))

    # dropped products
    d["dropped"]=[k for k,v in sub.items() if v["cy"]==0 and v["ly"]>0]
    d["never"]  =[k for k,v in sub.items() if v["cy"]==0 and v["ly"]==0]

    # star product list per station helper
    def prod_str(cnum):
        p=pm.get(cnum,{})
        parts=[]
        if p.get("has_hsd"): parts.append("HSD")
        if p.get("has_ldo"): parts.append("LDO")
        if p.get("has_pmg"): parts.append("PMG")
        if p.get("has_r95"): parts.append("R95")
        if p.get("has_lubes"):parts.append("LUBES")
        return "/".join(parts) or "—"
    d["prod_str_fn"]=prod_str

    return d

print("Computing city data …")
city_data = {}
for city in TOP5:
    city_data[city] = compute(retail[retail["CityNorm"]==city].copy())
    d=city_data[city]
    print(f"  {city}: {d['tot_cy']:.1f}ML  {d['n80']}/{d['n_tot']} stns for 80%  "
          f"pot +{d['tot_pot']:.1f}ML")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

# ── COVER ─────────────────────────────────────────────────────────────────────
def make_cover(prs):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl,Inches(0),Inches(0),SW,SH,NAVY)
    rect(sl,Inches(0),Inches(2.9),SW,Inches(0.06),GOLD)
    rect(sl,Inches(0),Inches(4.3),SW,Inches(0.06),GOLD)
    txt(sl,Inches(0.8),Inches(1.1),Inches(11.7),Inches(0.9),
        "PSO RETAIL — CITY VOLUME ANALYSIS",size=24,bold=True,fg=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,Inches(0.8),Inches(2.0),Inches(11.7),Inches(0.7),
        f"Top 5 Cities  |  Volume Deep-Dive  |  10M FY26 vs 10M FY25",
        size=14,fg=GOLD,align=PP_ALIGN.CENTER)
    txt(sl,Inches(0.8),Inches(3.05),Inches(11.7),Inches(1.1),
        "Volume Scorecard  ·  80% Station Rule  ·  Star Performers  ·  Product Gaps\n"
        "Sub-Product Mix  ·  Lost Volume Recovery  ·  Growth Roadmap",
        size=11,fg=LGREY,align=PP_ALIGN.CENTER)
    txt(sl,Inches(0.8),Inches(4.45),Inches(11.7),Inches(0.5),
        "  ·  ".join(TOP5),size=12,fg=GOLD,bold=True,align=PP_ALIGN.CENTER)
    txt(sl,Inches(0.8),Inches(5.2),Inches(11.7),Inches(0.4),
        "Objective: Identify actionable volume growth levers from station-level data",
        size=10,fg=MGREY,align=PP_ALIGN.CENTER,italic=True)
    txt(sl,Inches(0.8),Inches(6.6),Inches(11.7),Inches(0.3),
        "For PSO Management — Internal Use Only  |  Source: Working File Retail Fuels",
        size=8.5,fg=MGREY,align=PP_ALIGN.CENTER,italic=True)
    set_notes(sl,
        "COVER SLIDE — HOW TO PRESENT THIS DECK\n\n"
        "This deck covers the top 5 PSO retail cities by volume: "+", ".join(TOP5)+".\n\n"
        "Each city section has 7 slides:\n"
        "  1. Volume Scorecard — the headline numbers: where does this city stand?\n"
        "  2. The 80% Rule — which stations are actually carrying this city?\n"
        "  3. Star Performers — what does a great station look like? Can others match it?\n"
        "  4. Product Gaps — which stations are missing which products, and what volume is that costing?\n"
        "  5. Sub-Product Mix — inside each segment, what is growing and what is declining?\n"
        "  6. Lost Volume — inactive stations and severe under-performers. Recoverable volume.\n"
        "  7. Growth Roadmap — four prioritised levers with targets and timelines.\n\n"
        "PRESENTING TIPS:\n"
        "• Open with slide 1 to set context, then jump to slide 7 (roadmap) so the audience "
        "knows where you are going. Then walk through slides 2–6 as evidence.\n"
        "• All numbers come directly from the Working File. No projections, no assumptions "
        "except where explicitly stated (recovery rates, add-on potentials).\n"
        "• The objective is to increase volume. Every slide should end with an action question: "
        "'Who owns this lever, and by when?'\n"
    )

# ── CITY DIVIDER ──────────────────────────────────────────────────────────────
def make_divider(prs, city, d, rank):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl,Inches(0),Inches(0),SW,SH,NAVY)
    rect(sl,Inches(0),Inches(3.5),SW,Inches(0.06),GOLD)
    txt(sl,Inches(0.6),Inches(0.5),Inches(11),Inches(0.5),
        f"CITY {rank} OF 5",size=13,fg=GOLD,align=PP_ALIGN.CENTER)
    txt(sl,Inches(0.6),Inches(1.0),Inches(11),Inches(1.2),
        city.upper(),size=38,bold=True,fg=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,Inches(0.6),Inches(2.2),Inches(11),Inches(0.5),
        f"Volume: {fv(d['tot_cy'])} ML  |  {d['n_tot']} Stations  |  "
        f"National Share: {d['nat_sh']:.1f}%  |  YoY: {fp(d['tot_chg'])}",
        size=12,fg=GOLD,align=PP_ALIGN.CENTER)
    # 4 quick stats
    qs=[("Total Volume CY",f"{fv(d['tot_cy'])} ML"),
        ("Stations for 80%",f"{d['n80']} of {d['n_tot']}"),
        ("Inactive Stations",str(d['n_inact'])),
        ("Vol Opportunity",f"+{fv(d['tot_pot'])} ML")]
    qw=SW/4
    for qi,(lbl,val) in enumerate(qs):
        rect(sl,qi*qw+Inches(0.05),Inches(3.7),qw-Inches(0.1),Inches(1.0),DBLUE)
        txt(sl,qi*qw+Inches(0.1),Inches(3.75),qw-Inches(0.15),Inches(0.3),
            lbl,size=8,fg=LGREY,align=PP_ALIGN.CENTER)
        txt(sl,qi*qw+Inches(0.1),Inches(4.0),qw-Inches(0.15),Inches(0.55),
            val,size=16,bold=True,fg=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,Inches(0.6),Inches(5.0),Inches(11),Inches(0.4),
        "Slides in this section: Scorecard → 80% Rule → Stars → Product Gaps → Sub-Products → Lost Volume → Roadmap",
        size=9,fg=MGREY,align=PP_ALIGN.CENTER,italic=True)
    set_notes(sl,
        f"CITY DIVIDER — {city.upper()}\n\n"
        f"City rank: #{rank} by volume among top 5 PSO retail cities.\n"
        f"Total CY volume: {fv(d['tot_cy'])}ML  |  LY: {fv(d['tot_ly'])}ML  |  Change: {fp(d['tot_chg'])}\n"
        f"Stations: {d['n_tot']} total, {d['n_active']} active, {d['n_inact']} inactive\n"
        f"Volume per station CY: {d['vol_per_stn']:.2f}ML  |  10-city avg: {t10_avg_per_stn:.2f}ML\n"
        f"Identifiable volume opportunity: +{fv(d['tot_pot'])}ML (+{d['tot_pot']/d['tot_cy']*100:.0f}% on CY base)\n\n"
        "TALKING POINT: Before going into the detail, tell the audience the headline for this city. "
        "Is it a growing city or a declining one? Is the volume problem in Diesel, Petrol or Lubes? "
        "The divider gives you the top-line facts to anchor the story.\n"
    )

# ── SLIDE 1: SCORECARD ────────────────────────────────────────────────────────
def make_s1(prs, city, d, period):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    hdr(sl,f"{city.upper()} — VOLUME SCORECARD",
        "Headline performance: total volume CY vs LY, segment breakdown, and how this city compares to its peers.")

    # KPI strip
    chg_c=GREEN if (d["tot_chg"] or 0)>=0 else RED
    kpis=[
        ("Total Vol CY (ML)",  fv(d["tot_cy"]),   f"LY: {fv(d['tot_ly'])} ML",  WHITE,NAVY),
        ("YoY Volume Change",  fp(d["tot_chg"]),   "vs same period LY",          WHITE,chg_c),
        ("Natl Vol Share",     f"{d['nat_sh']:.1f}%","of PSO all-retail",        WHITE,DBLUE),
        ("Vol / Station CY",   f"{d['vol_per_stn']:.2f}ML",
                               f"10-city avg: {t10_avg_per_stn:.2f}ML",          WHITE,
                               GREEN if d["vol_per_stn"]>=t10_avg_per_stn else RED),
        ("Active / Total Stns",f"{d['n_active']}/{d['n_tot']}",
                               f"{d['n_inact']} inactive",                       WHITE,DBLUE),
    ]
    kw=Inches(2.5)
    for ki,(lbl,val,sub_,vc,bg) in enumerate(kpis):
        kpi(sl,Inches(0.18)+ki*kw,Inches(0.83),kw-Inches(0.07),Inches(0.95),lbl,val,sub_,bg)

    # bar chart
    bar_chart(sl,Inches(0.18),Inches(1.92),Inches(5.5),Inches(2.8),
              SEGS,[d["segs"][s]["cy"] for s in SEGS],[d["segs"][s]["ly"] for s in SEGS],period)
    txt(sl,Inches(0.18),Inches(1.84),Inches(5.5),Inches(0.22),
        "Volume by Segment (ML) — CY vs LY",size=8,bold=True,fg=NAVY)

    # segment table
    sec_bar(sl,Inches(5.9),Inches(1.84),Inches(7.25),"Segment Detail")
    t=sl.shapes.add_table(5,6,Inches(5.9),Inches(2.14),Inches(7.25),Inches(1.7)).table
    for ci,h in enumerate(["Segment","Stns","Vol CY (ML)","Vol LY (ML)","Change","City Share"]):
        cstyle(t.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
    for ci,w in enumerate([Inches(1.4),Inches(0.85),Inches(1.25),Inches(1.25),Inches(1.0),Inches(1.5)]):
        t.columns[ci].width=w
    t.rows[0].height=Inches(0.27)
    for ri,seg in enumerate(SEGS+["Total"],1):
        if seg=="Total": sv=dict(cy=d["tot_cy"],ly=d["tot_ly"],chg=d["tot_chg"],sh=100,stns=d["n_tot"]); bg=MGREY
        else: sv=d["segs"][seg]; bg=LGREY if ri%2==0 else WHITE
        vals=[seg,str(sv["stns"]),fv(sv["cy"]),fv(sv["ly"]),fp(sv["chg"]),f"{sv.get('sh',sv['cy']/d['tot_cy']*100 if d['tot_cy'] else 0):.1f}%"]
        for ci,v in enumerate(vals):
            cstyle(t.cell(ri,ci),v,size=8.5,bold=(seg=="Total"),fg=NAVY,
                   bg=chg_bg(sv["chg"]) if ci==4 else bg,
                   align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
        t.rows[ri].height=Inches(0.34)

    # peer comparison
    sec_bar(sl,Inches(5.9),Inches(4.05),Inches(7.25),"City Benchmarks — vs 10-City Average")
    peers=[
        ("Total Volume (ML)",        fv(d["tot_cy"]),      f"10-city avg: {t10_avg_city:.1f}ML",
         LGREEN if d["tot_cy"]>=t10_avg_city else LRED),
        ("Vol per Station (ML)",     f"{d['vol_per_stn']:.2f}", f"10-city avg: {t10_avg_per_stn:.2f}ML",
         LGREEN if d["vol_per_stn"]>=t10_avg_per_stn else LRED),
        ("Active station avg (ML)",  f"{d['active_avg']:.2f}", "per active station in city",LGREY),
        ("Top-25% station avg (ML)", f"{d['top25_avg']:.2f}",  "the star-performer benchmark",LGREY),
        ("Median active stn (ML)",   f"{d['median_vol']:.2f}", "50th percentile — midpoint",LGREY),
    ]
    pt=sl.shapes.add_table(6,3,Inches(5.9),Inches(4.34),Inches(7.25),Inches(2.1)).table
    for ci,h in enumerate(["Metric","Value","Context"]):
        cstyle(pt.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
    pt.columns[0].width=Inches(3.8); pt.columns[1].width=Inches(1.3); pt.columns[2].width=Inches(2.15)
    pt.rows[0].height=Inches(0.26)
    for ri,(lbl,val,ctx,bg) in enumerate(peers,1):
        cstyle(pt.cell(ri,0),lbl,size=8,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
        cstyle(pt.cell(ri,1),val,size=9,bold=True,fg=NAVY,bg=bg)
        cstyle(pt.cell(ri,2),ctx,size=7.5,fg=NAVY,bg=bg,italic=True,align=PP_ALIGN.LEFT)
        pt.rows[ri].height=Inches(0.33)

    # bottom insight
    above_avg = d["vol_per_stn"]>=t10_avg_per_stn
    dsl=d["segs"]["Diesel"]; pet=d["segs"]["Petrol"]; lub=d["segs"]["Lubricants"]
    insight_body=(
        f"{city} delivered {fv(d['tot_cy'])}ML in 10M FY26, a {fp(d['tot_chg'])} change vs LY. "
        f"Diesel dominates at {dsl['sh']:.0f}% of volume. "
        f"Petrol is {pet['sh']:.0f}% and Lubricants {lub['sh']:.1f}%. "
        f"At {d['vol_per_stn']:.2f}ML per station, this city is "
        f"{'above' if above_avg else 'BELOW'} the 10-city fleet average of {t10_avg_per_stn:.2f}ML."
    )
    insight(sl,Inches(0.18),Inches(4.88),Inches(5.5),Inches(2.0),
            "KEY TAKEAWAY",insight_body)

    footer(sl,"ML = Million Litres  |  CY = 10M FY26  |  LY = 10M FY25  |  Green=growth >2%  Red=decline >2%  Yellow=flat")
    set_notes(sl,
        f"SLIDE 1 — {city.upper()} VOLUME SCORECARD\n\n"
        f"PRESENTER TALKING POINTS:\n"
        f"Open by anchoring the audience on the total volume number: {fv(d['tot_cy'])}ML for 10 months.\n"
        f"{'Volume is growing ' + fp(d['tot_chg']) + ' — that is positive context.' if (d['tot_chg'] or 0)>0 else 'Volume is DECLINING ' + fp(d['tot_chg']) + ' — this is the central problem we need to address.'}\n\n"
        f"NUMBERS TO EMPHASISE:\n"
        f"• Diesel: {fv(dsl['cy'])}ML ({dsl['sh']:.0f}% of city) — {fp(dsl['chg'])} YoY\n"
        f"• Petrol: {fv(pet['cy'])}ML ({pet['sh']:.0f}% of city) — {fp(pet['chg'])} YoY\n"
        f"• Lubricants: {fv(lub['cy'])}ML ({lub['sh']:.1f}% of city) — {fp(lub['chg'])} YoY\n"
        f"• Vol/station: {d['vol_per_stn']:.2f}ML vs 10-city avg {t10_avg_per_stn:.2f}ML "
        f"— city is {'above' if above_avg else 'below'} average productivity\n\n"
        f"QUESTIONS YOU MIGHT GET:\n"
        f"• 'Why is diesel declining?' — go to Slide 5 for product mix and Slide 6 for inactive stations\n"
        f"• 'How do we compare with Lahore/Karachi?' — this slide shows the per-station benchmark\n"
        f"• 'What are we doing about the inactive stations?' — go to Slide 7 (roadmap, Lever 1)\n\n"
        f"WATCH OUT: If GRS is growing but volume is declining, the city is experiencing price-led revenue. "
        f"The volume base is eroding and this will catch up. Flag it.\n"
    )

# ── SLIDE 2: 80% RULE ────────────────────────────────────────────────────────
def make_s2(prs, city, d, period):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    stn=d["stn"]; n80=d["n80"]; n_tot=d["n_tot"]
    hdr(sl,f"THE 80% RULE: {n80} STATIONS CARRY {city.upper()}",
        f"{n80} of {n_tot} stations ({n80/n_tot*100:.0f}% of fleet) generate 80% of {city}'s total volume. "
        f"The other {n_tot-n80} stations share only 20%.")

    # KPI strip
    k2=[("Stns for 80% Vol",str(n80),f"{n80/n_tot*100:.0f}% of fleet",NAVY),
        ("Remaining stns",str(n_tot-n80),"share 20% of volume",RED),
        ("Top station vol",f"{stn.iloc[0]['ml_cy']:.1f}ML",stn.iloc[0]["name"][:22],DBLUE),
        ("Top 5 stns",f"{stn.head(5)['sh'].sum():.1f}%","of city volume",DBLUE),
        ("Top 10 stns",f"{stn.head(10)['sh'].sum():.1f}%","of city volume",DBLUE),
        ("Avg: top-80%",f"{stn.head(n80)['ml_cy'].mean():.2f}ML","per stn in top tier",GREEN),
        ("Avg: rest",f"{stn.tail(n_tot-n80)['ml_cy'].mean():.2f}ML","per stn in bottom tier",RED),]
    kw2=SW/len(k2)
    for ki,(lbl,val,sub_,bg) in enumerate(k2):
        rect(sl,ki*kw2,Inches(0.80),kw2,Inches(0.75),bg)
        txt(sl,ki*kw2+Inches(0.04),Inches(0.81),kw2-Inches(0.06),Inches(0.22),lbl,size=6.5,fg=MGREY,align=PP_ALIGN.CENTER)
        txt(sl,ki*kw2+Inches(0.03),Inches(1.02),kw2-Inches(0.05),Inches(0.30),val,size=12,bold=True,fg=WHITE,align=PP_ALIGN.CENTER)
        txt(sl,ki*kw2+Inches(0.03),Inches(1.31),kw2-Inches(0.05),Inches(0.20),sub_,size=6.5,fg=LGREY,align=PP_ALIGN.CENTER,italic=True)

    # ranked table
    sec_bar(sl,Inches(0.15),Inches(1.64),Inches(7.9),
            f"Station Volume Ranking — top {min(28,n_tot)} shown  |  Blue = within 80% threshold  |  Red = inactive")
    show=min(28,n_tot)
    rt=sl.shapes.add_table(show+1,6,Inches(0.15),Inches(1.93),Inches(7.9),Inches(5.0)).table
    for ci,h in enumerate(["#","Station Name","Vol CY (ML)","Vol LY (ML)","Chg%","Cum%"]):
        cstyle(rt.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=NAVY)
    rt.columns[0].width=Inches(0.3); rt.columns[1].width=Inches(2.85)
    for ci in range(2,6): rt.columns[ci].width=Inches(1.19)
    rt.rows[0].height=Inches(0.25)
    reached=False
    for ri,(idx,row) in enumerate(stn.head(show).iterrows(),1):
        crosses=(row["cumsh"]>80 and not reached)
        if crosses: reached=True
        is_key=row["cumsh"]<=80 or crosses
        bg=LBLUE if is_key and row["active"] else (LRED if not row["active"] else WHITE)
        vals=[str(ri),str(row["name"])[:38],fv(row["ml_cy"]),fv(row["ml_ly"]),
              fp(row["chg"]),"—" if not row["active"] else f"{row['cumsh']:.1f}%"]
        for ci,v in enumerate(vals):
            cstyle(rt.cell(ri,ci),v,size=7,bold=(ri<=5),
                   fg=RED if not row["active"] else NAVY,
                   bg=chg_bg(row["chg"]) if ci==4 and row["active"] else bg,
                   align=PP_ALIGN.LEFT if ci==1 else PP_ALIGN.CENTER)
        rt.rows[ri].height=Inches(0.17)

    # quintile panel (right)
    sec_bar(sl,Inches(8.2),Inches(1.64),Inches(4.95),"Volume by Station Tier (quintiles)")
    q_n=max(1,n_tot//5)
    qtiers=[("Top 20%",stn.head(q_n)),("21–40%",stn.iloc[q_n:2*q_n]),
            ("41–60%",stn.iloc[2*q_n:3*q_n]),("61–80%",stn.iloc[3*q_n:4*q_n]),
            ("Bottom 20%",stn.iloc[4*q_n:])]
    qt=sl.shapes.add_table(6,4,Inches(8.2),Inches(1.93),Inches(4.95),Inches(1.55)).table
    for ci,h in enumerate(["Tier","Count","Vol (ML)","Share"]):
        cstyle(qt.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
    qt.rows[0].height=Inches(0.24)
    qbg=[NAVY,DBLUE,LBLUE,LGREY,WHITE]; qfg=[WHITE,WHITE,NAVY,NAVY,NAVY]
    for ri,(lbl,g) in enumerate(qtiers,1):
        gv=g["ml_cy"].sum(); gs=gv/d["tot_cy"]*100 if d["tot_cy"] else 0
        cstyle(qt.cell(ri,0),lbl,size=8,bold=True,fg=qfg[ri-1],bg=qbg[ri-1],align=PP_ALIGN.LEFT)
        cstyle(qt.cell(ri,1),str(len(g)),size=8,fg=qfg[ri-1],bg=qbg[ri-1])
        cstyle(qt.cell(ri,2),fv(gv),size=8,bold=True,fg=qfg[ri-1],bg=qbg[ri-1])
        cstyle(qt.cell(ri,3),f"{gs:.1f}%",size=8,bold=True,fg=qfg[ri-1],bg=qbg[ri-1])
        qt.rows[ri].height=Inches(0.24)
    for ci,w in enumerate([Inches(1.2),Inches(0.7),Inches(1.5),Inches(1.55)]): qt.columns[ci].width=w

    bot_avg=stn.tail(n_tot-n80)["ml_cy"].mean()
    top_avg=stn.head(n80)["ml_cy"].mean()
    insight(sl,Inches(8.2),Inches(3.65),Inches(4.95),Inches(3.3),
        "WHAT TO DO WITH THIS",
        (f"The top {n80} stations are high-value accounts. "
         f"They average {top_avg:.2f}ML each — they need relationship management, not activation.\n\n"
         f"The bottom {n_tot-n80} stations average only {bot_avg:.2f}ML each. "
         f"A 30% lift across this group adds +{stn.tail(n_tot-n80)['ml_cy'].sum()*0.3:.1f}ML "
         f"to {city}'s total — achievable through product additions and re-engagement.\n\n"
         f"RISK: Top 5 stations = {stn.head(5)['sh'].sum():.0f}% of city volume. "
         f"Loss of even one is a city-level event. These stations need formal retention agreements."))

    footer(sl,f"Blue = stations within 80% vol threshold  |  Red = inactive (zero CY volume)  |  "
              f"Top:bottom tier volume ratio = {top_avg/max(bot_avg,0.01):.0f}x")
    set_notes(sl,
        f"SLIDE 2 — {city.upper()} 80% RULE\n\n"
        f"PRESENTER TALKING POINTS:\n"
        f"The key message is: {n80} stations ({n80/n_tot*100:.0f}% of the fleet) are doing the heavy lifting. "
        f"Ask the audience: 'Does your field team know these {n80} stations by name?'\n\n"
        f"NUMBERS TO EMPHASISE:\n"
        f"• {n80} stations → 80% of {city} volume\n"
        f"• Top station: '{stn.iloc[0]['name']}' at {stn.iloc[0]['ml_cy']:.1f}ML "
        f"({stn.iloc[0]['sh']:.1f}% of city)\n"
        f"• Top 5 stations combined: {stn.head(5)['sh'].sum():.1f}% of city — extreme concentration\n"
        f"• Top tier avg: {top_avg:.2f}ML per station vs bottom tier {bot_avg:.2f}ML — "
        f"{top_avg/max(bot_avg,0.01):.0f}x ratio\n"
        f"• Bottom {n_tot-n80} stations: +{stn.tail(n_tot-n80)['ml_cy'].sum()*0.3:.1f}ML potential "
        f"at 30% activation lift\n\n"
        f"QUESTIONS YOU MIGHT GET:\n"
        f"• 'Why are bottom-tier stations so low?' — slide 3 will show that single-product stations "
        f"average far less volume than multi-product ones. Product mix is the answer.\n"
        f"• 'Who are the inactive stations?' — covered in Slide 6.\n\n"
        f"QUINTILE BREAKDOWN:\n"
        f"  Top 20% ({qtiers[0][1]['ml_cy'].sum():.1f}ML) → "
        f"  21-40% ({qtiers[1][1]['ml_cy'].sum():.1f}ML) → "
        f"  41-60% ({qtiers[2][1]['ml_cy'].sum():.1f}ML) → "
        f"  61-80% ({qtiers[3][1]['ml_cy'].sum():.1f}ML) → "
        f"  Bottom 20% ({qtiers[4][1]['ml_cy'].sum():.1f}ML)\n"
    )

# ── SLIDE 3: STAR PERFORMERS ──────────────────────────────────────────────────
def make_s3(prs, city, d, period):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    stn=d["stn"]; pm=d["pm"]; pf=d["prod_str_fn"]
    top25n=max(1,d["n_tot"]//4)
    hdr(sl,f"STAR PERFORMERS — WHAT GOOD LOOKS LIKE IN {city.upper()}",
        f"Top 25% of stations ({top25n} outlets) average {d['top25_avg']:.2f}ML each. "
        f"Multi-product stations average {d['mul_avg']:.1f}x more volume than single-product ones.")

    # top 15 table
    sec_bar(sl,Inches(0.15),Inches(0.85),Inches(8.1),"Top 15 Stations — Volume Champions (products = actual CY sales)")
    stars=stn.head(15).copy()
    st=sl.shapes.add_table(16,7,Inches(0.15),Inches(1.13),Inches(8.1),Inches(5.85)).table
    for ci,h in enumerate(["#","Station","Vol CY (ML)","Vol LY (ML)","YoY Chg","Products Sold","Cum%"]):
        cstyle(st.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=NAVY)
    st.columns[0].width=Inches(0.3); st.columns[1].width=Inches(2.85)
    st.columns[2].width=Inches(1.1); st.columns[3].width=Inches(1.1)
    st.columns[4].width=Inches(0.9); st.columns[5].width=Inches(1.15); st.columns[6].width=Inches(0.7)
    st.rows[0].height=Inches(0.28)
    for ri,(idx,row) in enumerate(stars.iterrows(),1):
        sp=pf(row["Customer Number"])
        bg=LGREEN if (row["chg"] or 0)>5 else (LYELL if (row["chg"] or 0)>0 else (LRED if (row["chg"] or 0)<-5 else LGREY))
        for ci,v in enumerate([str(ri),str(row["name"])[:38],fv(row["ml_cy"]),
                                fv(row["ml_ly"]),fp(row["chg"]),sp,f"{row['cumsh']:.0f}%"]):
            cstyle(st.cell(ri,ci),v,size=7,bold=(ri<=5),fg=NAVY,
                   bg=chg_bg(row["chg"]) if ci==4 else bg,
                   align=PP_ALIGN.LEFT if ci in(1,5) else PP_ALIGN.CENTER)
        st.rows[ri].height=Inches(0.37)

    # product breadth table
    sec_bar(sl,Inches(8.3),Inches(0.85),Inches(4.85),"Product Breadth vs Volume — the multiplier effect")
    pv=d["pv"]
    pv_rows=[(k,np.mean(v),len(v)) for k,v in pv.items() if v]
    pvt=sl.shapes.add_table(len(pv_rows)+1,3,Inches(8.3),Inches(1.13),Inches(4.85),Inches(1.3)).table
    for ci,h in enumerate(["Product Groups","Avg Vol/Station","# Stations"]):
        cstyle(pvt.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY,align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
    pvt.columns[0].width=Inches(2.1); pvt.columns[1].width=Inches(1.5); pvt.columns[2].width=Inches(1.25)
    pvt.rows[0].height=Inches(0.26)
    pv_bgs=[LRED,LYELL,LGREEN]
    for ri,(k,avg,cnt) in enumerate(pv_rows,1):
        bg=pv_bgs[ri-1] if ri<=3 else LGREY
        cstyle(pvt.cell(ri,0),k,size=8.5,bold=True,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
        cstyle(pvt.cell(ri,1),f"{avg:.2f}ML",size=9,bold=True,fg=NAVY,bg=bg)
        cstyle(pvt.cell(ri,2),str(cnt),size=8.5,fg=NAVY,bg=bg)
        pvt.rows[ri].height=Inches(0.32)

    mul=d["mul_avg"]; one=d["one_avg"]
    insight(sl,Inches(8.3),Inches(2.58),Inches(4.85),Inches(1.85),
        "THE MULTIPLIER FINDING",
        (f"Stations with 3+ product groups average {mul:.2f}ML — "
         f"that is {mul/max(one,0.001):.1f}x more than single-product stations "
         f"at {one:.2f}ML.\n\n"
         f"Every time you add a product group to a station, volume goes up. "
         f"This is your fastest organic growth lever — no new outlets needed."))

    # benchmarks
    sec_bar(sl,Inches(8.3),Inches(4.6),Inches(4.85),"Volume Benchmarks — Active Stations")
    bm=[(f"Top 10% avg ({max(1,d['n_tot']//10)} stns)",
         f"{stn.head(max(1,d['n_tot']//10))['ml_cy'].mean():.2f}ML","Star target"),
        (f"Top 25% avg ({top25n} stns)",f"{d['top25_avg']:.2f}ML","Aspiration"),
        ("Median active station",f"{d['median_vol']:.2f}ML","Mid-tier baseline"),
        ("10-city fleet average",f"{t10_avg_per_stn:.2f}ML","Peer benchmark"),
        (f"Bottom 25% avg ({max(1,d['n_tot']//4)} stns)",
         f"{stn.tail(max(1,d['n_tot']//4))['ml_cy'].mean():.2f}ML","Needs activation")]
    bmt=sl.shapes.add_table(6,3,Inches(8.3),Inches(4.88),Inches(4.85),Inches(1.9)).table
    for ci,h in enumerate(["Benchmark","Vol","Context"]):
        cstyle(bmt.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
    bmt.columns[0].width=Inches(2.1); bmt.columns[1].width=Inches(1.3); bmt.columns[2].width=Inches(1.45)
    bmt.rows[0].height=Inches(0.25)
    bm_bgs=[LGREEN,LGREEN,LGREY,LYELL,LRED]
    for ri,(lbl,val,ctx) in enumerate(bm,1):
        cstyle(bmt.cell(ri,0),lbl,size=8,fg=NAVY,bg=bm_bgs[ri-1],align=PP_ALIGN.LEFT)
        cstyle(bmt.cell(ri,1),val,size=9,bold=True,fg=NAVY,bg=bm_bgs[ri-1])
        cstyle(bmt.cell(ri,2),ctx,size=7.5,fg=NAVY,bg=bm_bgs[ri-1],italic=True,align=PP_ALIGN.LEFT)
        bmt.rows[ri].height=Inches(0.31)

    footer(sl,"Green = CY > LY (growing)  |  Red = CY < LY (declining)  |  Products = actual sold in CY, from Working File")
    set_notes(sl,
        f"SLIDE 3 — {city.upper()} STAR PERFORMERS\n\n"
        f"PRESENTER TALKING POINTS:\n"
        f"This slide answers: 'What does a great station look like in {city}?' "
        f"and 'Why do some stations dramatically outperform others?'\n\n"
        f"THE CORE FINDING:\n"
        f"Stations with 3 or more product groups average {mul:.2f}ML, "
        f"versus {one:.2f}ML for single-product stations. "
        f"That is a {mul/max(one,0.001):.1f}x difference — and the only variable is product breadth.\n"
        f"This is not a market effect, it is an operational choice. "
        f"Field teams can change this.\n\n"
        f"TOP 5 STATIONS:\n"
        +"\n".join(f"  #{i+1}: {r['name']} — {r['ml_cy']:.2f}ML ({fp(r['chg'])} YoY) — Products: {pf(r['Customer Number'])}"
                   for i,(_,r) in enumerate(stn.head(5).iterrows()))
        +f"\n\nBENCHMARK CONVERSATION:\n"
        f"Ask: 'If all {d['n_active']} active stations hit the median benchmark of {d['median_vol']:.2f}ML, "
        f"what would total city volume be?' Answer: {d['median_vol']*d['n_active']:.1f}ML "
        f"vs current {d['tot_cy']:.1f}ML. "
        f"That gap is the volume locked in under-performing stations.\n\n"
        f"QUESTIONS YOU MIGHT GET:\n"
        f"• 'Can we replicate what the top stations do?' — Yes. The key is product breadth. "
        f"See Slide 4 for the specific gaps.\n"
        f"• 'Which stations are declining despite being big?' — Check the red-highlighted rows.\n"
    )

# ── SLIDE 4: PRODUCT GAPS ─────────────────────────────────────────────────────
def make_s4(prs, city, d, period):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    hdr(sl,f"PRODUCT GAPS — VOLUME SITTING ON THE TABLE IN {city.upper()}",
        f"{d['n_fl']} stations sell fuel but no lubricants. "
        f"{d['n_d_no_p']} sell diesel but no petrol. Total identifiable vol gap: +{fv(d['tot_pot'])}ML.")

    # coverage breakdown
    sec_bar(sl,Inches(0.15),Inches(0.85),Inches(5.9),"Station Product Coverage — What Each Station Sells")
    cats=[("Diesel + Petrol + Lubes",d["n_all3"],LGREEN,"Complete — maximum vol potential"),
          ("Diesel + Petrol only",   d["n_d_p"], LYELL, "Missing lubes — gap to close"),
          ("Diesel only",            d["n_d_only"],LRED,"Missing petrol & lubes — largest gap"),
          ("Petrol only",            d["n_p_only"],LORANG,"Missing diesel & lubes"),
          ("Inactive (zero vol)",    d["n_inact"],LRED, "No sales — immediate action needed")]
    ct=sl.shapes.add_table(6,4,Inches(0.15),Inches(1.13),Inches(5.9),Inches(2.0)).table
    for ci,h in enumerate(["Category","Count","Fleet%","Status"]):
        cstyle(ct.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
    ct.columns[0].width=Inches(2.55); ct.columns[1].width=Inches(0.75)
    ct.columns[2].width=Inches(0.95); ct.columns[3].width=Inches(1.65)
    ct.rows[0].height=Inches(0.26)
    for ri,(lbl,cnt,bg,st) in enumerate(cats,1):
        cstyle(ct.cell(ri,0),lbl,size=8.5,bold=True,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
        cstyle(ct.cell(ri,1),str(cnt),size=9,bold=True,fg=NAVY,bg=bg)
        cstyle(ct.cell(ri,2),f"{cnt/d['n_tot']*100:.0f}%",size=9,bold=True,fg=NAVY,bg=bg)
        cstyle(ct.cell(ri,3),st,size=7.5,fg=NAVY,bg=bg,italic=True,align=PP_ALIGN.LEFT)
        ct.rows[ri].height=Inches(0.33)

    # product penetration
    sec_bar(sl,Inches(0.15),Inches(3.27),Inches(5.9),"Individual Product Penetration — Stations Selling Each Product")
    pens=[("HSD",d["n_diesel"],"Core diesel product"),
          ("PMG (petrol)",d["n_petrol"],"Standard petrol"),
          ("R95 (hi-octane)",d["n_r95"],"Premium — low penetration = opportunity"),
          ("Any Lubricant",d["n_lubes"],"Highest vol multiplier"),
          ("DEO",d["sub"]["DEO"]["stns"],"Commercial engine oil"),
          ("PCMO",d["sub"]["PCMO"]["stns"],"Passenger car oil"),
          ("MCO",d["sub"]["MCO"]["stns"],"Motorcycle oil — high pull"),
          ("LOW GRADE",d["sub"]["LOW GRADE"]["stns"] if "LOW GRADE" in d["sub"] else 0,
           "Target shift to DEO")]
    pt=sl.shapes.add_table(9,4,Inches(0.15),Inches(3.56),Inches(5.9),Inches(3.3)).table
    for ci,h in enumerate(["Product","Stns Selling","% Fleet","Note"]):
        cstyle(pt.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
    pt.columns[0].width=Inches(1.8); pt.columns[1].width=Inches(1.0)
    pt.columns[2].width=Inches(1.0); pt.columns[3].width=Inches(2.1)
    pt.rows[0].height=Inches(0.24)
    for ri,(lbl,cnt,note) in enumerate(pens,1):
        pct=cnt/d["n_tot"]*100
        bg=LGREEN if pct>70 else (LYELL if pct>30 else LRED)
        cstyle(pt.cell(ri,0),lbl,size=8,bold=True,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
        cstyle(pt.cell(ri,1),str(cnt),size=9,bold=True,fg=NAVY,bg=bg)
        cstyle(pt.cell(ri,2),f"{pct:.0f}%",size=9,bold=True,fg=NAVY,bg=bg)
        cstyle(pt.cell(ri,3),note,size=7.5,fg=NAVY,bg=bg,italic=True,align=PP_ALIGN.LEFT)
        pt.rows[ri].height=Inches(0.36)

    # right: gap quantification
    sec_bar(sl,Inches(6.1),Inches(0.85),Inches(7.05),"Volume Opportunity — Closing Each Gap (directional estimates)")
    gaps=[
        ("Diesel stations with no Petrol",f"{d['n_d_no_p']} stations",
         f"Base diesel vol: {fv(d['dnp_vol'])}ML",f"+{fv(d['rec_petrol'])}ML",
         "Adding PMG/R95 at 25% of their diesel vol",LRED),
        ("Fuel stations with no Lubricants",f"{d['n_fl']} stations",
         f"Fuel vol at these stns: {fv(d['fl_vol'])}ML",f"+{fv(d['rec_lub'])}ML",
         f"Avg lube vol per existing lube stn: {d['avg_lub_per_stn']*1000:.0f} litres",LRED),
        ("Inactive stations — recover LY vol",f"{d['n_inact']} stations",
         f"Their LY volume: {fv(d['inact_ly_vol'])}ML",f"+{fv(d['rec_inact'])}ML",
         "65% reactivation assumption — field programme",LRED),
        ("Under-performers < 50% of own LY",f"{d['n_undp']} stations",
         f"Current vol gap vs LY",f"+{fv(d['rec_undp'])}ML",
         "Diagnostic visit + targeted support to reach 75% of LY",LYELL),
    ]
    gh=Inches(6.15)/len(gaps)
    for gi,(title,count,base,potential,detail,bg) in enumerate(gaps):
        y=Inches(1.13)+gi*gh
        rect(sl,Inches(6.1),y,Inches(7.05),gh-Inches(0.06),bg)
        rect(sl,Inches(6.1),y,Inches(0.22),gh-Inches(0.06),RED if bg==LRED else YELL)
        txt(sl,Inches(6.4),y+Inches(0.04),Inches(3.5),Inches(0.24),title,size=9,bold=True,fg=NAVY)
        txt(sl,Inches(6.4),y+Inches(0.27),Inches(2.5),Inches(0.2),count,size=8.5,fg=NAVY)
        txt(sl,Inches(6.4),y+Inches(0.46),Inches(4.5),Inches(0.2),detail,size=7.5,fg=NAVY,italic=True)
        txt(sl,Inches(10.1),y+Inches(0.1),Inches(2.9),Inches(0.5),potential,size=12,bold=True,
            fg=RED if bg==LRED else GREEN,align=PP_ALIGN.CENTER)
        txt(sl,Inches(10.1),y+Inches(0.55),Inches(2.9),Inches(0.2),base,size=7.5,
            fg=NAVY,align=PP_ALIGN.CENTER,italic=True)

    tot_p=d["tot_pot"]
    rect(sl,Inches(6.1),Inches(7.02),Inches(7.05),Inches(0.13),GOLD)
    txt(sl,Inches(6.2),Inches(7.03),Inches(6.9),Inches(0.12),
        f"TOTAL IDENTIFIABLE VOLUME OPPORTUNITY: +{fv(tot_p)}ML  "
        f"(+{tot_p/d['tot_cy']*100:.0f}% on current base of {fv(d['tot_cy'])}ML)",
        size=8.5,bold=True,fg=NAVY)

    footer(sl,"Potential estimates are directional. Lube: avg vol per existing lube stn. "
              "Petrol: 25% of diesel base vol. Inactive: 65% of LY vol assumed recoverable.")
    set_notes(sl,
        f"SLIDE 4 — {city.upper()} PRODUCT GAPS\n\n"
        f"PRESENTER TALKING POINTS:\n"
        f"This slide makes the volume opportunity concrete. "
        f"There are {d['n_fl']} stations in {city} that sell fuel but no lubricants. "
        f"That is {d['n_fl']/d['n_tot']*100:.0f}% of the fleet leaving lube revenue uncaptured.\n\n"
        f"STATION COVERAGE NUMBERS:\n"
        f"• Complete (D+P+Lubes): {d['n_all3']} stations ({d['n_all3']/d['n_tot']*100:.0f}%)\n"
        f"• Diesel + Petrol only: {d['n_d_p']} stations ({d['n_d_p']/d['n_tot']*100:.0f}%)\n"
        f"• Diesel only: {d['n_d_only']} stations ({d['n_d_only']/d['n_tot']*100:.0f}%)\n"
        f"• Inactive: {d['n_inact']} stations ({d['n_inact']/d['n_tot']*100:.0f}%)\n\n"
        f"VOLUME OPPORTUNITIES:\n"
        f"• Petrol at diesel-only stns: +{fv(d['rec_petrol'])}ML "
        f"({d['n_d_no_p']} stns × 25% of their diesel vol)\n"
        f"• Lubes at fuel-only stns: +{fv(d['rec_lub'])}ML "
        f"({d['n_fl']} stns × avg lube vol {d['avg_lub_per_stn']*1000:.0f}L/stn)\n"
        f"• Reactivate inactive: +{fv(d['rec_inact'])}ML\n"
        f"• Fix under-performers: +{fv(d['rec_undp'])}ML\n"
        f"• TOTAL: +{fv(d['tot_pot'])}ML\n\n"
        f"QUESTIONS YOU MIGHT GET:\n"
        f"• 'Why don't all stations sell all products?' — Storage, capex, or dealer preference. "
        f"Field team needs to screen each station for barriers.\n"
        f"• 'Is the lube potential realistic?' — Based on actual average vol per existing lube station. "
        f"Not an aggressive assumption.\n"
        f"• 'Which product to prioritise?' — MCO (motorcycle oil) is highest consumer pull "
        f"and lowest capex to introduce. Start there.\n"
    )

# ── SLIDE 5: SUB-PRODUCT MIX ──────────────────────────────────────────────────
def make_s5(prs, city, d, period):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    hdr(sl,f"SUB-PRODUCT VOLUME — WHAT IS GROWING, WHAT IS DECLINING IN {city.upper()}",
        "Every product within Diesel, Petrol and Lubricants. CY vs LY. Which products are losing ground and need attention?")

    active=[(k,v) for k,v in d["sub"].items() if v["cy"]>0 or v["ly"]>0]
    active.sort(key=lambda x: x[1]["cy"],reverse=True)
    grp_bgs={"Diesel":LBLUE,"Petrol":LGREEN,"Lubricants":LORANG}
    grp_tot={"Diesel":d["segs"]["Diesel"]["cy"],"Petrol":d["segs"]["Petrol"]["cy"],
             "Lubricants":d["segs"]["Lubricants"]["cy"]}

    sec_bar(sl,Inches(0.15),Inches(0.85),Inches(12.97),
            "Volume by Sub-Product — CY vs LY  |  Only products with actual sales shown  |  "
            "Red rows = had LY volume, now zero — urgent recovery needed")
    nr=len(active)+1
    t=sl.shapes.add_table(nr,8,Inches(0.15),Inches(1.13),Inches(12.97),
                           min(Inches(5.85),Inches(0.3)*nr)).table
    hdrs=["Product","Group","Stns Selling","Vol CY (ML)","Vol LY (ML)","Chg%","Seg Share%","Status"]
    cws=[Inches(1.8),Inches(1.1),Inches(0.9),Inches(1.4),Inches(1.4),Inches(1.0),Inches(1.65),Inches(3.67)]
    for ci,(h,w) in enumerate(zip(hdrs,cws)):
        cstyle(t.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
        t.columns[ci].width=w
    t.rows[0].height=Inches(0.26)
    last_grp=None
    for ri,(key,sv) in enumerate(active,1):
        is_new=sv["grp"]!=last_grp; last_grp=sv["grp"]
        bg=grp_bgs.get(sv["grp"],LGREY)
        ss=sv["cy"]/grp_tot.get(sv["grp"],1)*100 if grp_tot.get(sv["grp"],0)>0 else 0
        if sv["cy"]==0 and sv["ly"]>0: status="⚠ DROPPED — selling LY, zero CY"; sbg=LRED
        elif (sv["chg"] or 0)>5:       status=f"Growing {fp(sv['chg'])} YoY"; sbg=LGREEN
        elif (sv["chg"] or 0)<-5:      status=f"Declining {fp(sv['chg'])} — investigate cause"; sbg=LRED
        elif sv["cy"]==0:               status="Not sold here"; sbg=MGREY
        else:                           status="Stable (within ±5%)"; sbg=LYELL
        row=[key,sv["grp"],str(sv["stns"]) if sv["stns"] else "—",
             fv(sv["cy"]),fv(sv["ly"]),fp(sv["chg"]),f"{ss:.1f}%",status]
        for ci,v in enumerate(row):
            cstyle(t.cell(ri,ci),v,size=8,bold=(is_new and ci==0),fg=NAVY,
                   bg=chg_bg(sv["chg"]) if ci==5 else (sbg if ci==7 else bg),
                   align=PP_ALIGN.LEFT if ci in(0,1,7) else PP_ALIGN.CENTER)
        t.rows[ri].height=Inches(0.28)

    dropped=d["dropped"]
    if dropped:
        foot_note=f"DROPPED products (CY=0, had LY volume): {', '.join(dropped)} — these represent recoverable volume."
    else:
        foot_note="No products dropped since LY — all previously sold products are still active in CY."
    footer(sl,foot_note)

    set_notes(sl,
        f"SLIDE 5 — {city.upper()} SUB-PRODUCT DETAIL\n\n"
        f"PRESENTER TALKING POINTS:\n"
        f"This slide goes inside each segment to show which specific products are driving "
        f"or dragging performance.\n\n"
        f"KEY NUMBERS:\n"
        f"  HSD: {fv(d['sub']['HSD']['cy'])}ML CY vs {fv(d['sub']['HSD']['ly'])}ML LY ({fp(d['sub']['HSD']['chg'])})\n"
        f"  PMG: {fv(d['sub']['PMG']['cy'])}ML CY vs {fv(d['sub']['PMG']['ly'])}ML LY ({fp(d['sub']['PMG']['chg'])})\n"
        f"  R95: {fv(d['sub']['R95']['cy'])}ML CY vs {fv(d['sub']['R95']['ly'])}ML LY ({fp(d['sub']['R95']['chg'])})\n"
        f"  DEO: {fv(d['sub']['DEO']['cy'])}ML  PCMO: {fv(d['sub']['PCMO']['cy'])}ML  "
        f"MCO: {fv(d['sub']['MCO']['cy'])}ML  LOW GRADE: {fv(d['sub'].get('LOW GRADE',{}).get('cy',0))}ML\n\n"
        f"PRODUCTS TO CALL OUT:\n"
        +(f"• DROPPED SINCE LY: {', '.join(dropped)} — supply chain or market exit. "
           f"Needs immediate investigation.\n" if dropped else "• No products dropped since LY.\n")
        +f"\nDECLINING PRODUCTS (CY < LY by >5%):\n"
        +"\n".join(f"  • {k}: {fp(v['chg'])} ({fv(v['cy'])}ML CY vs {fv(v['ly'])}ML LY)"
                   for k,v in d["sub"].items() if (v["chg"] or 0)<-5 and v["cy"]>0)
        +f"\n\nQUESTIONS YOU MIGHT GET:\n"
        f"• 'Why is HSD flat/declining?' — Check if competitor pricing is more aggressive "
        f"on commercial corridors in this city. Also check if inactive stations were heavy HSD sellers.\n"
        f"• 'Why is LOW GRADE still being sold?' — It should be migrated to DEO. "
        f"LOW GRADE has lowest margin and should be phased out.\n"
        f"• 'Can we grow R95?' — R95 has highest margin. Every PMG station that adds R95 "
        f"is a no-capex volume and margin win.\n"
    )

# ── SLIDE 6: LOST VOLUME ──────────────────────────────────────────────────────
def make_s6(prs, city, d, period):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    stn=d["stn"]; inact=stn[~stn["active"]]; undp=d["undp_df"]
    hdr(sl,f"LOST VOLUME — {d['n_inact']} INACTIVE + {d['n_undp']} UNDER-PERFORMING STATIONS IN {city.upper()}",
        f"Inactive stations carried {fv(d['inact_ly_vol'])}ML in LY and now show zero. "
        f"Under-performers lost {fv((undp['ml_ly']-undp['ml_cy']).sum())}ML vs their own LY. "
        f"Combined recoverable: ~{fv(d['rec_inact']+d['rec_undp'])}ML.")

    # inactive (left)
    sec_bar(sl,Inches(0.15),Inches(0.85),Inches(6.45),
            f"Inactive Stations ({d['n_inact']}) — Zero CY Volume, Had Volume in LY  |  Sorted by LY vol desc")
    if not inact.empty:
        show=min(len(inact),15)
        it=sl.shapes.add_table(show+1,5,Inches(0.15),Inches(1.13),Inches(6.45),Inches(3.65)).table
        for ci,h in enumerate(["Station Name","LY Vol (ML)","Products (LY)","LY Rank","Priority"]):
            cstyle(it.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=RED)
        it.columns[0].width=Inches(2.4); it.columns[1].width=Inches(1.1)
        it.columns[2].width=Inches(1.4); it.columns[3].width=Inches(0.8); it.columns[4].width=Inches(0.75)
        it.rows[0].height=Inches(0.24)
        pf=d["prod_str_fn"]
        for ri,(idx,row) in enumerate(inact.sort_values("ml_ly",ascending=False).head(show).iterrows(),1):
            ly_r=stn.index[stn["Customer Number"]==row["Customer Number"]].tolist()
            pr="HIGH" if row["ml_ly"]>1 else "MED" if row["ml_ly"]>0.3 else "LOW"
            for ci,v in enumerate([str(row["name"])[:34],fv(row["ml_ly"]),
                                    pf(row["Customer Number"]),
                                    str(ly_r[0]+1) if ly_r else "—",pr]):
                cstyle(it.cell(ri,ci),v,size=7,fg=RED if ci==4 else NAVY,bg=LRED,
                       align=PP_ALIGN.LEFT if ci in(0,2) else PP_ALIGN.CENTER,bold=(ci==4))
            it.rows[ri].height=Inches(0.215)
    else:
        txt(sl,Inches(0.2),Inches(1.2),Inches(6.3),Inches(0.4),
            "No inactive stations — all stations in this city have CY volume.",
            size=9,fg=GREEN,bold=True)

    # under-performers (left bottom)
    sec_bar(sl,Inches(0.15),Inches(4.94),Inches(6.45),
            f"Severely Under-Performing ({d['n_undp']}) — CY Volume < 50% of Their Own LY Baseline")
    if not undp.empty:
        show2=min(len(undp),7)
        ut=sl.shapes.add_table(show2+1,5,Inches(0.15),Inches(5.22),Inches(6.45),Inches(1.65)).table
        for ci,h in enumerate(["Station","Vol CY","Vol LY","Vol Lost","Drop%"]):
            cstyle(ut.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=RED)
        for ci,w in enumerate([Inches(2.4),Inches(1.0),Inches(1.0),Inches(1.0),Inches(1.05)]): ut.columns[ci].width=w
        ut.rows[0].height=Inches(0.22)
        for ri,(idx,row) in enumerate(undp.sort_values("ml_ly",ascending=False).head(show2).iterrows(),1):
            loss=row["ml_ly"]-row["ml_cy"]
            for ci,v in enumerate([str(row["name"])[:34],fv(row["ml_cy"]),fv(row["ml_ly"]),f"-{loss:.2f}",fp(row["chg"])]):
                cstyle(ut.cell(ri,ci),v,size=7,bold=(ci==3),
                       fg=RED if ci in(3,4) else NAVY,bg=LYELL,
                       align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
            ut.rows[ri].height=Inches(0.20)
    else:
        txt(sl,Inches(0.2),Inches(5.3),Inches(6.3),Inches(0.4),
            "No stations below 50% of their own LY baseline — under-performance is distributed.",
            size=9,fg=GREEN)

    # right panel
    insight(sl,Inches(6.72),Inches(0.85),Inches(6.43),Inches(2.0),
        "INACTIVE STATIONS — IMMEDIATE ACTION",
        (f"  {d['n_inact']} stations had {fv(d['inact_ly_vol'])}ML in LY. Now zero.\n"
         f"  At 65% reactivation: +{fv(d['rec_inact'])}ML back.\n\n"
         f"  CLASSIFY EACH STATION (field visit within 2 weeks):\n"
         f"  • Temp closure (supply, credit, maintenance) → reactivate now\n"
         f"  • Dealer exit → find replacement or surrender licence\n"
         f"  • Competitive loss → pricing & relationship intervention\n"
         f"  HIGH-PRIORITY: stations with LY vol > 1ML — each matters individually."),
        tbg=RED,bbg=LRED,tfg=WHITE,bfg=NAVY)

    insight(sl,Inches(6.72),Inches(3.0),Inches(6.43),Inches(1.8),
        "UNDER-PERFORMERS — STATION DIAGNOSTIC",
        (f"  {d['n_undp']} stations fell below 50% of their own previous-year volume.\n"
         f"  This is NOT market decline — it is station-specific failure.\n"
         f"  Root causes: equipment failure, credit limits, competitor opened nearby,\n"
         f"  management change at dealer level, product supply disruption.\n"
         f"  At 75% of LY target: +{fv(d['rec_undp'])}ML recoverable."),
        tbg=ORANGE,bbg=LYELL,tfg=NAVY,bfg=NAVY)

    insight(sl,Inches(6.72),Inches(4.95),Inches(6.43),Inches(1.93),
        "COMBINED RECOVERY SUMMARY",
        (f"  Inactive reactivation:      +{fv(d['rec_inact'])}ML\n"
         f"  Under-performer recovery:   +{fv(d['rec_undp'])}ML\n"
         f"  ─────────────────────────────────────────\n"
         f"  TOTAL RECOVERY TARGET:      +{fv(d['rec_inact']+d['rec_undp'])}ML\n\n"
         f"  This is {(d['rec_inact']+d['rec_undp'])/d['tot_cy']*100:.0f}% of {city}'s current CY volume. "
         f"All from existing outlets. No new investment needed.\n"
         f"  Timeline: 60–90 days with dedicated field programme."),
        tbg=GREEN,bbg=LGREEN,tfg=WHITE,bfg=NAVY)

    footer(sl,"Under-performing = active stations where CY vol < 50% of their own LY. "
              "Recovery figures are directional — validate with field team.")
    set_notes(sl,
        f"SLIDE 6 — {city.upper()} LOST VOLUME\n\n"
        f"PRESENTER TALKING POINTS:\n"
        f"This is the most urgent slide in the deck. These are stations that used to contribute "
        f"and have gone quiet. The volume loss is not a market problem — it is a "
        f"station-management problem.\n\n"
        f"INACTIVE STATIONS — THE FACTS:\n"
        f"• {d['n_inact']} stations have zero CY volume\n"
        f"• Their combined LY volume was {fv(d['inact_ly_vol'])}ML\n"
        f"• At 65% reactivation: +{fv(d['rec_inact'])}ML\n"
        f"• Top priority: stations that had > 1ML in LY\n\n"
        f"UNDER-PERFORMERS — THE FACTS:\n"
        f"• {d['n_undp']} active stations are below 50% of their own LY\n"
        f"• Combined volume gap: {fv((undp['ml_ly']-undp['ml_cy']).sum() if not undp.empty else 0)}ML\n"
        f"• At 75% of LY recovery: +{fv(d['rec_undp'])}ML\n\n"
        f"HOW TO PRESENT:\n"
        f"Read out the name of the highest-LY inactive station. Ask: "
        f"'Does anyone in this room know why this station went to zero?' "
        f"That creates immediate accountability.\n\n"
        f"QUESTIONS YOU MIGHT GET:\n"
        f"• 'Is the inactive data accurate?' — Yes, from the Working File directly. "
        f"Zero SalesLtr_CY with non-zero SalesLtr_LY.\n"
        f"• 'What's the typical reason for going inactive?' — Most common: dealer credit frozen, "
        f"operational breakdown, or dealer switched to competitor on a better deal.\n"
        f"• 'Can we recover all of them?' — Field experience suggests 50–70% reactivation rate. "
        f"Some may be permanently closed and should be written off the active fleet.\n"
    )

# ── SLIDE 7: ROADMAP ──────────────────────────────────────────────────────────
def make_s7(prs, city, d, period):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    tot_p=d["tot_pot"]
    hdr(sl,f"{city.upper()} VOLUME GROWTH ROADMAP — +{fv(tot_p)}ML IDENTIFIABLE OPPORTUNITY",
        f"Four levers, prioritised by speed and certainty. All from existing outlets. "
        f"No new stations needed. {fv(d['tot_cy'])}ML today → {fv(d['tot_cy']+tot_p)}ML potential.")

    levers=[
        ("LEVER 1","Reactivate Inactive Stations","Fastest win — lost volume recovery",
         f"{d['n_inact']} stations",f"+{fv(d['rec_inact'])}ML","30–60 days",RED,LRED,
         [f"Field visit to all {d['n_inact']} inactive stations within 2 weeks",
          "Classify each: temp closure / dealer exit / credit issue / competitor loss",
          f"Stations with LY vol > 1ML are highest priority — each has material impact",
          "Target: reactivate 65% of fleet → recover 65% of lost {:.1f}ML".format(d['inact_ly_vol']),
          "Weekly progress tracker vs reactivation targets — escalate at week 3 if no movement"]),
        ("LEVER 2","Fix Under-Performing Stations","Station-specific failure — diagnostic required",
         f"{d['n_undp']} stations",f"+{fv(d['rec_undp'])}ML","60–90 days",ORANGE,LYELL,
         ["Station-by-station field visit — what specifically happened since LY?",
          "Compare CY vs LY by product: which product dropped? Was it supply or demand?",
          "Check: competitor proximity, credit limit hit, equipment failure, dealer change",
          f"Set 90-day volume recovery target at 75% of each station's LY baseline",
          "Weekly check-in with regional managers — escalate at 4-week mark if no recovery"]),
        ("LEVER 3","Add Lubricants to Fuel-Only Stations","Highest vol multiplier per addition",
         f"{d['n_fl']} eligible stns",f"+{fv(d['rec_lub'])}ML","90–120 days",DBLUE,LBLUE,
         [f"Prioritise: stations in top-80% GRS group that have zero lube sales today",
          f"Avg lube vol per existing lube seller: {d['avg_lub_per_stn']*1000:.0f} litres/station in 10M",
          "Start with MCO (motorcycle oil) — easiest consumer pull, lowest storage requirement",
          "Provide: starter inventory on credit, display stand, staff product training",
          f"Trial: 10 stations Month 1, full rollout at Month 3 if >60% conversion rate"]),
        ("LEVER 4","Add Petrol to Diesel-Only Stations","Infrastructure-dependent but high impact",
         f"{d['n_d_no_p']} eligible stns",f"+{fv(d['rec_petrol'])}ML","120–180 days",GREEN,LGREEN,
         [f"{d['n_d_no_p']} stations sell diesel but not petrol — urban ones are highest value",
          "Screen stations: does site have petrol storage? Is location on consumer corridor?",
          "Volume assumption: 25% of station's current diesel vol could convert to PMG/R95",
          "Prioritise urban/peri-urban locations over highway diesel-only outlets",
          "This lever requires capex — build the business case per station before committing"]),
    ]
    lh=(SH-Inches(1.0))/len(levers)
    for li,(lever,title,tag,scope,potential,timeline,hbg,bbg,actions) in enumerate(levers):
        y=Inches(0.85)+li*lh
        rect(sl,Inches(0),y,SW,Inches(0.32),hbg)
        txt(sl,Inches(0.18),y+Inches(0.04),Inches(1.3),Inches(0.25),lever,size=9,bold=True,fg=WHITE)
        txt(sl,Inches(1.6),y+Inches(0.04),Inches(6.5),Inches(0.25),title,size=10,bold=True,fg=WHITE)
        txt(sl,Inches(8.2),y+Inches(0.04),Inches(1.8),Inches(0.25),tag,size=7.5,fg=LGREY,italic=True)
        txt(sl,Inches(10.2),y+Inches(0.04),Inches(1.5),Inches(0.25),timeline,size=8,fg=GOLD,align=PP_ALIGN.CENTER)
        rect(sl,Inches(11.85),y,Inches(1.48),Inches(0.32),GOLD)
        txt(sl,Inches(11.87),y+Inches(0.03),Inches(1.44),Inches(0.25),potential,size=10,bold=True,fg=NAVY,align=PP_ALIGN.CENTER)
        rect(sl,Inches(0),y+Inches(0.32),SW,lh-Inches(0.34),bbg)
        n_act=len(actions)
        ah=(lh-Inches(0.38))/n_act
        for ai,action in enumerate(actions):
            ay=y+Inches(0.35)+ai*ah
            txt(sl,Inches(0.25),ay,Inches(13.0),ah,f"• {action}",size=8,fg=NAVY)

    rect(sl,Inches(0),SH-Inches(0.38),SW,Inches(0.38),NAVY)
    summary=(f"TOTAL: +{fv(tot_p)}ML  =  "
             f"Inactive {fv(d['rec_inact'])}ML  +  Under-perf {fv(d['rec_undp'])}ML  +  "
             f"Lubes {fv(d['rec_lub'])}ML  +  Petrol add-on {fv(d['rec_petrol'])}ML  |  "
             f"Current base: {fv(d['tot_cy'])}ML  →  Potential: {fv(d['tot_cy']+tot_p)}ML  "
             f"(+{tot_p/d['tot_cy']*100:.0f}%)")
    txt(sl,Inches(0.15),SH-Inches(0.35),SW-Inches(0.3),Inches(0.30),
        summary,size=8.5,bold=True,fg=GOLD,align=PP_ALIGN.CENTER)

    set_notes(sl,
        f"SLIDE 7 — {city.upper()} VOLUME GROWTH ROADMAP\n\n"
        f"PRESENTER TALKING POINTS:\n"
        f"End here. This is the 'so what' slide. Every number on this slide came from the data "
        f"in the previous six slides. There are no assumptions except those clearly stated.\n\n"
        f"HOW TO PRESENT THE FOUR LEVERS:\n"
        f"Lever 1 (Inactive, +{fv(d['rec_inact'])}ML): "
        f"'This is money we have already earned and lost. "
        f"{d['n_inact']} stations went dark. We need {d['n_inact']} field visits in 2 weeks.'\n\n"
        f"Lever 2 (Under-performers, +{fv(d['rec_undp'])}ML): "
        f"'These {d['n_undp']} stations are still alive but running at less than half their own LY. "
        f"Something broke. We need a diagnosis, not a target.'\n\n"
        f"Lever 3 (Lubes, +{fv(d['rec_lub'])}ML): "
        f"'{d['n_fl']} stations have zero lube sales. "
        f"Our data shows that stations with lubes average {d['mul_avg']:.1f}x more total volume. "
        f"This lever costs trial inventory and training — no major capex.'\n\n"
        f"Lever 4 (Petrol add-on, +{fv(d['rec_petrol'])}ML): "
        f"'{d['n_d_no_p']} diesel stations not selling petrol. "
        f"This one needs physical assessment — but the opportunity is real.'\n\n"
        f"CLOSE WITH:\n"
        f"'These four levers together are +{fv(tot_p)}ML on a current base of {fv(d['tot_cy'])}ML. "
        f"That is a {tot_p/d['tot_cy']*100:.0f}% volume uplift. "
        f"Before we leave this room: who owns Lever 1? Who owns Lever 3? "
        f"What is the commitment date?'\n\n"
        f"NUMBERS SUMMARY:\n"
        f"  L1 Inactive: {d['n_inact']} stns, {fv(d['inact_ly_vol'])}ML LY → +{fv(d['rec_inact'])}ML at 65%\n"
        f"  L2 Under-perf: {d['n_undp']} stns → +{fv(d['rec_undp'])}ML at 75% of LY\n"
        f"  L3 Lubes: {d['n_fl']} stns × {d['avg_lub_per_stn']*1000:.0f}L/stn → +{fv(d['rec_lub'])}ML\n"
        f"  L4 Petrol: {d['n_d_no_p']} stns × 25% diesel vol → +{fv(d['rec_petrol'])}ML\n"
        f"  TOTAL: +{fv(tot_p)}ML (+{tot_p/d['tot_cy']*100:.0f}%)\n"
    )

# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
print("Building presentation …")
prs = Presentation()
prs.slide_width=SW; prs.slide_height=SH
BL = prs.slide_layouts[6]

make_cover(prs)

for rank, city in enumerate(TOP5, 1):
    d = city_data[city]
    print(f"  [{rank}/5] {city} …")
    make_divider(prs, city, d, rank)
    make_s1(prs, city, d, period)
    make_s2(prs, city, d, period)
    make_s3(prs, city, d, period)
    make_s4(prs, city, d, period)
    make_s5(prs, city, d, period)
    make_s6(prs, city, d, period)
    make_s7(prs, city, d, period)

out = Path("reports/PSO_Top5_Cities_Volume.pptx")
out.parent.mkdir(exist_ok=True)
prs.save(str(out))
print(f"\nSaved → {out.resolve()}")
print(f"Total slides: {len(prs.slides)}")
for rank, city in enumerate(TOP5,1):
    d=city_data[city]
    print(f"  {city}: {fv(d['tot_cy'])}ML  pot +{fv(d['tot_pot'])}ML  ({d['n_inact']} inactive, {d['n80']}/{d['n_tot']} for 80%)")
