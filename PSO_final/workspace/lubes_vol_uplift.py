"""
PSO Lubricants — Volume Uplift Potential from 'Where to Focus' initiatives.
Produces:
  1. Word table (top 20 cities × 4 initiatives × conservative & optimal)
  2. Single PPTX summary slide
"""
import sys, os, io
sys.path.insert(0,'src')
sys.stdout.reconfigure(encoding='utf-8')

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import parse_xml

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

from pso import ingest
from _pso_common import INPUT_PATH, get_period_label, out_path

# ── scenario parameters ───────────────────────────────────────────────────────
# Full potential = all gap stations reach category average volume
# Conservative  = 25% of full potential (realistic 12-month target)
# Optimal       = 55% of full potential (stretch 18-month target)
CONSERV = 0.25
OPTIMAL = 0.55

# ── colours ───────────────────────────────────────────────────────────────────
W_BLUE   = RGBColor(0x0A,0x16,0x28)
W_TEAL   = RGBColor(0x14,0x8F,0x77)
W_GREEN  = RGBColor(0x1E,0x84,0x49)
W_RED    = RGBColor(0xC0,0x00,0x00)
W_ORANGE = RGBColor(0xE6,0x7E,0x22)
W_WHITE  = RGBColor(0xFF,0xFF,0xFF)
W_GREY   = RGBColor(0x7F,0x8C,0x8D)
W_DARK   = RGBColor(0x1B,0x2A,0x4A)
W_GOLD   = RGBColor(0xD4,0xAC,0x0D)

def set_bg(cell, hex_str):
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.append(parse_xml(
        f'<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
        f'w:val="clear" w:color="auto" w:fill="{hex_str}"/>'))

def ct(cell, text, bold=False, size=8, color=None,
       align=WD_ALIGN_PARAGRAPH.CENTER, italic=False):
    cell.text = ''
    p = cell.paragraphs[0]; p.alignment = align
    r = p.add_run(str(text))
    r.bold=bold; r.italic=italic; r.font.size=Pt(size); r.font.name='Arial'
    if color: r.font.color.rgb = color

def fmt(v):
    kl = v/1000
    if kl >= 100: return f"{kl:.1f}"
    if kl >= 10:  return f"{kl:.2f}"
    if kl >= 1:   return f"{kl:.3f}"
    return "—"

def chg(v): return f"+{v:.1f}%" if v >= 0 else f"{v:.1f}%"

# ── load & compute ────────────────────────────────────────────────────────────
print("Loading data…")
df, _ = ingest.load(INPUT_PATH)
REPORT_DATE = get_period_label(df)
retail    = df[df['IsRetail'] & ~df['IsInternational']].copy()
lubes_all = retail[retail['FuelSegment']=='Lubricants'].copy()

top20 = (lubes_all.groupby('CityNorm')['SalesLtr_CY'].sum()
         .sort_values(ascending=False).head(20).index.tolist())

rows = []
for city in top20:
    cl = lubes_all[lubes_all['CityNorm']==city]
    ca = retail[retail['CityNorm']==city]

    vol_cy      = cl['SalesLtr_CY'].sum()
    vol_ly      = cl['SalesLtr_LY'].sum()
    vol_sply    = cl['SalesLtr_SPLY'].sum()
    vol_chg_pct = (vol_cy-vol_sply)/vol_sply*100 if vol_sply else 0

    stn_all  = ca['Customer Number'].nunique()
    stn_tot  = cl.groupby('Customer Number')['SalesLtr_CY'].sum()
    n_active = (stn_tot > 0).sum()
    n_zero   = stn_all - n_active
    avg_stn  = vol_cy / stn_all if stn_all else 0
    med_stn  = stn_tot[stn_tot>0].median() if n_active else 0

    # low-tier avg
    low_stns  = stn_tot[(stn_tot>0) & (stn_tot<=med_stn)]
    n_low     = len(low_stns)
    avg_low   = low_stns.mean() if n_low else 0

    # per category
    cat = {}
    for c in ['LOW GRADE','PCMO']:
        dc      = cl[cl['LubeCategory']==c]
        sv      = dc.groupby('Customer Number')['SalesLtr_CY'].sum()
        n_sell  = (sv>0).sum()
        avg_c   = sv[sv>0].mean() if (sv>0).sum() else 0
        cat[c]  = dict(n_gap=stn_all-n_sell, avg=avg_c)

    # ── full potentials (litres) ───────────────────────────────────────────
    fp1 = cat['LOW GRADE']['n_gap'] * cat['LOW GRADE']['avg']   # activate LG
    fp2 = cat['PCMO']['n_gap']      * cat['PCMO']['avg']        # activate PCMO
    fp3 = max(0, (med_stn - avg_low)) * n_low                   # upgrade low-tier
    fp4 = n_zero * avg_stn                                       # convert zero stns

    rows.append(dict(
        city     = city,
        vol_cy   = vol_cy,
        vol_chg  = vol_chg_pct,
        n_stns   = stn_all,
        n_zero   = n_zero,
        n_low    = n_low,
        lg_gap   = cat['LOW GRADE']['n_gap'],
        pcmo_gap = cat['PCMO']['n_gap'],
        fp1=fp1, fp2=fp2, fp3=fp3, fp4=fp4,
        fp_total = fp1+fp2+fp3+fp4,
        c1=fp1*CONSERV, c2=fp2*CONSERV, c3=fp3*CONSERV, c4=fp4*CONSERV,
        o1=fp1*OPTIMAL, o2=fp2*OPTIMAL, o3=fp3*OPTIMAL, o4=fp4*OPTIMAL,
        c_total  = (fp1+fp2+fp3+fp4)*CONSERV,
        o_total  = (fp1+fp2+fp3+fp4)*OPTIMAL,
        c_pct    = (fp1+fp2+fp3+fp4)*CONSERV/vol_cy*100 if vol_cy else 0,
        o_pct    = (fp1+fp2+fp3+fp4)*OPTIMAL/vol_cy*100 if vol_cy else 0,
    ))

# national sums
nat_vol    = sum(r['vol_cy']   for r in rows)
nat_c_tot  = sum(r['c_total']  for r in rows)
nat_o_tot  = sum(r['o_total']  for r in rows)
nat_fp_tot = sum(r['fp_total'] for r in rows)
nat_c1 = sum(r['c1'] for r in rows); nat_o1 = sum(r['o1'] for r in rows)
nat_c2 = sum(r['c2'] for r in rows); nat_o2 = sum(r['o2'] for r in rows)
nat_c3 = sum(r['c3'] for r in rows); nat_o3 = sum(r['o3'] for r in rows)
nat_c4 = sum(r['c4'] for r in rows); nat_o4 = sum(r['o4'] for r in rows)

# ══════════════════════════════════════════════════════════════════════════════
# WORD TABLE
# ══════════════════════════════════════════════════════════════════════════════
print("Building Word table…")
doc = Document()
sec = doc.sections[0]
sec.page_width=Cm(29.7); sec.page_height=Cm(21.0)
sec.left_margin=sec.right_margin=Cm(1.3)
sec.top_margin=sec.bottom_margin=Cm(1.1)

# Title
p=doc.add_paragraph(); p.paragraph_format.space_after=Pt(2)
r=p.add_run('PSO Lubricants — Volume Uplift Potential from Sales Initiatives (Top 20 Cities)')
r.bold=True; r.font.size=Pt(15); r.font.color.rgb=W_BLUE; r.font.name='Arial'

p2=doc.add_paragraph(); p2.paragraph_format.space_before=Pt(0); p2.paragraph_format.space_after=Pt(4)
r2=p2.add_run(f'Volume Edition  |  {REPORT_DATE}  |  All volumes in KL (kilolitres = litres ÷ 1,000)')
r2.font.size=Pt(9); r2.font.color.rgb=W_TEAL; r2.font.name='Arial'

# Assumption box
p3=doc.add_paragraph(); p3.paragraph_format.space_after=Pt(6)
for seg, bd in [
    ('Scenario assumptions:  ', True),
    ('Conservative (25% of full potential) — ', False),
    ('25% of identified gap stations activated, each achieving 50% of category avg volume.  ', False),
    ('Optimal (55% of full potential) — ', False),
    ('55% of gap stations activated at 75%+ of category avg.  |  ', False),
    ('Full Potential = 100% of gap stations reach category average.', False),
]:
    rr=p3.add_run(seg); rr.bold=bd; rr.font.size=Pt(8.5); rr.font.color.rgb=W_DARK; rr.font.name='Arial'

# ─── column definitions ───────────────────────────────────────────────────────
HDR = [
    '#','City','Current\nVol (KL)','vs SPLY',
    # I1
    'I1: Low Grade\nGap Stns','I1 Full\nPotential','I1\nConserv','I1\nOptimal',
    # I2
    'I2: PCMO\nGap Stns',     'I2 Full\nPotential','I2\nConserv','I2\nOptimal',
    # I3
    'I3: Low-Tier\nStns',     'I3 Full\nPotential','I3\nConserv','I3\nOptimal',
    # I4
    'I4: Zero\nStns',         'I4 Full\nPotential','I4\nConserv','I4\nOptimal',
    # totals
    'Full\nPotential','Conservative\nTotal','Optimal\nTotal',
    'Conserv\nUplift %','Optimal\nUplift %',
]
N = len(HDR)

# group colours
GRP = {  # col_index → bg_hex (header)
    0:'0A1628', 1:'0A1628', 2:'0A1628', 3:'0A1628',
    4:'0B5E32', 5:'0B5E32', 6:'0B5E32', 7:'0B5E32',   # I1 Low Grade green
    8:'5B2C6F', 9:'5B2C6F',10:'5B2C6F',11:'5B2C6F',   # I2 PCMO purple
   12:'784212',13:'784212',14:'784212',15:'784212',    # I3 low-tier brown
   16:'8B0000',17:'8B0000',18:'8B0000',19:'8B0000',   # I4 zero red
   20:'1A3A6E',21:'145A32',22:'7D6608',               # totals
   23:'145A32',24:'7D6608',                            # uplift %
}

tbl = doc.add_table(rows=len(rows)+3, cols=N)
tbl.style='Table Grid'; tbl.alignment=WD_TABLE_ALIGNMENT.CENTER

# header row 0
for ci,h in enumerate(HDR):
    set_bg(tbl.cell(0,ci), GRP.get(ci,'0A1628'))
    ct(tbl.cell(0,ci), h, bold=True, size=7, color=W_WHITE)

# sub-header row 1
SUB = ['','','KL','vs SPLY',
       'count','KL','KL','KL',
       'count','KL','KL','KL',
       'count','KL','KL','KL',
       'count','KL','KL','KL',
       'KL','KL','KL','%','%']
for ci,s in enumerate(SUB):
    set_bg(tbl.cell(1,ci), GRP.get(ci,'0A1628'))
    ct(tbl.cell(1,ci), s, size=6, color=W_TEAL, italic=True)

# data rows
for ri,d in enumerate(rows, 2):
    bg = 'F2FFF4' if ri%2==0 else 'FFFFFF'
    c_col = W_GREEN; o_col = W_GOLD

    vals = [
        str(ri-1), d['city'],
        fmt(d['vol_cy']), chg(d['vol_chg']),
        str(d['lg_gap']),   fmt(d['fp1']), fmt(d['c1']), fmt(d['o1']),
        str(d['pcmo_gap']), fmt(d['fp2']), fmt(d['c2']), fmt(d['o2']),
        str(d['n_low']),    fmt(d['fp3']), fmt(d['c3']), fmt(d['o3']),
        str(d['n_zero']),   fmt(d['fp4']), fmt(d['c4']), fmt(d['o4']),
        fmt(d['fp_total']),
        fmt(d['c_total']),
        fmt(d['o_total']),
        f"+{d['c_pct']:.1f}%",
        f"+{d['o_pct']:.1f}%",
    ]
    for ci,v in enumerate(vals):
        set_bg(tbl.cell(ri,ci), bg)
        fc   = None
        bold = ci in [1,2,21,22]
        if ci==3:  fc = W_GREEN if d['vol_chg']>=0 else W_RED
        if ci==21: fc = W_GREEN
        if ci==22: fc = W_GOLD
        if ci==23: fc = W_GREEN
        if ci==24: fc = W_GOLD
        al = WD_ALIGN_PARAGRAPH.LEFT if ci==1 else WD_ALIGN_PARAGRAPH.CENTER
        ct(tbl.cell(ri,ci), v, bold=bold, size=7.5, color=fc, align=al)

# totals row
TR = len(rows)+2
set_bg(tbl.cell(TR,0), '0A1628')
ct(tbl.cell(TR,0), '', bold=True, size=7.5, color=W_WHITE)
set_bg(tbl.cell(TR,1), '0A1628')
ct(tbl.cell(TR,1), 'NATIONAL TOTAL', bold=True, size=7.5, color=W_WHITE,
   align=WD_ALIGN_PARAGRAPH.LEFT)
set_bg(tbl.cell(TR,2), '0A1628')
ct(tbl.cell(TR,2), fmt(nat_vol), bold=True, size=7.5, color=W_WHITE)
set_bg(tbl.cell(TR,3), '0A1628')
ct(tbl.cell(TR,3), '', size=7.5, color=W_WHITE)

tot_vals = [
    fmt(sum(r['lg_gap'] for r in rows)), fmt(sum(r['fp1'] for r in rows)), fmt(nat_c1), fmt(nat_o1),
    fmt(sum(r['pcmo_gap'] for r in rows)), fmt(sum(r['fp2'] for r in rows)), fmt(nat_c2), fmt(nat_o2),
    fmt(sum(r['n_low'] for r in rows)),  fmt(sum(r['fp3'] for r in rows)), fmt(nat_c3), fmt(nat_o3),
    fmt(sum(r['n_zero'] for r in rows)), fmt(sum(r['fp4'] for r in rows)), fmt(nat_c4), fmt(nat_o4),
    fmt(nat_fp_tot), fmt(nat_c_tot), fmt(nat_o_tot),
    f"+{nat_c_tot/nat_vol*100:.1f}%", f"+{nat_o_tot/nat_vol*100:.1f}%",
]
for ci,v in enumerate(tot_vals, 4):
    set_bg(tbl.cell(TR,ci), GRP.get(ci,'0A1628'))
    fc = None
    if ci==21: fc=W_TEAL
    if ci==22: fc=W_GOLD
    ct(tbl.cell(TR,ci), v, bold=True, size=7.5, color=W_WHITE if ci<21 else fc)

# column widths
CW = [0.42,2.70,1.25,0.85,
      0.80,1.15,1.10,1.10,
      0.80,1.15,1.10,1.10,
      0.80,1.15,1.10,1.10,
      0.80,1.15,1.10,1.10,
      1.15,1.30,1.30,0.95,0.95]
for ci,w in enumerate(CW):
    for row in tbl.rows:
        row.cells[ci].width=Cm(w)

# legend note
doc.add_paragraph()
fn=doc.add_paragraph()
fn.paragraph_format.space_before=Pt(2)
notes = [
    ('Initiative definitions:  ', True),
    ('I1 = Activate stations not selling Low Grade (fastest growing category).  '
     'I2 = Activate stations not selling PCMO (highest margin/litre).  '
     'I3 = Upgrade low-volume stations (below median) toward city median.  '
     'I4 = Convert zero-lubes stations (active fuel stations selling no lubricants).  '
     'Full Potential assumes 100% activation at category average.  '
     'Conservative = 25% of full potential.  Optimal = 55% of full potential.  |  ', False),
    (REPORT_DATE, True),
]
for seg,bd in notes:
    rr=fn.add_run(seg); rr.bold=bd; rr.font.size=Pt(7.5)
    rr.font.color.rgb=W_GREY; rr.font.name='Arial'

out_word = out_path('PSO_Lubes_Vol_Uplift_Table', 'docx', df)
doc.save(out_word)
print(f"Word saved: {out_word}")

# ══════════════════════════════════════════════════════════════════════════════
# PPTX SLIDE
# ══════════════════════════════════════════════════════════════════════════════
print("Building PPTX slide…")

def prgb(h): return PRGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

SW=PInches(10); SH=PInches(5.625)
PC={'NAVY':'0F2035','DNAVY':'0A1628','TEAL':'148F77','TEAL_L':'A2D9CE',
    'GREEN':'1E8449','GREEN_L':'EAFAF1','GOLD':'D4AC0D','GOLD_L':'FEF9E7',
    'BLUE':'2E86C1','RED':'C00000','ORANGE':'E67E22','ORANGE_L':'E59866',
    'GREY':'7F8C8D','GREY2':'BFBFBF','WHITE':'FFFFFF','DARK':'1B2A4A',
    'AMBER':'D35400'}

def prect(sl,x,y,w,h,fill):
    s=sl.shapes.add_shape(1,PInches(x),PInches(y),PInches(w),PInches(h))
    s.line.fill.background(); f=s.fill; f.solid()
    f.fore_color.rgb = fill if isinstance(fill,PRGBColor) else prgb(fill)
    return s

def ptxt(sl,x,y,w,h,text,size=9,bold=False,color='FFFFFF',
         align=PP_ALIGN.LEFT,wrap=True):
    t=sl.shapes.add_textbox(PInches(x),PInches(y),PInches(w),PInches(h))
    tf=t.text_frame; tf.word_wrap=wrap; tf.auto_size=None
    pa=tf.paragraphs[0]; pa.alignment=align
    rn=pa.add_run(); rn.text=text
    rn.font.size=PPt(size); rn.font.bold=bold
    rn.font.color.rgb=prgb(color); rn.font.name='Arial'
    return t

def set_pptx_bg(sl,fill):
    bg=sl.background; f=bg.fill; f.solid()
    f.fore_color.rgb=prgb(fill)

def img_buf(fig):
    buf=io.BytesIO()
    fig.savefig(buf,format='png',dpi=160,bbox_inches='tight')
    buf.seek(0); plt.close(fig); return buf

# ── waterfall / bar chart ─────────────────────────────────────────────────────
fig,axes=plt.subplots(1,2,figsize=(9.5,3.8))
fig.patch.set_facecolor('white')

initiatives = ['I1\nLow Grade','I2\nPCMO','I3\nLow Tier','I4\nZero Stns']
c_vals = [nat_c1/1000, nat_c2/1000, nat_c3/1000, nat_c4/1000]
o_vals = [nat_o1/1000, nat_o2/1000, nat_o3/1000, nat_o4/1000]
cols   = ['#1E8449','#8E44AD','#D35400','#C00000']

ax=axes[0]
x=np.arange(4); bw=0.35
ax.bar(x-bw/2, c_vals, bw, color=[c+'AA' for c in cols], label='Conservative', zorder=3)
ax.bar(x+bw/2, o_vals, bw, color=cols,                    label='Optimal',      zorder=3, alpha=0.9)
ax.set_xticks(x); ax.set_xticklabels(initiatives, fontsize=8)
ax.set_ylabel('Additional Volume (KL)', fontsize=8)
ax.set_title('Uplift by Initiative', fontsize=9, fontweight='bold', color='#0A1628', pad=6)
ax.yaxis.grid(True,linestyle='--',alpha=0.4); ax.set_axisbelow(True)
ax.spines[['top','right']].set_visible(False)
ax.legend(fontsize=7.5)
for xi,(cv,ov) in enumerate(zip(c_vals,o_vals)):
    ax.annotate(f'{cv:.0f}',xy=(xi-bw/2,cv),xytext=(0,3),textcoords='offset points',
                ha='center',fontsize=7,color='#148F77')
    ax.annotate(f'{ov:.0f}',xy=(xi+bw/2,ov),xytext=(0,3),textcoords='offset points',
                ha='center',fontsize=7,color='#D4AC0D')

ax2=axes[1]
categories=['Current\nVolume','+ Conservative','+ Optimal']
values=[nat_vol/1000,(nat_vol+nat_c_tot)/1000,(nat_vol+nat_o_tot)/1000]
bar_colors=['#0A1628','#148F77','#D4AC0D']
bars=ax2.bar(categories,values,color=bar_colors,width=0.5,zorder=3,alpha=0.9)
ax2.set_ylabel('Total Volume (KL)', fontsize=8)
ax2.set_title('Projected Total Volume — Top 20 Cities', fontsize=9,
              fontweight='bold', color='#0A1628', pad=6)
ax2.yaxis.grid(True,linestyle='--',alpha=0.4); ax2.set_axisbelow(True)
ax2.spines[['top','right']].set_visible(False)
for bar,val,lbl in zip(bars,values,
    [f'{nat_vol/1000:.0f} KL\n(Base)',
     f'{(nat_vol+nat_c_tot)/1000:.0f} KL\n(+{nat_c_tot/nat_vol*100:.1f}%)',
     f'{(nat_vol+nat_o_tot)/1000:.0f} KL\n(+{nat_o_tot/nat_vol*100:.1f}%)']):
    ax2.annotate(lbl,xy=(bar.get_x()+bar.get_width()/2, bar.get_height()),
                 xytext=(0,6),textcoords='offset points',ha='center',fontsize=8,
                 fontweight='bold',color=bar.get_facecolor())

plt.tight_layout()
chart_buf=img_buf(fig)

# ── build slide ───────────────────────────────────────────────────────────────
prs=Presentation()
prs.slide_width=SW; prs.slide_height=SH

sl=prs.slides.add_slide(prs.slide_layouts[6])
set_pptx_bg(sl, PC['NAVY'])

# header
prect(sl,0,0,10,0.58,PC['DNAVY'])
ptxt(sl,0.25,0.13,7.5,0.32,
     'Volume Uplift Potential — Where to Focus Initiatives',
     size=17,bold=True,color=PC['WHITE'])
ptxt(sl,7.5,0.18,2.30,0.20,
     f'Top 20 Cities  |  {REPORT_DATE}',
     size=8.5,color=PC['ORANGE_L'],align=PP_ALIGN.RIGHT)

# chart (left side)
sl.shapes.add_picture(chart_buf, PInches(0.18), PInches(0.66),
                      PInches(5.70), PInches(3.10))

# right side: scenario panels
# ── CONSERVATIVE panel ────────────────────────────────────────────────────────
cx=6.02; cy=0.66; cw=1.84; ch=4.68
prect(sl,cx,cy,cw,ch, PC['TEAL'])
ptxt(sl,cx+0.12,cy+0.10,cw-0.20,0.24,
     'CONSERVATIVE', size=9,bold=True,color=PC['WHITE'])
ptxt(sl,cx+0.12,cy+0.10+0.26,cw-0.20,0.14,
     '25% of full potential', size=7,color=PC['TEAL_L'])

# big number
ptxt(sl,cx+0.08,cy+0.56,cw-0.14,0.45,
     fmt(nat_c_tot)+' KL', size=20,bold=True,color=PC['WHITE'])
ptxt(sl,cx+0.08,cy+1.04,cw-0.14,0.18,
     f"+{nat_c_tot/nat_vol*100:.1f}% uplift on current",
     size=8.5,bold=True,color=PC['TEAL_L'])

ptxt(sl,cx+0.08,cy+1.30,cw-0.14,0.14,
     f"New total:  {fmt(nat_vol+nat_c_tot)} KL",
     size=8,color=PC['WHITE'])

# breakdown
prect(sl,cx+0.08,cy+1.58,cw-0.16,0.02,PC['TEAL_L'])
lines_c=[
    ('I1 Low Grade', fmt(nat_c1)),
    ('I2 PCMO',      fmt(nat_c2)),
    ('I3 Low-Tier',  fmt(nat_c3)),
    ('I4 Zero Stns', fmt(nat_c4)),
]
for ii,(lbl,val) in enumerate(lines_c):
    iy=cy+1.66+ii*0.44
    prect(sl,cx+0.08,iy,cw-0.16,0.38,PC['DNAVY'])
    ptxt(sl,cx+0.14,iy+0.04,cw-0.28,0.16,lbl, size=7,color=PC['TEAL_L'])
    ptxt(sl,cx+0.14,iy+0.20,cw-0.28,0.16,f"+{val} KL", size=8,bold=True,color=PC['WHITE'])

# ── OPTIMAL panel ─────────────────────────────────────────────────────────────
ox=cx+cw+0.12; oy=cy
prect(sl,ox,oy,cw,ch,PC['GOLD'])
ptxt(sl,ox+0.12,oy+0.10,cw-0.20,0.24,
     'OPTIMAL', size=9,bold=True,color=PC['DARK'])
ptxt(sl,ox+0.12,oy+0.10+0.26,cw-0.20,0.14,
     '55% of full potential', size=7,color=PC['DARK'])

ptxt(sl,ox+0.08,oy+0.56,cw-0.14,0.45,
     fmt(nat_o_tot)+' KL', size=20,bold=True,color=PC['DARK'])
ptxt(sl,ox+0.08,oy+1.04,cw-0.14,0.18,
     f"+{nat_o_tot/nat_vol*100:.1f}% uplift on current",
     size=8.5,bold=True,color=PC['DARK'])

ptxt(sl,ox+0.08,oy+1.30,cw-0.14,0.14,
     f"New total:  {fmt(nat_vol+nat_o_tot)} KL",
     size=8,color=PC['DARK'])

prect(sl,ox+0.08,oy+1.58,cw-0.16,0.02,PC['DARK'])
lines_o=[
    ('I1 Low Grade', fmt(nat_o1)),
    ('I2 PCMO',      fmt(nat_o2)),
    ('I3 Low-Tier',  fmt(nat_o3)),
    ('I4 Zero Stns', fmt(nat_o4)),
]
for ii,(lbl,val) in enumerate(lines_o):
    iy=oy+1.66+ii*0.44
    prect(sl,ox+0.08,iy,cw-0.16,0.38,PC['DNAVY'])
    ptxt(sl,ox+0.14,iy+0.04,cw-0.28,0.16,lbl, size=7,color=PC['GOLD'])
    ptxt(sl,ox+0.14,iy+0.20,cw-0.28,0.16,f"+{val} KL", size=8,bold=True,color=PC['WHITE'])

# footer
prect(sl,0,5.38,10,0.24,PC['DNAVY'])
ptxt(sl,0.20,5.40,9.60,0.18,
     f"Assumes: Conservative = 25% of gap stations activated at 50% of category avg  |  "
     f"Optimal = 55% at 75% of avg  |  Full potential (100%) = {fmt(nat_fp_tot)} KL  |  {REPORT_DATE}",
     size=7,color=PC['ORANGE_L'])

out_pptx = out_path('PSO_Lubes_Uplift_Scenarios', 'pptx', df)
prs.save(out_pptx)
print(f"PPTX saved: {out_pptx}")

# print summary
print(f"\n{'─'*55}")
print(f"  Current volume (top 20):  {nat_vol/1000:>8.1f} KL")
print(f"  Full potential:           {nat_fp_tot/1000:>8.1f} KL  (+{nat_fp_tot/nat_vol*100:.1f}%)")
print(f"  Conservative uplift:      {nat_c_tot/1000:>8.1f} KL  (+{nat_c_tot/nat_vol*100:.1f}%)")
print(f"  Optimal uplift:           {nat_o_tot/1000:>8.1f} KL  (+{nat_o_tot/nat_vol*100:.1f}%)")
print(f"{'─'*55}")
