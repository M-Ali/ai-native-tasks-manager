"""
PSO City Fuels Profile Generator
Mirrors workspace/city_profiles.py (Lubricants), adapted to Fuels.
Generates one PPTX per city with 5 slides:
  1. Cover
  2. At a Glance
  3. Category Performance
  4. Station Performance
  5. Priorities / Where to Focus (R95 upsell, Diesel/HSD cross-sell, low-tier
     upgrade, zero-fuels conversion — mirrors workspace/fuels_vol_uplift.py's
     I1-I4 initiatives)
"""
import sys, os
sys.path.insert(0, 'src')
sys.stdout.reconfigure(encoding='utf-8')

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pso import ingest
from _pso_common import INPUT_PATH, get_period_label, out_path

# ── constants ─────────────────────────────────────────────────────────────────
SW = Inches(10); SH = Inches(5.625)

def rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C = dict(
    NAVY='0F2035', DNAVY='0A1628', DNAV2='1B2A4A',
    ORANGE='E67E22', ORANGE_L='E59866',
    BLUE='2E86C1', BLUE_L='7FB3D3', BLUE2='1A5276',
    GREEN='1E8449', GREEN_L='EAFAF1',
    RED='E74C3C', RED_D='922B21', RED_BG='FDEDEC',
    PURPLE='8E44AD',
    AMBER='D35400',
    GOLD='D4AC0D',
    GREY='7F8C8D', GREY2='5D6D7E', GREY3='95A5A6',
    WHITE='FFFFFF',
    BG_BLUE='EAF2FF', BG_GREEN='EAFAF1',
    BG_YELLOW='FEF9E7', BG_RED='FDEDEC', BG_LIGHT='EBF5FB',
    KPI_BG='F4F6F8',
    BROWN='784212',
    TRACK='EAF2FF',
)

CAT_COLORS = {
    'Diesel':      C['BLUE'],
    'Petrol':      C['ORANGE'],
    'Other Fuels': C['PURPLE'],
    'LPG':         C['GREEN'],
}
CAT_ORDER = ['Diesel', 'Petrol', 'Other Fuels', 'LPG']

# ── helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill_hex):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    fill = shape.fill; fill.solid()
    fill.fore_color.rgb = fill_hex if isinstance(fill_hex, RGBColor) else rgb(fill_hex)
    return shape

def txt(slide, x, y, w, h, text, size=9, bold=False, color='FFFFFF',
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = wrap
    tf.auto_size = None
    para = tf.paragraphs[0]; para.alignment = align
    run = para.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = rgb(color); run.font.name = 'Arial'
    run.font.italic = italic
    return txb

def set_bg(slide, fill_hex):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = rgb(fill_hex)

def fmt_m(v):
    if v >= 1e9: return f"PKR {v/1e9:.2f}B"
    if v >= 1e6: return f"PKR {v/1e6:.2f}M"
    if v >= 1e3: return f"PKR {v/1e3:.0f}K"
    return f"PKR {v:.0f}"

def fmt_b(v):
    return f"PKR {v/1e9:.2f}B"

# ── data loader ───────────────────────────────────────────────────────────────
print("Loading data…")
df, _ = ingest.load(INPUT_PATH)
REPORT_PERIOD = get_period_label(df)
retail = df[df['IsRetail'] & ~df['IsInternational']].copy()
fuels_all = retail[retail['FuelSegment'].isin(CAT_ORDER)].copy()

def compute_city(city_norm):
    city_fuels = fuels_all[fuels_all['CityNorm'] == city_norm].copy()
    city_all = retail[retail['CityNorm'] == city_norm].copy()

    total_rev = city_all['SalesGRS_CY'].sum()
    fuels_rev = city_fuels['SalesGRS_CY'].sum()
    fuels_rev_sply = city_fuels['SalesGRS_SPLY'].sum()

    stn_total = (city_fuels.groupby('Customer Number')['SalesGRS_CY'].sum().sort_values(ascending=False))
    all_stns = city_all['Customer Number'].nunique()

    n_stns = all_stns
    n_active = (stn_total > 0).sum()
    n_zero = n_stns - n_active
    avg_per_stn = fuels_rev / n_stns if n_stns else 0
    med_per_stn = stn_total[stn_total > 0].median() if n_active else 0
    min_per_stn = stn_total[stn_total > 0].min() if n_active else 0
    max_per_stn = stn_total.max() if n_active else 0

    top10_rev = stn_total.head(10).sum()
    top10_pct = top10_rev / fuels_rev * 100 if fuels_rev else 0

    active_cats = [c for c in CAT_ORDER if city_fuels[city_fuels['FuelSegment'] == c]['SalesGRS_CY'].sum() > 0]
    n_cats = len(active_cats)

    cat_stats = {}
    for cat in CAT_ORDER:
        df_cat = city_fuels[city_fuels['FuelSegment'] == cat]
        stn_cat = df_cat.groupby('Customer Number')['SalesGRS_CY'].sum()
        stn_cat_active = stn_cat[stn_cat > 0]
        rev = stn_cat.sum()
        n_sell = (stn_cat > 0).sum()
        cat_stats[cat] = dict(
            rev=rev, rev_pct=rev/fuels_rev*100 if fuels_rev else 0,
            n_selling=n_sell, n_pct=n_sell/n_stns*100 if n_stns else 0,
            avg=stn_cat_active.mean() if len(stn_cat_active) else 0,
            median=stn_cat_active.median() if len(stn_cat_active) else 0,
            min_val=stn_cat_active.min() if len(stn_cat_active) else 0,
            max_val=stn_cat_active.max() if len(stn_cat_active) else 0,
        )

    p75 = stn_total[stn_total > 0].quantile(0.75) if n_active else 0
    n_high = (stn_total > p75).sum()
    n_mid = ((stn_total > med_per_stn) & (stn_total <= p75)).sum()
    n_low = ((stn_total > 0) & (stn_total <= med_per_stn)).sum()
    rev_high = stn_total[stn_total > p75].sum()
    rev_mid = stn_total[(stn_total > med_per_stn) & (stn_total <= p75)].sum()
    rev_low = stn_total[(stn_total > 0) & (stn_total <= med_per_stn)].sum()

    top10_codes = stn_total.head(10).index.tolist()
    top10_detail = []
    for code in top10_codes:
        stn_data = city_fuels[city_fuels['Customer Number'] == code]
        name = stn_data['Name 1'].iloc[0] if len(stn_data) else code
        total = stn_total[code]
        by_cat = {cat: stn_data[stn_data['FuelSegment'] == cat]['SalesGRS_CY'].sum() for cat in CAT_ORDER}
        top10_detail.append({'code': code, 'name': name, 'total': total, 'by_cat': by_cat})

    # R95 / Diesel(HSD) gap — mirrors fuels_vol_uplift.py's I1/I2
    petrol = city_fuels[city_fuels['FuelSegment'] == 'Petrol']
    r95_sv = petrol[petrol['ProductCategory'] == 'R95'].groupby('Customer Number')['SalesGRS_CY'].sum()
    pmg_sv = petrol[petrol['ProductCategory'] == 'PMG'].groupby('Customer Number')['SalesGRS_CY'].sum()
    n_petrol_stns = pmg_sv.index.union(r95_sv.index).nunique() if (len(pmg_sv) or len(r95_sv)) else 0
    no_r95 = max(0, n_petrol_stns - (r95_sv > 0).sum())
    r95_avg = r95_sv[r95_sv > 0].mean() if (r95_sv > 0).sum() else 0

    hsd = city_fuels[city_fuels['ProductCategory'] == 'HSD']
    hsd_sv = hsd.groupby('Customer Number')['SalesGRS_CY'].sum()
    no_hsd = max(0, n_stns - (hsd_sv > 0).sum())
    hsd_avg = hsd_sv[hsd_sv > 0].mean() if (hsd_sv > 0).sum() else 0

    return dict(
        city=city_norm, n_stns=n_stns, n_active=n_active, n_zero=n_zero,
        total_rev=total_rev, fuels_rev=fuels_rev, fuels_rev_sply=fuels_rev_sply,
        avg_per_stn=avg_per_stn, med_per_stn=med_per_stn,
        min_per_stn=min_per_stn, max_per_stn=max_per_stn,
        top10_rev=top10_rev, top10_pct=top10_pct,
        n_cats=n_cats, active_cats=active_cats, cat=cat_stats,
        n_high=n_high, n_mid=n_mid, n_low=n_low, n_zero_stn=n_zero,
        rev_high=rev_high, rev_mid=rev_mid, rev_low=rev_low, p75=p75,
        top10=top10_detail,
        no_r95=no_r95, r95_avg=r95_avg, no_hsd=no_hsd, hsd_avg=hsd_avg,
        fuels_share=fuels_rev/total_rev*100 if total_rev else 0,
    )

# ── SLIDE 1: Cover ────────────────────────────────────────────────────────────
def build_slide1(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C['NAVY'])
    city = d['city']
    rect(sl, 0.90, 1.54, 0.50, 0.06, C['ORANGE'])
    txt(sl, 1.54, 1.47, 6.0, 0.20, f"PSO {city.upper()} — RETAIL FUELS DIVISION", size=9, color=C['BLUE_L'])
    rect(sl, 0.93, 1.81, 0.035, 0.90, C['ORANGE'])
    txt(sl, 1.27, 1.75, 7.98, 0.60, f"{city} Station Profile", size=34, bold=True, color=C['WHITE'])
    txt(sl, 1.27, 2.44, 7.98, 0.28, "Fuels Performance Analysis", size=14, color=C['ORANGE_L'])
    rect(sl, 0.90, 3.09, 8.19, 0.018, C['GREY'])
    txt(sl, 0.90, 3.19, 4.0, 0.14, REPORT_PERIOD, size=8, color=C['GREY'])

    kpis = [
        (0.90, 1.01, f"{d['n_stns']}", "Stations"),
        (2.20, 2.05, fmt_b(d['fuels_rev']), "Total Fuels Revenue"),
        (4.55, 1.37, str(d['n_cats']), "Active Categories"),
        (6.22, 2.08, fmt_m(d['avg_per_stn']), "Avg per Station"),
    ]
    for bx, bw, val, lbl in kpis:
        rect(sl, bx, 3.35, bw, 0.80, C['DNAVY'])
        txt(sl, bx+0.18, 3.46, bw-0.22, 0.36, val, size=22, bold=True, color=C['ORANGE'])
        txt(sl, bx+0.18, 3.86, bw-0.22, 0.16, lbl, size=8, color=C['BLUE_L'])

# ── SLIDE 2: At a Glance ─────────────────────────────────────────────────────
def build_slide2(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, 10, 0.61, C['DNAVY'])
    txt(sl, 0.25, 0.14, 6.5, 0.32, f"{d['city']} — Fuels At a Glance", size=17, bold=True, color=C['WHITE'])
    txt(sl, 6.5, 0.20, 3.25, 0.20,
        f"{d['n_stns']} Stations  |  {fmt_b(d['total_rev'])} Total Revenue  |  {REPORT_PERIOD}",
        size=8.5, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    lx, lw = 0.22, 3.06
    boxes = [
        (0.75, 0.86, fmt_b(d['fuels_rev']), 'TOTAL FUELS REVENUE', C['ORANGE'],
         f"{d['fuels_share']:.1f}% of all {d['city']} product revenue ({fmt_b(d['total_rev'])})"),
        (1.73, 0.86, fmt_m(d['avg_per_stn']), 'AVERAGE FUELS PER STATION', C['BLUE'],
         f"Median: {fmt_m(d['med_per_stn'])}  |  Range: {fmt_m(d['min_per_stn'])} - {fmt_m(d['max_per_stn'])}"),
        (2.70, 0.40, str(d['n_active']), 'ACTIVE SELLERS', C['GREEN'],
         f"{d['n_active']/d['n_stns']*100:.1f}% of stations" if d['n_stns'] else "N/A"),
        (3.20, 0.40, str(d['n_zero']), 'ZERO FUELS', C['RED'], "Not selling any fuel"),
        (3.70, 0.40, f"{d['top10_pct']:.1f}%", 'TOP 10 STATION CONCENTRATION', C['PURPLE'],
         f"Top 10 stations = {fmt_m(d['top10_rev'])} of total fuels"),
    ]
    for by, bh, val, lbl, vcol, sub in boxes:
        rect(sl, lx, by, lw, bh, C['KPI_BG'])
        rect(sl, lx+0.03, by, 0.02, bh, vcol)
        txt(sl, lx+0.20, by+0.09, lw-0.26, bh*0.35, val, size=20 if by < 2.5 else 16, bold=True, color=vcol)
        txt(sl, lx+0.20, by+0.09+bh*0.35, lw-0.26, 0.14, lbl, size=7.5, color=C['GREY'])
        txt(sl, lx+0.20, by+0.09+bh*0.35+0.15, lw-0.26, 0.15, sub, size=8, color=C['GREY2'])

    rx = 3.47
    txt(sl, rx, 0.73, 6.5, 0.15, "REVENUE CONTRIBUTION BY CATEGORY", size=8, bold=True, color=C['GREY'])
    active = [c for c in CAT_ORDER if d['cat'][c]['rev'] > 0]
    cat_max_rev = max((d['cat'][c]['rev'] for c in active), default=1) or 1
    track_w = 3.90

    for i, cat in enumerate(active):
        by = 0.90 + i * 0.55
        cs = d['cat'][cat]
        bar_w = max(0.05, cs['rev']/cat_max_rev * track_w)
        col = CAT_COLORS.get(cat, C['GREY3'])
        txt(sl, rx, by+0.03, 1.05, 0.14, cat, size=8, color=C['DNAV2'])
        rect(sl, rx+1.10, by, track_w, 0.19, C['BG_BLUE'])
        rect(sl, rx+1.10, by+0.01, bar_w, 0.17, col)
        txt(sl, rx+1.10+track_w+0.06, by+0.03, 1.25, 0.14,
            f"{cs['rev_pct']:.1f}%  |  {cs['n_selling']} stns ({cs['n_pct']:.0f}%)", size=7, color=C['DNAV2'])

    ins_y = max(0.90 + len(active) * 0.55, 3.90)
    rect(sl, rx, ins_y, 6.38, 1.05, C['BG_YELLOW'])
    rect(sl, rx+0.02, ins_y, 0.02, 1.05, C['BROWN'])
    txt(sl, rx+0.18, ins_y+0.08, 6.10, 0.15, "KEY INSIGHT", size=7.5, bold=True, color=C['BROWN'])

    top2 = sorted(active, key=lambda c: d['cat'][c]['rev'], reverse=True)[:2]
    top2_pct = sum(d['cat'][c]['rev_pct'] for c in top2)
    widest_cat = max(active, key=lambda c: d['cat'][c]['n_selling']) if active else 'Diesel'
    insight = (
        f"{' + '.join(top2)} = {top2_pct:.0f}% of {d['city']} fuels revenue. "
        f"{widest_cat} is the widest-reaching ({d['cat'][widest_cat]['n_pct']:.0f}% of stations). "
        f"{d['no_r95']} petrol-selling stations don't carry R95 (premium upsell opportunity). "
        f"{d['no_hsd']} stations don't carry HSD (diesel cross-sell opportunity)."
    )
    txt(sl, rx+0.18, ins_y+0.26, 6.10, 0.65, insight, size=9.5, color=C['DNAV2'], wrap=True)

# ── SLIDE 3: Category Performance ────────────────────────────────────────────
def build_slide3(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, 10, 0.58, C['DNAVY'])
    txt(sl, 0.25, 0.14, 7.2, 0.30, "Category Performance — Avg, Min & Max per Station",
        size=17, bold=True, color=C['WHITE'])
    txt(sl, 7.25, 0.20, 2.55, 0.18, f"{d['n_stns']} Stations Total  |  All values in PKR Million",
        size=9, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    active_cats = [c for c in CAT_ORDER if d['cat'][c]['rev'] > 0 and d['cat'][c]['n_selling'] > 0]
    chart_top = 0.68; chart_h = 1.72
    row_h = chart_h / max(len(active_cats), 1)
    max_max = max((d['cat'][c]['max_val'] for c in active_cats), default=1) or 1
    track_x = 1.50; track_w = 6.10

    for i, cat in enumerate(active_cats):
        cs = d['cat'][cat]
        col = CAT_COLORS.get(cat, C['GREY3'])
        ry = chart_top + i * row_h + 0.04
        rh = row_h - 0.08
        rect(sl, 0.22, ry, 1.24, rh, C['KPI_BG'])
        txt(sl, 0.25, ry+0.02, 1.18, rh-0.04, cat, size=8, bold=True, color=col)
        txt(sl, 0.25, ry+rh*0.45, 1.18, rh*0.50, f"{cs['n_selling']} stns ({cs['n_pct']:.0f}%)", size=7, color=C['GREY'])
        rect(sl, track_x, ry+rh*0.2, track_w, rh*0.45, C['TRACK'])
        max_w = max(0.05, cs['max_val']/max_max * track_w)
        rect(sl, track_x, ry+rh*0.22, max_w, rh*0.41, C['KPI_BG'])
        avg_w = max(0.03, cs['avg']/max_max * track_w)
        rect(sl, track_x, ry+rh*0.22, avg_w, rh*0.41, col)
        if cs['median'] > 0:
            med_x = track_x + cs['median']/max_max * track_w
            rect(sl, med_x-0.015, ry+rh*0.18, 0.03, rh*0.49, C['DNAVY'])
        vx = track_x + track_w + 0.12
        txt(sl, vx, ry, 2.04, rh*0.50, f"Avg  {fmt_m(cs['avg'])}", size=8, bold=True, color=col)
        txt(sl, vx, ry+rh*0.45, 2.04, rh*0.55, f"Med {fmt_m(cs['median'])}  Max {fmt_m(cs['max_val'])}", size=7, color=C['GREY2'])

    # 3 insight boxes — dynamically pick highest-margin-proxy, most-reliable, least-covered category
    by_rev = sorted(active_cats, key=lambda c: d['cat'][c]['rev'], reverse=True)
    by_avg = sorted(active_cats, key=lambda c: d['cat'][c]['avg'], reverse=True)
    by_cov = sorted(active_cats, key=lambda c: d['cat'][c]['n_pct'])

    insight_data = []
    if by_avg:
        top_avg_cat = by_avg[0]
        cs = d['cat'][top_avg_cat]
        insight_data.append((
            C['RED_BG'], C['RED_D'], f"{top_avg_cat.upper()} — HIGHEST REVENUE PER STATION",
            f"Median {fmt_m(cs['median'])} vs Mean {fmt_m(cs['avg'])} per active station. "
            f"{'High variance — a few stations earn most of this category revenue.' if cs['max_val'] > 3*cs['avg'] else 'Fairly consistent across stations.'} "
            f"Prioritise stations not yet carrying it."
        ))
    if len(by_cov) > 1:
        reliable_cat = by_cov[-1]
        cs = d['cat'][reliable_cat]
        insight_data.append((
            C['BG_GREEN'], C['BLUE2'], f"{reliable_cat.upper()} — THE MOST RELIABLE CATEGORY",
            f"{cs['n_pct']:.0f}% station coverage with "
            f"{'tight' if cs['max_val'] < 5*cs['avg'] else 'moderate'} distribution "
            f"(median {fmt_m(cs['median'])}, avg {fmt_m(cs['avg'])}). Most stations sell it at consistent levels."
        ))
    if by_cov:
        low_cov_cat = by_cov[0]
        cs = d['cat'][low_cov_cat]
        not_selling = d['n_stns'] - cs['n_selling']
        insight_data.append((
            C['BG_YELLOW'], C['BROWN'], f"{low_cov_cat.upper()} — {not_selling} STATIONS NOT SELLING",
            f"{cs['n_pct']:.0f}% coverage leaves {not_selling} stations not stocking it. "
            f"If these stations achieve the city avg of {fmt_m(cs['avg'])}, "
            f"potential uplift ~{fmt_m(not_selling * cs['avg'])}."
        ))

    bx_w = 3.11
    for ii, (bg, tc, title, body) in enumerate(insight_data):
        bx = 0.19 + ii*(bx_w+0.22)
        by = 2.51
        rect(sl, bx, by, bx_w, 0.98, bg)
        rect(sl, bx+0.03, by, 0.02, 0.98, rgb(tc))
        txt(sl, bx+0.18, by+0.07, bx_w-0.22, 0.15, title, size=7.5, bold=True, color=tc)
        txt(sl, bx+0.18, by+0.25, bx_w-0.22, 0.68, body, size=8, color=C['DNAV2'], wrap=True)

# ── SLIDE 4: Station Performance ──────────────────────────────────────────────
def build_slide4(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, 10, 0.58, C['DNAVY'])
    txt(sl, 0.25, 0.14, 6.5, 0.30, "Station Performance — Who Is Delivering", size=17, bold=True, color=C['WHITE'])
    txt(sl, 6.5, 0.20, 3.25, 0.18, f"Tiers based on Total Fuels Revenue  |  Median {fmt_m(d['med_per_stn'])}",
        size=8.5, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    txt(sl, 0.19, 0.68, 3.26, 0.14, "STATION PERFORMANCE TIERS", size=7.5, bold=True, color=C['GREY'])
    tiers = [
        (C['GREEN'], f"HIGH PERFORMERS    {d['n_high']} stations", fmt_m(d['rev_high']),
         f"{d['rev_high']/d['fuels_rev']*100:.1f}% of {d['city']} Fuels  |  Above {fmt_m(d['p75'])} each" if d['fuels_rev'] else "N/A"),
        (C['BLUE'], f"MID PERFORMERS    {d['n_mid']} stations", fmt_m(d['rev_mid']),
         f"{d['rev_mid']/d['fuels_rev']*100:.1f}% of {d['city']} Fuels  |  {fmt_m(d['med_per_stn'])} - {fmt_m(d['p75'])}" if d['fuels_rev'] else "N/A"),
        (C['AMBER'], f"LOW PERFORMERS    {d['n_low']} stations", fmt_m(d['rev_low']),
         f"{d['rev_low']/d['fuels_rev']*100:.1f}% of {d['city']} Fuels  |  Below {fmt_m(d['med_per_stn'])}" if d['fuels_rev'] else "N/A"),
        (C['RED'], f"ZERO FUELS    {d['n_zero_stn']} stations", "PKR 0", "Not selling any fuel category"),
    ]
    for ti, (col, title, val, sub) in enumerate(tiers):
        by = 0.86 + ti*0.97
        rect(sl, 0.19, by, 3.19, 0.83, col)
        txt(sl, 0.33, by+0.08, 2.98, 0.18, title, size=9, bold=True, color=C['WHITE'])
        txt(sl, 0.33, by+0.28, 2.98, 0.30, val, size=18, bold=True, color=C['WHITE'])
        txt(sl, 0.33, by+0.62, 2.98, 0.16, sub, size=8, color=C['WHITE'])

    txt(sl, 3.56, 0.68, 6.35, 0.14, "TOP 10 STATIONS BY TOTAL FUELS REVENUE (PKR M)", size=7.5, bold=True, color=C['GREY'])
    top10 = d['top10']
    chart_x = 3.56; chart_y = 0.86
    max_rev = top10[0]['total'] if top10 else 1
    bar_max_w = 5.60
    row_h = 0.24; row_gap = 0.03

    for si, stn in enumerate(top10[:10]):
        ry = chart_y + si * (row_h + row_gap)
        name_short = str(stn['name'])[:22]
        txt(sl, chart_x, ry+0.02, 1.75, row_h-0.04, name_short, size=7, color=C['DNAV2'])
        bx_start = chart_x + 1.78
        for cat in CAT_ORDER:
            rev = stn['by_cat'].get(cat, 0)
            if rev > 0:
                seg_w = rev/max_rev * bar_max_w
                rect(sl, bx_start, ry+0.02, seg_w, row_h-0.04, CAT_COLORS.get(cat, C['GREY3']))
                bx_start += seg_w
        txt(sl, bx_start+0.06, ry+0.02, 0.90, row_h-0.04, fmt_m(stn['total']), size=7, color=C['GREY2'])

    leg_y = chart_y + 10*(row_h+row_gap) + 0.04
    for li, cat in enumerate(CAT_ORDER):
        lx = chart_x + li * 1.25
        rect(sl, lx, leg_y, 0.12, 0.10, CAT_COLORS.get(cat, C['GREY3']))
        txt(sl, lx+0.15, leg_y, 1.08, 0.12, cat, size=6.5, color=C['GREY'])

    ins_y = leg_y + 0.18
    rect(sl, 3.56, ins_y, 6.25, 0.52, C['BG_YELLOW'])
    rect(sl, 3.58, ins_y, 0.02, 0.52, C['BROWN'])
    top10_names = [s['name'][:20] for s in top10[:2]] if top10 else []
    ins_txt = (f"Top 10 stations = {fmt_m(d['top10_rev'])} ({d['top10_pct']:.1f}% of all {d['city']} fuels). "
               f"{top10_names[0] if top10_names else ''} leads at {fmt_m(top10[0]['total']) if top10 else 'N/A'} — "
               f"study this station's category mix as a best-practice model for the network.")
    txt(sl, 3.72, ins_y+0.06, 6.06, 0.40, ins_txt, size=8, color=C['DNAV2'], wrap=True)

# ── SLIDE 5: Priorities ───────────────────────────────────────────────────────
def build_slide5(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, 10, 0.58, C['DNAVY'])
    txt(sl, 0.25, 0.14, 7.5, 0.30, f"{d['city']} — Where to Focus & What to Fix", size=17, bold=True, color=C['WHITE'])
    txt(sl, 7.5, 0.20, 2.30, 0.18, "4 Priorities  |  Estimated Revenue Upside", size=8.5, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    no_r95 = d['no_r95']; r95_avg = d['r95_avg']
    no_hsd = d['no_hsd']; hsd_avg = d['hsd_avg']
    n_low = d['n_low']; n_zero = d['n_zero_stn']
    med = d['med_per_stn']
    avg_low = d['rev_low'] / d['n_low'] if d['n_low'] else 0
    potential3 = max(0, (med - avg_low) * min(n_low, max(int(n_low*0.35), 5)))

    priorities = [
        (C['BG_LIGHT'], C['BLUE'], C['BLUE2'],
         '01', f"Activate {no_r95} Petrol Stations Not Selling R95",
         f"R95 carries the highest margin per litre of any fuel product but {no_r95} petrol-selling stations — "
         f"are not stocking it. Average R95 revenue per active station is {fmt_m(r95_avg)}.",
         fmt_m(no_r95 * r95_avg * 0.5),
         "HIGH MARGIN  Immediate Action"),

        (C['BG_GREEN'], C['GREEN'], C['BLUE2'],
         '02', f"Activate {no_hsd} Stations Not Selling Diesel/HSD",
         f"Diesel is a core volume driver nationally. {no_hsd} stations in {d['city']} are missing out. "
         f"Avg HSD revenue per active station: {fmt_m(hsd_avg)}.",
         fmt_m(no_hsd * hsd_avg * 0.5),
         "VOLUME GROWTH  Cross-Sell"),

        (C['BG_YELLOW'], C['GOLD'], C['BROWN'],
         '03', f"Upgrade {n_low} Low-Tier Stations to Mid-Tier",
         f"{n_low} stations earn below the city median of {fmt_m(med)}. "
         f"If {min(n_low, max(int(n_low*0.35), 5))} of them were lifted to {fmt_m(med)} "
         f"from avg of {fmt_m(avg_low)}, the revenue gain would be significant.",
         fmt_m(potential3) + "+",
         "WALLET SHARE  Sales Force Priority"),

        (C['RED_BG'], C['RED'], C['RED_D'],
         '04', f"Convert {n_zero} Zero-Fuels Stations",
         f"{n_zero} stations are selling ZERO fuels despite being active retail stations. "
         f"Investigate: missing product license, no shelf space, competitor brand only. "
         f"Each conversion at city avg {fmt_m(d['avg_per_stn'])} generates meaningful revenue.",
         fmt_m(n_zero * d['avg_per_stn'] * 0.4) + "+",
         "QUICK WIN  Investigate & Convert"),
    ]

    col_w = 2.28
    for pi, (bg, num_col, title_col, num, title, body, pot, tag) in enumerate(priorities):
        px = 0.19 + pi*(col_w+0.14)
        rect(sl, px, 0.70, col_w, 4.82, bg)
        rect(sl, px+0.03, 0.72, col_w, 0.02, rgb(num_col))
        txt(sl, px+0.17, 0.85, col_w-0.20, 0.44, num, size=28, bold=True, color=num_col)
        txt(sl, px+0.17, 1.38, col_w-0.20, 0.36, title, size=9, bold=True, color=title_col, wrap=True)
        txt(sl, px+0.17, 1.80, col_w-0.20, 1.30, body, size=7.5, color=C['DNAV2'], wrap=True)
        txt(sl, px+0.17, 3.18, col_w-0.20, 0.22, f"Potential: {pot}", size=10, bold=True, color=num_col)
        txt(sl, px+0.17, 3.44, col_w-0.20, 0.22, tag, size=6.5, bold=True, color=C['GREY'])

# ── MAIN ─────────────────────────────────────────────────────────────────────
top_vol = (fuels_all.groupby('CityNorm')['SalesLtr_CY'].sum().sort_values(ascending=False))
target_cities = top_vol.index.tolist()[:10]

print(f"Generating profiles for: {target_cities}")
profiles_dir = os.path.join(os.environ.get('PSO_OUTDIR', 'reports'), 'fuels_city_profiles')
os.makedirs(profiles_dir, exist_ok=True)

for city in target_cities:
    print(f"  Processing {city}…", end=' ', flush=True)
    d = compute_city(city)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    build_slide1(prs, d)
    build_slide2(prs, d)
    build_slide3(prs, d)
    build_slide4(prs, d)
    build_slide5(prs, d)

    out = out_path(f"PSO_{city.replace(' ', '_')}_Fuels_Profile", 'pptx', df,
                    out_dir=profiles_dir)
    prs.save(out)
    print(f"saved → {out}")

print("\nAll profiles generated.")
