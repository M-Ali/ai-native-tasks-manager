"""
Karachi Volume Focus Deck — 7 slides
One story per slide. Objective: increase volumes.
"""
import sys, os
sys.path.insert(0, 'src')
os.environ['PYTHONIOENCODING'] = 'utf-8'

import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.oxml import parse_xml
from pso import ingest

# ── colours ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B,0x2A,0x4A)
GOLD   = RGBColor(0xC9,0xA0,0x30)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LGREY  = RGBColor(0xF2,0xF4,0xF8)
MGREY  = RGBColor(0xD9,0xDC,0xE3)
DBLUE  = RGBColor(0x2E,0x5B,0x9A)
LBLUE  = RGBColor(0xBD,0xD7,0xEE)
GREEN  = RGBColor(0x37,0x5B,0x25)
LGREEN = RGBColor(0xC6,0xEF,0xCE)
RED    = RGBColor(0xC0,0x00,0x00)
LRED   = RGBColor(0xFF,0xC7,0xCE)
YELL   = RGBColor(0xFF,0xC0,0x00)
LYELL  = RGBColor(0xFF,0xEB,0x9C)
ORANGE = RGBColor(0xFF,0x82,0x00)

SW, SH = Inches(13.33), Inches(7.50)

def rhex(c): return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

# ── primitive helpers ─────────────────────────────────────────────────────────
def rect(slide, l,t,w,h, fill, line=None):
    s = slide.shapes.add_shape(1,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def txt(slide, l,t,w,h, text, size=10, bold=False, fg=NAVY, bg=None,
        align=PP_ALIGN.LEFT, italic=False):
    if bg: rect(slide, l,t,w,h, bg)
    tb = slide.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.word_wrap=True
    p  = tf.paragraphs[0]; p.alignment=align
    r  = p.add_run(); r.text=str(text)
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=fg; r.font.italic=italic
    return tb

def cell_style(cell, text, size=8, bold=False, fg=NAVY, bg=None,
               align=PP_ALIGN.CENTER, italic=False):
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]; p.alignment=align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=fg; r.font.italic=italic
    if bg:
        pr = cell._tc.get_or_add_tcPr()
        sf = parse_xml(
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{rhex(bg)}"/></a:solidFill>')
        for ch in list(pr):
            if 'Fill' in ch.tag or 'fill' in ch.tag.lower(): pr.remove(ch)
        pr.append(sf)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def hdr(slide, title, subtitle=""):
    rect(slide, Inches(0),Inches(0),SW,Inches(0.75),NAVY)
    rect(slide, Inches(0),Inches(0.75),SW,Inches(0.04),GOLD)
    txt(slide, Inches(0.2),Inches(0.06),Inches(10),Inches(0.38),
        title, size=15, bold=True, fg=GOLD)
    if subtitle:
        txt(slide, Inches(0.2),Inches(0.44),Inches(11),Inches(0.28),
            subtitle, size=9, fg=MGREY)
    txt(slide, Inches(10.8),Inches(0.24),Inches(2.3),Inches(0.25),
        "PSO Karachi  |  10M FY26", size=7.5, fg=MGREY, align=PP_ALIGN.RIGHT)

def footer(slide, text):
    rect(slide, Inches(0),Inches(7.18),SW,Inches(0.32),NAVY)
    txt(slide, Inches(0.15),Inches(7.20),Inches(13),Inches(0.26),
        text, size=7.5, fg=MGREY, italic=True)

def section_bar(slide, l,t,w, label, bg=DBLUE):
    rect(slide, l,t,w,Inches(0.27),bg)
    txt(slide, l+Inches(0.08),t+Inches(0.02),w-Inches(0.1),Inches(0.24),
        label, size=8.5, bold=True, fg=WHITE)

def insight_box(slide, l,t,w,h, title, body, title_bg=GOLD, body_bg=LYELL):
    rect(slide, l,t,w,Inches(0.28),title_bg)
    txt(slide, l+Inches(0.08),t+Inches(0.02),w-Inches(0.1),Inches(0.24),
        title, size=8.5, bold=True, fg=NAVY)
    rect(slide, l,t+Inches(0.28),w,h-Inches(0.28),body_bg)
    txt(slide, l+Inches(0.1),t+Inches(0.32),w-Inches(0.15),h-Inches(0.38),
        body, size=8.5, fg=NAVY)

def kpi_box(slide, l,t,w,h, label, value, sub="", val_color=WHITE, bg=DBLUE):
    rect(slide, l,t,w,h,bg)
    txt(slide, l+Inches(0.07),t+Inches(0.04),w-Inches(0.1),Inches(0.22),
        label, size=7, fg=LGREY, align=PP_ALIGN.CENTER)
    txt(slide, l+Inches(0.05),t+Inches(0.24),w-Inches(0.08),Inches(0.38),
        value, size=15, bold=True, fg=val_color, align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, l+Inches(0.05),t+Inches(0.60),w-Inches(0.08),Inches(0.22),
            sub, size=7.5, fg=LGREY, align=PP_ALIGN.CENTER, italic=True)

def bar_chart(slide, l,t,w,h, cats, cy, ly, cy_lbl="CY 10M FY26", ly_lbl="LY 10M FY25"):
    cd = ChartData()
    cd.categories = cats
    cd.add_series(cy_lbl, cy)
    cd.add_series(ly_lbl, ly)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,l,t,w,h,cd)
    ch = cf.chart
    ch.has_legend=True
    ch.legend.position=XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout=False
    try:
        ch.series[0].format.fill.solid()
        ch.series[0].format.fill.fore_color.rgb = NAVY
        ch.series[1].format.fill.solid()
        ch.series[1].format.fill.fore_color.rgb = MGREY
        ch.category_axis.tick_labels.font.size = Pt(8)
        ch.value_axis.tick_labels.font.size    = Pt(7)
        ch.has_title=False
    except: pass
    return cf

def fmt_vol(v, dec=1): return f"{v:,.{dec}f}" if v is not None else "—"
def fmt_pct(v, dec=1):
    if v is None: return "—"
    return f"+{v:.{dec}f}%" if v>=0 else f"{v:.{dec}f}%"
def chg_bg(v):
    if v is None: return MGREY
    if v >  2: return LGREEN
    if v < -2: return LRED
    return LYELL
def safe_pct(cy,ly): return (cy-ly)/abs(ly)*100 if ly and ly!=0 else None

# ═══════════════════════════════════════════════════════════════════════════════
# LOAD & PREPARE
# ═══════════════════════════════════════════════════════════════════════════════
print("Loading …")
df, _ = ingest.load(DATA_FILE := "data/input/Working File Retail Fuels Data.xlsx")
retail  = df[df["IsRetail"] & ~df["IsInternational"]].copy()
karachi = retail[retail["CityNorm"]=="Karachi"].copy()

# top-10 cities for benchmark
top10 = retail.groupby("CityNorm")["SalesLtr_CY"].sum().nlargest(10).index.tolist()
retail10 = retail[retail["CityNorm"].isin(top10)]

# ── city-level volumes (ML) ────────────────────────────────────────────────────
tot_cy = karachi["SalesLtr_CY"].sum()/1e6
tot_ly = karachi["SalesLtr_LY"].sum()/1e6
tot_chg = safe_pct(tot_cy, tot_ly)
nat_vol  = retail["SalesLtr_CY"].sum()/1e6
k_nat_sh = tot_cy/nat_vol*100

SEGS = ["Diesel","Petrol","Lubricants"]
seg_data = {}
for seg in SEGS:
    s = karachi[karachi["FuelSegment"]==seg]
    cy = s["SalesLtr_CY"].sum()/1e6
    ly = s["SalesLtr_LY"].sum()/1e6
    seg_data[seg] = dict(cy=cy, ly=ly, chg=safe_pct(cy,ly),
                         sh=cy/tot_cy*100 if tot_cy else 0,
                         stations=s["Customer Number"].nunique())

# ── sub-product volumes (ML) ───────────────────────────────────────────────────
DIESEL_P = ["HSD","LDO","SLUDGE-HSD","SLUDGE"]
PETROL_P = ["PMG","R95"]
LUBE_C   = ["DEO","PCMO","MCO","LOW GRADE","INDUSTRIAL GRADE","Greases","OTHERS"]

sub = {}
for p in DIESEL_P:
    s = karachi[(karachi["FuelSegment"]=="Diesel")&(karachi["ProductCategory"]==p)]
    cy=s["SalesLtr_CY"].sum()/1e6; ly=s["SalesLtr_LY"].sum()/1e6
    sub[p]=dict(grp="Diesel",cy=cy,ly=ly,chg=safe_pct(cy,ly),stns=s["Customer Number"].nunique())
for p in PETROL_P:
    s = karachi[(karachi["FuelSegment"]=="Petrol")&(karachi["ProductCategory"]==p)]
    cy=s["SalesLtr_CY"].sum()/1e6; ly=s["SalesLtr_LY"].sum()/1e6
    sub[p]=dict(grp="Petrol",cy=cy,ly=ly,chg=safe_pct(cy,ly),stns=s["Customer Number"].nunique())
for c in LUBE_C:
    s = karachi[(karachi["FuelSegment"]=="Lubricants")&(karachi["LubeCategory"]==c)]
    cy=s["SalesLtr_CY"].sum()/1e6; ly=s["SalesLtr_LY"].sum()/1e6
    sub[c]=dict(grp="Lubricants",cy=cy,ly=ly,chg=safe_pct(cy,ly),stns=s["Customer Number"].nunique())

# ── station-level ─────────────────────────────────────────────────────────────
stn = (karachi.groupby("Customer Number",as_index=False)
       .agg(name=("Name 1","first"),
            vol_cy=("SalesLtr_CY","sum"),
            vol_ly=("SalesLtr_LY","sum"))
       .sort_values("vol_cy",ascending=False).reset_index(drop=True))
stn["vol_cy_ml"] = stn["vol_cy"]/1e6
stn["vol_ly_ml"] = stn["vol_ly"]/1e6
stn["chg"]       = stn.apply(lambda r: safe_pct(r["vol_cy"],r["vol_ly"]),axis=1)
stn["sh"]        = stn["vol_cy"]/stn["vol_cy"].sum()*100
stn["cum_sh"]    = stn["sh"].cumsum()
stn["active"]    = stn["vol_cy"]>0

n_tot   = len(stn)
n_inact = int((~stn["active"]).sum())
n_active= n_tot - n_inact
n80     = int((stn["cum_sh"]<80).sum())+1  # stations needed for 80%

# ── product map per station ───────────────────────────────────────────────────
stn_prods = {}
for cnum, grp in karachi.groupby("Customer Number"):
    fp = set(grp[grp["FuelSegment"]!="Lubricants"]["ProductCategory"].dropna())
    lp = set(grp[grp["FuelSegment"]=="Lubricants"]["LubeCategory"].dropna())
    stn_prods[cnum] = dict(
        has_hsd  ="HSD" in fp,
        has_ldo  ="LDO" in fp,
        has_pmg  ="PMG" in fp,
        has_r95  ="R95" in fp,
        has_diesel=bool({"HSD","LDO"}&fp),
        has_petrol=bool({"PMG","R95"}&fp),
        has_lubes =bool(lp),
        has_deo   ="DEO"  in lp,
        has_pcmo  ="PCMO" in lp,
        has_mco   ="MCO"  in lp,
        has_lg    ="LOW GRADE" in lp,
        lube_cats =lp, fuel_prods=fp,
    )

# ── product coverage counts ───────────────────────────────────────────────────
pp = list(stn_prods.values())
n_diesel      = sum(1 for p in pp if p["has_diesel"])
n_petrol      = sum(1 for p in pp if p["has_petrol"])
n_lubes       = sum(1 for p in pp if p["has_lubes"])
n_all3        = sum(1 for p in pp if p["has_diesel"] and p["has_petrol"] and p["has_lubes"])
n_d_and_p     = sum(1 for p in pp if p["has_diesel"] and p["has_petrol"] and not p["has_lubes"])
n_d_only      = sum(1 for p in pp if p["has_diesel"] and not p["has_petrol"] and not p["has_lubes"])
n_p_only      = sum(1 for p in pp if p["has_petrol"] and not p["has_diesel"] and not p["has_lubes"])
n_d_no_p      = sum(1 for p in pp if p["has_diesel"] and not p["has_petrol"])
n_p_no_r95    = sum(1 for p in pp if p["has_petrol"] and not p["has_r95"] and p["has_pmg"])
n_fuel_no_lub = sum(1 for p in pp if (p["has_diesel"] or p["has_petrol"]) and not p["has_lubes"])
n_r95         = sum(1 for p in pp if p["has_r95"])
n_pmg_only    = sum(1 for p in pp if p["has_pmg"] and not p["has_r95"])

# ── volume lost: diesel-only stations (no petrol) ────────────────────────────
diesel_no_petrol_custs = [c for c,p in stn_prods.items() if p["has_diesel"] and not p["has_petrol"]]
d_no_p_stns = stn[stn["Customer Number"].isin(diesel_no_petrol_custs)]
d_no_p_vol  = d_no_p_stns["vol_cy_ml"].sum()

# ── volume in lubes-missing stations ─────────────────────────────────────────
fuel_no_lub_custs = [c for c,p in stn_prods.items()
                     if (p["has_diesel"] or p["has_petrol"]) and not p["has_lubes"]]
fnl_vol = stn[stn["Customer Number"].isin(fuel_no_lub_custs)]["vol_cy_ml"].sum()

# avg lube vol per lube-selling station (for potential calc)
lub_custs = [c for c,p in stn_prods.items() if p["has_lubes"]]
lube_vols = karachi[(karachi["FuelSegment"]=="Lubricants") &
                    (karachi["Customer Number"].isin(lub_custs))]
avg_lube_per_stn = lube_vols["SalesLtr_CY"].sum()/1e6 / max(len(lub_custs),1)

# ── benchmarks ───────────────────────────────────────────────────────────────
# Top-10 city averages
t10_vols  = {c: retail[retail["CityNorm"]==c]["SalesLtr_CY"].sum()/1e6 for c in top10}
t10_stns  = {c: retail[retail["CityNorm"]==c]["Customer Number"].nunique() for c in top10}
t10_avg_city_vol = np.mean(list(t10_vols.values()))
t10_avg_per_stn  = np.mean([t10_vols[c]/t10_stns[c] for c in top10])
k_vol_per_stn    = tot_cy / n_tot

# active-station average
k_active_vol_per_stn = stn[stn["active"]]["vol_cy_ml"].mean()
top25_threshold = stn.head(max(1,n_tot//4))["vol_cy_ml"].min()
top25_avg = stn.head(max(1,n_tot//4))["vol_cy_ml"].mean()
median_vol = stn[stn["active"]]["vol_cy_ml"].median()

# ── inactive recovery potential ──────────────────────────────────────────────
inact = stn[~stn["active"]]
inact_ly_vol = inact["vol_ly_ml"].sum()
recovery_inact = inact_ly_vol * 0.65  # 65% realistic reactivation

# under-performers (active but vol_cy < 50% of vol_ly)
underperf = stn[(stn["active"]) & (stn["vol_ly"]>0) & (stn["vol_cy"]<stn["vol_ly"]*0.5)]
underperf_gap = (underperf["vol_ly_ml"]*0.75 - underperf["vol_cy_ml"]).sum()

# lube cross-sell potential
lube_potential = len(fuel_no_lub_custs) * avg_lube_per_stn

# petrol add-on potential for diesel-only stations
petrol_avg_at_dual = stn[stn["Customer Number"].isin(
    [c for c,p in stn_prods.items() if p["has_diesel"] and p["has_petrol"]])]["vol_cy_ml"].mean()
# assume diesel-only stations could add petrol at ~30% of current vol
petrol_addon_potential = d_no_p_vol * 0.25

total_potential = recovery_inact + underperf_gap + lube_potential + petrol_addon_potential

print(f"Karachi: {tot_cy:.1f}ML CY | {n80} stns for 80% | "
      f"{n_inact} inactive | potential {total_potential:.1f}ML")

# ═══════════════════════════════════════════════════════════════════════════════
# BUILD PPT
# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width=SW; prs.slide_height=SH
BL = prs.slide_layouts[6]  # blank

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Volume Scorecard
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BL)
hdr(sl,"KARACHI — VOLUME SCORECARD",
    "How does Karachi stand today? Is volume growing or shrinking, and where?")

# top KPIs
chg_col = LGREEN if (tot_chg or 0)>0 else LRED
kpis=[
    ("Total Vol CY (ML)",  f"{fmt_vol(tot_cy)}",   f"LY: {fmt_vol(tot_ly)} ML",   WHITE, NAVY),
    ("YoY Volume Change",  fmt_pct(tot_chg),         "vs same period last year",   WHITE if (tot_chg or 0)>=0 else WHITE,
                                                      GREEN if (tot_chg or 0)>=0 else RED),
    ("National Vol Share", f"{k_nat_sh:.1f}%",       "of PSO retail (all cities)", WHITE, DBLUE),
    ("Vol / Station CY",   f"{k_vol_per_stn:.2f}ML", f"10-city avg: {t10_avg_per_stn:.2f}ML", WHITE,
                            GREEN if k_vol_per_stn>=t10_avg_per_stn else RED),
    ("Active Stations",    f"{n_active}/{n_tot}",    f"{n_inact} inactive", WHITE, DBLUE),
]
kw = Inches(2.5); kh = Inches(0.95)
for ki,(lbl,val,sub_,vc,bg) in enumerate(kpis):
    kpi_box(sl, Inches(0.18)+ki*kw, Inches(0.83), kw-Inches(0.07), kh,
            lbl, val, sub_, WHITE, bg)

# segment bar chart (left)
bar_chart(sl,
    Inches(0.18),Inches(1.92),Inches(5.5),Inches(2.9),
    SEGS,
    [seg_data[s]["cy"] for s in SEGS],
    [seg_data[s]["ly"] for s in SEGS],
)
txt(sl,Inches(0.18),Inches(1.85),Inches(5.5),Inches(0.22),
    "Volume by Segment (ML) — CY vs LY",size=8,bold=True,fg=NAVY)

# segment detail table (right)
section_bar(sl,Inches(5.9),Inches(1.85),Inches(7.25),"Segment Detail")
tbl=sl.shapes.add_table(5,6,Inches(5.9),Inches(2.15),Inches(7.25),Inches(1.75)).table
for ci,h in enumerate(["Segment","Stations","Vol CY (ML)","Vol LY (ML)","Change","City Share%"]):
    cell_style(tbl.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
cws=[Inches(1.4),Inches(0.9),Inches(1.3),Inches(1.3),Inches(1.0),Inches(1.35)]
for ci,w in enumerate(cws): tbl.columns[ci].width=w
for ri,seg in enumerate(SEGS+["Total"],1):
    if seg=="Total":
        d=dict(cy=tot_cy,ly=tot_ly,chg=tot_chg,sh=100,stations=n_tot)
        bg=MGREY
    else:
        d=seg_data[seg]; bg=LGREY if ri%2==0 else WHITE
    vals=[seg,str(d["stations"]),fmt_vol(d["cy"]),fmt_vol(d["ly"]),fmt_pct(d["chg"]),
          f"{d.get('sh',d['cy']/tot_cy*100 if tot_cy else 0):.1f}%"]
    for ci,v in enumerate(vals):
        cbg = chg_bg(d["chg"]) if ci==4 else bg
        cell_style(tbl.cell(ri,ci),v,size=8.5,bold=(seg=="Total"),
                   fg=NAVY,bg=cbg,align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
    tbl.rows[ri].height=Inches(0.34)
tbl.rows[0].height=Inches(0.28)

# peer comparison table
section_bar(sl,Inches(5.9),Inches(4.1),Inches(7.25),"Karachi vs Top-10 City Benchmarks")
peers=[
    ("Karachi Total Volume (ML)",        f"{tot_cy:.1f}", f"10-city avg: {t10_avg_city_vol:.1f}ML",
     LGREEN if tot_cy>=t10_avg_city_vol else LRED),
    ("Vol per Station (ML)",             f"{k_vol_per_stn:.2f}",f"10-city avg: {t10_avg_per_stn:.2f}ML",
     LGREEN if k_vol_per_stn>=t10_avg_per_stn else LRED),
    ("Active station avg vol (ML)",      f"{k_active_vol_per_stn:.2f}", "per active station",LGREY),
    ("Top-25% station avg vol (ML)",     f"{top25_avg:.2f}",  "benchmark for star performers",LGREY),
    ("Median active station vol (ML)",   f"{median_vol:.2f}", "50th percentile baseline",LGREY),
]
ptbl=sl.shapes.add_table(len(peers)+1,3,Inches(5.9),Inches(4.4),Inches(7.25),Inches(2.1)).table
for ci,h in enumerate(["Metric","Karachi","Context"]):
    cell_style(ptbl.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
ptbl.columns[0].width=Inches(3.8); ptbl.columns[1].width=Inches(1.4); ptbl.columns[2].width=Inches(2.05)
for ri,(lbl,val,ctx,bg) in enumerate(peers,1):
    cell_style(ptbl.cell(ri,0),lbl,size=8,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
    cell_style(ptbl.cell(ri,1),val,size=9,bold=True,fg=NAVY,bg=bg)
    cell_style(ptbl.cell(ri,2),ctx,size=7.5,fg=NAVY,bg=bg,italic=True)
    ptbl.rows[ri].height=Inches(0.33)
ptbl.rows[0].height=Inches(0.27)

insight_box(sl,Inches(0.18),Inches(4.92),Inches(5.5),Inches(1.95),
    "READ THIS FIRST",
    (f"Karachi delivered {fmt_vol(tot_cy)}ML in 10M FY26 vs {fmt_vol(tot_ly)}ML in LY "
     f"— a {fmt_pct(tot_chg)} change. "
     f"Diesel dominates at {seg_data['Diesel']['sh']:.0f}% of volume. "
     f"Petrol is {seg_data['Petrol']['sh']:.0f}% — important because petrol yields better margin/ltr. "
     f"Lubes contribute only {seg_data['Lubricants']['sh']:.1f}% despite being highest-value. "
     f"Per-station volume of {k_vol_per_stn:.2f}ML "
     f"{'exceeds' if k_vol_per_stn>=t10_avg_per_stn else 'is below'} the 10-city average of {t10_avg_per_stn:.2f}ML."))

footer(sl,"All volumes in Million Litres (ML).  CY = 10M FY26  |  LY = 10M FY25  |  Source: PSO Working File")
notes(sl,
    f"SLIDE 1 — KARACHI VOLUME SCORECARD\n\n"
    f"Total CY: {tot_cy:.2f}ML  |  LY: {tot_ly:.2f}ML  |  Chg: {fmt_pct(tot_chg)}\n"
    f"National share: {k_nat_sh:.1f}%\n"
    f"Diesel: {seg_data['Diesel']['cy']:.2f}ML ({seg_data['Diesel']['sh']:.1f}%) {fmt_pct(seg_data['Diesel']['chg'])}\n"
    f"Petrol: {seg_data['Petrol']['cy']:.2f}ML ({seg_data['Petrol']['sh']:.1f}%) {fmt_pct(seg_data['Petrol']['chg'])}\n"
    f"Lubes:  {seg_data['Lubricants']['cy']:.2f}ML ({seg_data['Lubricants']['sh']:.1f}%) {fmt_pct(seg_data['Lubricants']['chg'])}\n"
    f"Total stations: {n_tot}  |  Active: {n_active}  |  Inactive: {n_inact}\n"
    f"Vol/station: {k_vol_per_stn:.2f}ML vs 10-city avg {t10_avg_per_stn:.2f}ML\n"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The 80% Rule
# ─────────────────────────────────────────────────────────────────────────────
sl=prs.slides.add_slide(BL)
hdr(sl,f"THE 80% RULE: {n80} STATIONS CARRY KARACHI",
    f"Only {n80} of {n_tot} stations ({n80/n_tot*100:.0f}% of fleet) generate 80% of Karachi's volume. "
    f"The other {n_tot-n80} stations share the remaining 20%.")

# KPI strip
k2=[
    ("Stations for 80% Vol", str(n80),           f"{n80/n_tot*100:.0f}% of fleet"),
    ("Remaining stations",   str(n_tot-n80),      f"share only 20% of volume"),
    ("Top station vol",      f"{stn.iloc[0]['vol_cy_ml']:.1f}ML", stn.iloc[0]["name"][:25]),
    ("Top 5 stations",       f"{stn.head(5)['sh'].sum():.1f}%",   "of city volume"),
    ("Top 10 stations",      f"{stn.head(10)['sh'].sum():.1f}%",  "of city volume"),
    ("Avg vol: top 80%",     f"{stn.head(n80)['vol_cy_ml'].mean():.2f}ML","per station"),
    ("Avg vol: bottom 20%",  f"{stn.tail(n_tot-n80)['vol_cy_ml'].mean():.2f}ML","per station"),
]
kw2=SW/len(k2)
for ki,(lbl,val,sub_) in enumerate(k2):
    bg2=RED if ki==1 else (GREEN if ki==0 else DBLUE)
    rect(sl, ki*kw2,Inches(0.80),kw2,Inches(0.75),bg2)
    txt(sl, ki*kw2+Inches(0.04),Inches(0.81),kw2-Inches(0.06),Inches(0.22),lbl,size=6.5,fg=MGREY,align=PP_ALIGN.CENTER)
    txt(sl, ki*kw2+Inches(0.03),Inches(1.02),kw2-Inches(0.05),Inches(0.30),val,size=12,bold=True,fg=WHITE,align=PP_ALIGN.CENTER)
    txt(sl, ki*kw2+Inches(0.03),Inches(1.31),kw2-Inches(0.05),Inches(0.20),sub_,size=6.5,fg=LGREY,align=PP_ALIGN.CENTER,italic=True)

# Station ranking table (left — top 30)
section_bar(sl,Inches(0.15),Inches(1.65),Inches(7.9),"Station Ranking — Volume CY (top 30 shown, 80% threshold highlighted in blue)")
show=min(30,n_tot)
rtbl=sl.shapes.add_table(show+1,6,Inches(0.15),Inches(1.95),Inches(7.9),Inches(5.1)).table
for ci,h in enumerate(["#","Station","Vol CY (ML)","Vol LY (ML)","Chg%","Cum%"]):
    cell_style(rtbl.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=NAVY)
rtbl.columns[0].width=Inches(0.3); rtbl.columns[1].width=Inches(2.8)
for ci in range(2,6): rtbl.columns[ci].width=Inches(1.2)
rtbl.rows[0].height=Inches(0.25)
reached=False
for ri,(idx,row) in enumerate(stn.head(show).iterrows(),1):
    is_key = row["cum_sh"]<=80 or (not reached and row["cum_sh"]>80)
    if row["cum_sh"]>80 and not reached: reached=True
    bg=LBLUE if is_key and row["active"] else (LRED if not row["active"] else WHITE)
    vals=[str(ri),str(row["name"])[:38],
          fmt_vol(row["vol_cy_ml"]),fmt_vol(row["vol_ly_ml"]),
          fmt_pct(row["chg"]),"—" if not row["active"] else f"{row['cum_sh']:.1f}%"]
    for ci,v in enumerate(vals):
        cbg=chg_bg(row["chg"]) if ci==4 and row["active"] else bg
        cell_style(rtbl.cell(ri,ci),v,size=7,bold=(ri<=5),
                   fg=RED if not row["active"] else NAVY,
                   bg=cbg,align=PP_ALIGN.LEFT if ci==1 else PP_ALIGN.CENTER)
    rtbl.rows[ri].height=Inches(0.175)

# Quintile analysis panel (right)
section_bar(sl,Inches(8.2),Inches(1.65),Inches(4.95),"Volume Distribution by Station Tier")
q_n=max(1,n_tot//5)
qtiers=[("Top 20%",stn.head(q_n)),("21–40%",stn.iloc[q_n:2*q_n]),
        ("41–60%",stn.iloc[2*q_n:3*q_n]),("61–80%",stn.iloc[3*q_n:4*q_n]),
        ("Bottom 20%",stn.iloc[4*q_n:])]
qtbl=sl.shapes.add_table(6,4,Inches(8.2),Inches(1.95),Inches(4.95),Inches(1.6)).table
for ci,h in enumerate(["Tier","Stns","Vol (ML)","Vol Share"]):
    cell_style(qtbl.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
qtbl.rows[0].height=Inches(0.25)
qbgs=[NAVY,DBLUE,LBLUE,LGREY,WHITE]; qfgs=[WHITE,WHITE,NAVY,NAVY,NAVY]
for ri,(lbl,grp) in enumerate(qtiers,1):
    gv=grp["vol_cy_ml"].sum(); gsh=gv/tot_cy*100 if tot_cy else 0
    cell_style(qtbl.cell(ri,0),lbl,size=8,bold=True,fg=qfgs[ri-1],bg=qbgs[ri-1],align=PP_ALIGN.LEFT)
    cell_style(qtbl.cell(ri,1),str(len(grp)),size=8,fg=qfgs[ri-1],bg=qbgs[ri-1])
    cell_style(qtbl.cell(ri,2),fmt_vol(gv),size=8,bold=True,fg=qfgs[ri-1],bg=qbgs[ri-1])
    cell_style(qtbl.cell(ri,3),f"{gsh:.1f}%",size=8,bold=True,fg=qfgs[ri-1],bg=qbgs[ri-1])
    qtbl.rows[ri].height=Inches(0.25)
for ci,w in enumerate([Inches(1.2),Inches(0.7),Inches(1.5),Inches(1.55)]): qtbl.columns[ci].width=w

insight_box(sl,Inches(8.2),Inches(3.7),Inches(4.95),Inches(3.35),
    "WHAT THIS MEANS FOR VOLUME GROWTH",
    (f"The top {n80} stations are already managed — they will deliver predictable volume. "
     f"The growth opportunity lies in the BOTTOM {n_tot-n80} stations which together produce "
     f"only {stn.tail(n_tot-n80)['vol_cy_ml'].sum():.1f}ML.\n\n"
     f"ACTION: If the bottom {n_tot-n80} stations each lifted volume by just "
     f"{stn.tail(n_tot-n80)['vol_cy_ml'].mean()*0.3:.2f}ML (30%), "
     f"Karachi gains +{stn.tail(n_tot-n80)['vol_cy_ml'].sum()*0.3:.1f}ML.\n\n"
     f"RISK: Top {min(5,n80)} stations = {stn.head(5)['sh'].sum():.0f}% of city volume. "
     f"Losing even one top station is a major city-level event. Protect them first."))

footer(sl,f"Blue = stations within 80% volume threshold  |  Red = inactive  |  "
         f"Bottom {n_tot-n80} stations average {stn.tail(n_tot-n80)['vol_cy_ml'].mean():.2f}ML each")
notes(sl,
    f"SLIDE 2 — 80% RULE\n"
    f"Stations for 80% volume: {n80} ({n80/n_tot*100:.0f}% of fleet)\n"
    f"Top 5 stations: {stn.head(5)['sh'].sum():.1f}% of city vol\n"
    f"Top 10 stations: {stn.head(10)['sh'].sum():.1f}% of city vol\n"
    f"Avg vol top 80%: {stn.head(n80)['vol_cy_ml'].mean():.2f}ML\n"
    f"Avg vol bottom 20%: {stn.tail(n_tot-n80)['vol_cy_ml'].mean():.2f}ML\n"
    f"Ratio top:bottom tier = {stn.head(n80)['vol_cy_ml'].mean()/max(0.01,stn.tail(n_tot-n80)['vol_cy_ml'].mean()):.0f}x\n"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Star Performers
# ─────────────────────────────────────────────────────────────────────────────
sl=prs.slides.add_slide(BL)
top25n=max(1,n_tot//4)
hdr(sl,"STAR PERFORMERS — WHAT GOOD LOOKS LIKE",
    f"Top 25% stations ({top25n} outlets) average {top25_avg:.2f}ML each. "
    f"If all active stations hit just 60% of this benchmark, Karachi adds +{(top25_avg*0.6-median_vol)*n_active:.1f}ML.")

# top 15 table (left)
section_bar(sl,Inches(0.15),Inches(0.85),Inches(8.0),"Top 15 Stations — Volume Champions  (products shown = what they actually sell)")
stars=stn.head(15).copy()
stars["prods"]=stars["Customer Number"].map(
    lambda c: ("/".join(sorted(
        (["HSD"] if stn_prods.get(c,{}).get("has_hsd") else [])+
        (["LDO"] if stn_prods.get(c,{}).get("has_ldo") else [])+
        (["PMG"] if stn_prods.get(c,{}).get("has_pmg") else [])+
        (["R95"] if stn_prods.get(c,{}).get("has_r95") else [])+
        (["LUBES"] if stn_prods.get(c,{}).get("has_lubes") else [])
    )))
)
stbl=sl.shapes.add_table(16,7,Inches(0.15),Inches(1.13),Inches(8.0),Inches(5.75)).table
for ci,h in enumerate(["#","Station","Vol CY (ML)","Vol LY (ML)","YoY Chg","Products Sold","Cum%"]):
    cell_style(stbl.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=NAVY)
stbl.columns[0].width=Inches(0.3); stbl.columns[1].width=Inches(2.8)
stbl.columns[2].width=Inches(1.1); stbl.columns[3].width=Inches(1.1)
stbl.columns[4].width=Inches(0.9); stbl.columns[5].width=Inches(1.2); stbl.columns[6].width=Inches(0.6)
stbl.rows[0].height=Inches(0.28)
for ri,(idx,row) in enumerate(stars.iterrows(),1):
    prod_str=str(row.get("prods",""))
    n_prods=len([x for x in prod_str.split("/") if x])
    bg=LGREEN if row["chg"] and row["chg"]>5 else (LYELL if row["chg"] and row["chg"]>0 else LRED if row["chg"] and row["chg"]<-5 else LGREY)
    vals=[str(ri),str(row["name"])[:38],fmt_vol(row["vol_cy_ml"]),fmt_vol(row["vol_ly_ml"]),
          fmt_pct(row["chg"]),prod_str,f"{row['cum_sh']:.0f}%"]
    for ci,v in enumerate(vals):
        cbg=chg_bg(row["chg"]) if ci==4 else (bg if ci>0 else LGREY)
        cell_style(stbl.cell(ri,ci),v,size=7,bold=(ri<=5),
                   fg=NAVY,bg=cbg,align=PP_ALIGN.LEFT if ci in(1,5) else PP_ALIGN.CENTER)
    stbl.rows[ri].height=Inches(0.37)

# right panel — star analysis
section_bar(sl,Inches(8.3),Inches(0.85),Inches(4.85),"What Makes a Star Performer?")

# product count vs volume relationship
prod_vol={"1 product":[],"2 products":[],"3+ products":[]}
for _,row in stn.iterrows():
    if not row["active"]: continue
    sp=stn_prods.get(row["Customer Number"],{})
    n_p=(1 if sp.get("has_diesel") else 0)+(1 if sp.get("has_petrol") else 0)+(1 if sp.get("has_lubes") else 0)
    k="1 product" if n_p<=1 else "2 products" if n_p==2 else "3+ products"
    prod_vol[k].append(row["vol_cy_ml"])

pv_stats=[(k,np.mean(v),len(v)) for k,v in prod_vol.items() if v]
pvtbl=sl.shapes.add_table(len(pv_stats)+2,3,Inches(8.3),Inches(1.13),Inches(4.85),Inches(1.5)).table
cell_style(pvtbl.cell(0,0),"Product Breadth",size=8,bold=True,fg=WHITE,bg=NAVY,align=PP_ALIGN.LEFT)
cell_style(pvtbl.cell(0,1),"Avg Vol/Stn",size=8,bold=True,fg=WHITE,bg=NAVY)
cell_style(pvtbl.cell(0,2),"# Stations",size=8,bold=True,fg=WHITE,bg=NAVY)
pvtbl.rows[0].height=Inches(0.26)
for ri,(k,avg,cnt) in enumerate(pv_stats,1):
    bg=LGREEN if ri==len(pv_stats) else (LYELL if ri==2 else LRED)
    cell_style(pvtbl.cell(ri,0),k,size=8.5,bold=True,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
    cell_style(pvtbl.cell(ri,1),f"{avg:.2f}ML",size=8.5,bold=True,fg=NAVY,bg=bg)
    cell_style(pvtbl.cell(ri,2),str(cnt),size=8.5,fg=NAVY,bg=bg)
    pvtbl.rows[ri].height=Inches(0.32)
pvtbl.columns[0].width=Inches(2.2); pvtbl.columns[1].width=Inches(1.5); pvtbl.columns[2].width=Inches(1.15)

# separator insight
mul_prod_avg=np.mean(prod_vol.get("3+ products",[0]))
one_prod_avg=np.mean(prod_vol.get("1 product",[0.01]))
insight_box(sl,Inches(8.3),Inches(2.75),Inches(4.85),Inches(2.0),
    "KEY FINDING: MORE PRODUCTS = MORE VOLUME",
    (f"Stations selling 3+ product groups average "
     f"{mul_prod_avg:.2f}ML vs {one_prod_avg:.2f}ML for single-product stations.\n\n"
     f"That is a {mul_prod_avg/max(one_prod_avg,0.01):.1f}x volume multiplier.\n\n"
     f"Every single-product station that is converted to multi-product "
     f"is a direct volume uplift opportunity."))

# benchmark panel
section_bar(sl,Inches(8.3),Inches(4.85),Inches(4.85),"Volume Benchmarks (Active Stations)")
bmarks=[
    ("Top 10% stations avg",  f"{stn.head(max(1,n_tot//10))['vol_cy_ml'].mean():.2f}ML","Star target"),
    ("Top 25% stations avg",  f"{top25_avg:.2f}ML",                                      "Aspiration"),
    ("Median active station", f"{median_vol:.2f}ML",                                     "Mid-tier baseline"),
    ("10-city fleet average", f"{t10_avg_per_stn:.2f}ML",                               "Peer benchmark"),
    ("Bottom 25% avg",        f"{stn.tail(max(1,n_tot//4))['vol_cy_ml'].mean():.2f}ML", "Requires activation"),
]
bmt=sl.shapes.add_table(len(bmarks)+1,3,Inches(8.3),Inches(5.12),Inches(4.85),Inches(1.85)).table
for ci,h in enumerate(["Benchmark","Value","Meaning"]):
    cell_style(bmt.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
bmt.rows[0].height=Inches(0.26)
for ri,(lbl,val,ctx) in enumerate(bmarks,1):
    bg=LGREEN if ri<=2 else (LGREY if ri==3 else (LYELL if ri==4 else LRED))
    cell_style(bmt.cell(ri,0),lbl,size=8,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
    cell_style(bmt.cell(ri,1),val,size=8.5,bold=True,fg=NAVY,bg=bg)
    cell_style(bmt.cell(ri,2),ctx,size=7.5,fg=NAVY,bg=bg,italic=True,align=PP_ALIGN.LEFT)
    bmt.rows[ri].height=Inches(0.32)
bmt.columns[0].width=Inches(2.1); bmt.columns[1].width=Inches(1.3); bmt.columns[2].width=Inches(1.45)

footer(sl,"Green rows = volume growing YoY  |  Red = declining  |  Products sold based on CY actual sales data")
notes(sl,
    f"SLIDE 3 — STAR PERFORMERS\n"
    f"Top station: {stn.iloc[0]['name']} — {stn.iloc[0]['vol_cy_ml']:.2f}ML\n"
    f"Top 5 stations combined: {stn.head(5)['vol_cy_ml'].sum():.2f}ML\n"
    f"1-product stations avg vol: {one_prod_avg:.2f}ML\n"
    f"3+ product stations avg vol: {mul_prod_avg:.2f}ML  ({mul_prod_avg/max(one_prod_avg,0.01):.1f}x)\n"
    f"If all active stations reached median {median_vol:.2f}ML, total = {median_vol*n_active:.1f}ML "
    f"(current = {tot_cy:.1f}ML)\n"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Product Coverage Gaps
# ─────────────────────────────────────────────────────────────────────────────
sl=prs.slides.add_slide(BL)
hdr(sl,f"PRODUCT GAPS: {n_fuel_no_lub} STATIONS MISSING VOLUME OPPORTUNITIES",
    "Which stations are not selling which products — and how much volume is being left on the table?")

# coverage breakdown visual
section_bar(sl,Inches(0.15),Inches(0.85),Inches(5.8),"Station Product Coverage Breakdown")
cats=[
    ("Diesel + Petrol + Lubes", n_all3,     LGREEN,  "Complete offering"),
    ("Diesel + Petrol only",    n_d_and_p,  LYELL,   "Missing lubes"),
    ("Diesel only",             n_d_only,   LRED,    "Missing petrol & lubes"),
    ("Petrol only",             n_p_only,   ORANGE,  "Missing diesel & lubes"),
    ("Inactive (zero vol)",     n_inact,    LRED,    "No sales at all"),
]
ctbl=sl.shapes.add_table(len(cats)+1,4,Inches(0.15),Inches(1.13),Inches(5.8),Inches(2.0)).table
for ci,h in enumerate(["Category","Count","% of Fleet","Status"]):
    cell_style(ctbl.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
ctbl.rows[0].height=Inches(0.27)
for ri,(lbl,cnt,bg,status) in enumerate(cats,1):
    cell_style(ctbl.cell(ri,0),lbl,size=8.5,bold=True,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
    cell_style(ctbl.cell(ri,1),str(cnt),size=9,bold=True,fg=NAVY,bg=bg)
    cell_style(ctbl.cell(ri,2),f"{cnt/n_tot*100:.0f}%",size=9,bold=True,fg=NAVY,bg=bg)
    cell_style(ctbl.cell(ri,3),status,size=8,fg=NAVY,bg=bg,italic=True,align=PP_ALIGN.LEFT)
    ctbl.rows[ri].height=Inches(0.33)
ctbl.columns[0].width=Inches(2.5); ctbl.columns[1].width=Inches(0.75)
ctbl.columns[2].width=Inches(1.0); ctbl.columns[3].width=Inches(1.55)

# specific product penetration
section_bar(sl,Inches(0.15),Inches(3.28),Inches(5.8),"Product-Level Penetration Across All Stations")
prods_pen=[
    ("HSD (Diesel)",    n_diesel,   "Core product — broadest reach"),
    ("PMG (Petrol)",    n_petrol,   "Most stations sell some petrol"),
    ("R95 (Hi-Octane)", n_r95,      "Premium — significant gap"),
    ("Any Lubricant",   n_lubes,    "Low penetration — key opportunity"),
    ("DEO",             sub["DEO"]["stns"],  "Commercial lubricant"),
    ("PCMO",            sub["PCMO"]["stns"], "Passenger car oil"),
    ("MCO",             sub["MCO"]["stns"],  "Motorcycle oil"),
    ("LOW GRADE",       sub["LOW GRADE"]["stns"] if "LOW GRADE" in sub else 0, "Should be shifted to DEO"),
]
ptbl=sl.shapes.add_table(len(prods_pen)+1,4,Inches(0.15),Inches(3.56),Inches(5.8),Inches(3.3)).table
for ci,h in enumerate(["Product","Stns Selling","% of Fleet","Note"]):
    cell_style(ptbl.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
ptbl.rows[0].height=Inches(0.25)
for ri,(lbl,cnt,note) in enumerate(prods_pen,1):
    pct=cnt/n_tot*100
    bg=LGREEN if pct>70 else (LYELL if pct>30 else LRED)
    cell_style(ptbl.cell(ri,0),lbl,size=8,bold=True,fg=NAVY,bg=bg,align=PP_ALIGN.LEFT)
    cell_style(ptbl.cell(ri,1),str(cnt),size=8.5,bold=True,fg=NAVY,bg=bg)
    cell_style(ptbl.cell(ri,2),f"{pct:.0f}%",size=8.5,bold=True,fg=NAVY,bg=bg)
    cell_style(ptbl.cell(ri,3),note,size=7.5,fg=NAVY,bg=bg,italic=True,align=PP_ALIGN.LEFT)
    ptbl.rows[ri].height=Inches(0.37)
ptbl.columns[0].width=Inches(1.8); ptbl.columns[1].width=Inches(1.0)
ptbl.columns[2].width=Inches(1.0); ptbl.columns[3].width=Inches(2.0)

# Right — gap quantification
section_bar(sl,Inches(6.1),Inches(0.85),Inches(7.05),"Volume Opportunity from Closing Product Gaps")

gaps=[
    ("Diesel stations with no Petrol",
     f"{n_d_no_p} stations",
     f"{d_no_p_vol:.1f}ML base vol",
     f"+{petrol_addon_potential:.1f}ML potential",
     f"Adding petrol at {n_d_no_p} diesel-only outlets at ~25% of their diesel vol",
     LRED),
    ("Fuel stations with no Lubricants",
     f"{n_fuel_no_lub} stations",
     f"{fnl_vol:.1f}ML base fuel vol",
     f"+{lube_potential:.1f}ML potential",
     f"Avg lube vol/stn (current lube sellers): {avg_lube_per_stn:.3f}ML",
     LRED),
    ("PMG stations with no R95",
     f"{n_pmg_only} stations",
     f"PMG-only risk",
     "+volume uplift",
     "R95 adds premium tier — does not replace PMG vol, adds new vol",
     LYELL),
    ("Inactive stations (recover LY vol)",
     f"{n_inact} stations",
     f"Had {inact_ly_vol:.1f}ML in LY",
     f"+{recovery_inact:.1f}ML at 65% recovery",
     "Field reactivation programme targeting each inactive outlet",
     LRED),
]
gh=Inches(6.15)/len(gaps)
for gi,(title,count,base,potential,detail,bg) in enumerate(gaps):
    y=Inches(1.13)+gi*gh
    rect(sl,Inches(6.1),y,Inches(7.05),gh-Inches(0.06),bg)
    rect(sl,Inches(6.1),y,Inches(0.2),gh-Inches(0.06),RED if bg==LRED else YELL)
    txt(sl,Inches(6.38),y+Inches(0.04),Inches(3.5),Inches(0.22),title,size=9,bold=True,fg=NAVY)
    txt(sl,Inches(6.38),y+Inches(0.25),Inches(2.0),Inches(0.22),count,size=8.5,fg=NAVY)
    txt(sl,Inches(6.38),y+Inches(0.43),Inches(3.5),Inches(0.2),detail,size=7.5,fg=NAVY,italic=True)
    txt(sl,Inches(9.9),y+Inches(0.1),Inches(3.1),Inches(0.5),potential,size=12,bold=True,fg=RED if bg==LRED else GREEN,align=PP_ALIGN.CENTER)
    txt(sl,Inches(9.9),y+Inches(0.55),Inches(3.1),Inches(0.22),base,size=7.5,fg=NAVY,align=PP_ALIGN.CENTER,italic=True)

# total potential
total_vol_pot=recovery_inact+lube_potential+petrol_addon_potential
rect(sl,Inches(6.1),Inches(7.0),Inches(7.05),Inches(0.15),GOLD)
txt(sl,Inches(6.2),Inches(7.01),Inches(7.0),Inches(0.13),
    f"TOTAL IDENTIFIABLE VOLUME POTENTIAL: +{total_vol_pot:.1f}ML  "
    f"(+{total_vol_pot/tot_cy*100:.0f}% of current city volume)",
    size=8.5,bold=True,fg=NAVY)

footer(sl,"Potential estimates are directional. Lube: at avg vol/lube-selling stn. "
         "Petrol: 25% of base diesel vol. Inactive: 65% LY volume recovery assumed.")
notes(sl,
    f"SLIDE 4 — PRODUCT GAPS\n"
    f"Complete (D+P+L): {n_all3}  |  D+P only: {n_d_and_p}  |  D only: {n_d_only}  |  Inactive: {n_inact}\n"
    f"Diesel penetration: {n_diesel}/{n_tot}  Petrol: {n_petrol}/{n_tot}  Lubes: {n_lubes}/{n_tot}  R95: {n_r95}/{n_tot}\n"
    f"Potential: reactivation {recovery_inact:.1f}ML + lubes {lube_potential:.1f}ML + petrol {petrol_addon_potential:.1f}ML = {total_vol_pot:.1f}ML\n"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Sub-Product Deep Dive
# ─────────────────────────────────────────────────────────────────────────────
sl=prs.slides.add_slide(BL)
hdr(sl,"SUB-PRODUCT VOLUME — WHAT IS GROWING, WHAT IS DECLINING",
    "Every product within Diesel, Petrol, and Lubricants — CY vs LY, which products are losing ground?")

active_subs=[(k,v) for k,v in sub.items() if v["cy"]>0 or v["ly"]>0]
active_subs.sort(key=lambda x: x[1]["cy"],reverse=True)

section_bar(sl,Inches(0.15),Inches(0.85),Inches(12.97),"Volume by Sub-Product — CY vs LY (only products with actual sales shown)")
cols=["Product","Group","Stations","Vol CY (ML)","Vol LY (ML)","Chg%","Share of Segment","Direction"]
n_r=len(active_subs)+1
ptbl=sl.shapes.add_table(n_r,len(cols),Inches(0.15),Inches(1.13),Inches(12.97),Inches(5.85)).table
cw=[Inches(1.8),Inches(1.1),Inches(0.8),Inches(1.4),Inches(1.4),Inches(1.0),Inches(2.0),Inches(3.47)]
for ci,(h,w) in enumerate(zip(cols,cw)):
    cell_style(ptbl.cell(0,ci),h,size=8,bold=True,fg=WHITE,bg=NAVY)
    ptbl.columns[ci].width=w
ptbl.rows[0].height=Inches(0.27)

grp_totals={"Diesel":seg_data["Diesel"]["cy"],"Petrol":seg_data["Petrol"]["cy"],"Lubricants":seg_data["Lubricants"]["cy"]}
grp_bgs={"Diesel":LBLUE,"Petrol":LGREEN,"Lubricants":RGBColor(0xFC,0xE4,0xD6)}
last_grp=None
for ri,(key,sv) in enumerate(active_subs,1):
    is_new=sv["grp"]!=last_grp; last_grp=sv["grp"]
    bg=grp_bgs.get(sv["grp"],LGREY)
    seg_sh=sv["cy"]/grp_totals.get(sv["grp"],1)*100 if grp_totals.get(sv["grp"],0)>0 else 0
    if sv["cy"]==0 and sv["ly"]>0:
        direction="⚠ DROPPED — was selling LY"; dbg=LRED
    elif sv["chg"] and sv["chg"]>5:
        direction=f"Growing (+{sv['chg']:.1f}%)"; dbg=LGREEN
    elif sv["chg"] and sv["chg"]<-5:
        direction=f"Declining ({sv['chg']:.1f}%) — investigate"; dbg=LRED
    elif sv["cy"]==0:
        direction="Not sold in Karachi"; dbg=MGREY
    else:
        direction="Stable (within ±5%)"; dbg=LYELL
    vals=[key,sv["grp"],str(sv["stns"]) if sv["stns"] else "—",
          fmt_vol(sv["cy"]),fmt_vol(sv["ly"]),fmt_pct(sv["chg"]),
          f"{seg_sh:.1f}% of {sv['grp']}",direction]
    for ci,v in enumerate(vals):
        cbg=chg_bg(sv["chg"]) if ci==5 else (dbg if ci==7 else bg)
        cell_style(ptbl.cell(ri,ci),v,size=8,bold=(is_new and ci==0),
                   fg=NAVY,bg=cbg,align=PP_ALIGN.LEFT if ci in(0,1,7) else PP_ALIGN.CENTER)
    ptbl.rows[ri].height=Inches(0.28)

footer(sl,"Products with zero CY volume but LY volume are flagged — these represent recovered volume opportunities if reactivated")
notes(sl,
    "SLIDE 5 — SUB-PRODUCT DETAIL\n"+
    "\n".join(f"  {k}: CY {v['cy']:.3f}ML  LY {v['ly']:.3f}ML  Chg {fmt_pct(v['chg'])}" for k,v in active_subs)
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Inactive & Under-Performers
# ─────────────────────────────────────────────────────────────────────────────
sl=prs.slides.add_slide(BL)
inact_df=stn[~stn["active"]].copy()
undp=stn[(stn["active"])&(stn["vol_ly"]>0)&(stn["vol_cy"]<stn["vol_ly"]*0.5)].copy()
hdr(sl,f"LOST VOLUME: {n_inact} INACTIVE + {len(undp)} SEVERELY UNDER-PERFORMING STATIONS",
    f"Inactive stations had {inact_ly_vol:.1f}ML in LY. Under-performers (vol < 50% of LY) lost "
    f"{(undp['vol_ly_ml']-undp['vol_cy_ml']).sum():.1f}ML. Combined recoverable: ~{recovery_inact+underperf_gap:.1f}ML.")

# inactive table (left)
section_bar(sl,Inches(0.15),Inches(0.85),Inches(6.4),
    f"Inactive Stations ({n_inact}) — Zero Volume in CY, Had Volume in LY")
if not inact_df.empty:
    itbl=sl.shapes.add_table(min(len(inact_df),16)+1,5,Inches(0.15),Inches(1.13),Inches(6.4),Inches(3.7)).table
    for ci,h in enumerate(["Station","LY Vol (ML)","LY Rank","Products (LY)","Recovery"]):
        cell_style(itbl.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=RED)
    itbl.columns[0].width=Inches(2.3); itbl.columns[1].width=Inches(1.1)
    itbl.columns[2].width=Inches(0.8); itbl.columns[3].width=Inches(1.4); itbl.columns[4].width=Inches(0.8)
    itbl.rows[0].height=Inches(0.25)
    inact_sorted=inact_df.sort_values("vol_ly_ml",ascending=False)
    for ri,(idx,row) in enumerate(inact_sorted.head(15).iterrows(),1):
        ly_rank=stn.index[stn["Customer Number"]==row["Customer Number"]].tolist()
        rank_str=str(ly_rank[0]+1) if ly_rank else "—"
        sp=stn_prods.get(row["Customer Number"],{})
        prods=("/".join([p for p,f in [("HSD",sp.get("has_hsd")),("PMG",sp.get("has_pmg")),
                         ("R95",sp.get("has_r95")),("L",sp.get("has_lubes"))] if f]))
        priority="HIGH" if row["vol_ly_ml"]>1 else "MED"
        cell_style(itbl.cell(ri,0),str(row["name"])[:32],size=7,fg=NAVY,bg=LRED,align=PP_ALIGN.LEFT)
        cell_style(itbl.cell(ri,1),fmt_vol(row["vol_ly_ml"]),size=7.5,bold=True,fg=RED,bg=LRED)
        cell_style(itbl.cell(ri,2),rank_str,size=7,fg=NAVY,bg=LRED)
        cell_style(itbl.cell(ri,3),prods or "Unknown",size=7,fg=NAVY,bg=LRED)
        cell_style(itbl.cell(ri,4),priority,size=7.5,bold=True,
                   fg=RED if priority=="HIGH" else ORANGE,bg=LRED)
        itbl.rows[ri].height=Inches(0.215)

# under-performers (left, below)
section_bar(sl,Inches(0.15),Inches(5.0),Inches(6.4),
    f"Severely Under-Performing ({len(undp)}) — CY Volume < 50% of Their Own LY")
if not undp.empty:
    uptbl=sl.shapes.add_table(min(len(undp),8)+1,5,Inches(0.15),Inches(5.28),Inches(6.4),Inches(1.7)).table
    for ci,h in enumerate(["Station","Vol CY","Vol LY","Loss (ML)","Drop%"]):
        cell_style(uptbl.cell(0,ci),h,size=7.5,bold=True,fg=WHITE,bg=RED)
    for ci,w in enumerate([Inches(2.3),Inches(1.0),Inches(1.0),Inches(1.0),Inches(1.1)]): uptbl.columns[ci].width=w
    uptbl.rows[0].height=Inches(0.23)
    for ri,(idx,row) in enumerate(undp.sort_values("vol_ly_ml",ascending=False).head(7).iterrows(),1):
        loss=row["vol_ly_ml"]-row["vol_cy_ml"]
        cell_style(uptbl.cell(ri,0),str(row["name"])[:32],size=7,fg=NAVY,bg=LYELL,align=PP_ALIGN.LEFT)
        cell_style(uptbl.cell(ri,1),fmt_vol(row["vol_cy_ml"]),size=7,fg=NAVY,bg=LYELL)
        cell_style(uptbl.cell(ri,2),fmt_vol(row["vol_ly_ml"]),size=7,fg=NAVY,bg=LYELL)
        cell_style(uptbl.cell(ri,3),f"-{loss:.2f}",size=7.5,bold=True,fg=RED,bg=LYELL)
        cell_style(uptbl.cell(ri,4),fmt_pct(row["chg"]),size=7.5,bold=True,fg=RED,bg=LYELL)
        uptbl.rows[ri].height=Inches(0.21)

# right panel — recovery analysis
section_bar(sl,Inches(6.7),Inches(0.85),Inches(6.45),"Recovery Potential — What These Stations Could Deliver")
insight_box(sl,Inches(6.7),Inches(1.13),Inches(6.45),Inches(2.0),
    "INACTIVE STATIONS — PRIORITY REACTIVATION",
    (f"  {n_inact} stations had a combined {inact_ly_vol:.1f}ML in LY.\n"
     f"  At 65% reactivation rate: +{recovery_inact:.1f}ML recoverable.\n\n"
     f"  TOP PRIORITY: Stations with LY vol > 1ML — these had real scale and are now dark. "
     f"Each one matters individually. Assign a field rep to each within 2 weeks.\n\n"
     f"  LOW PRIORITY: Stations with LY vol < 0.1ML — may not be worth reactivation cost. "
     f"Classify as permanently closed and remove from active fleet."))

insight_box(sl,Inches(6.7),Inches(3.25),Inches(6.45),Inches(1.85),
    "UNDER-PERFORMING STATIONS — DIAGNOSTIC REQUIRED",
    (f"  {len(undp)} stations fell to below 50% of their own LY volume.\n"
     f"  This is NOT market decline — it is station-specific failure.\n"
     f"  Each station needs a root-cause visit: equipment issue? credit problem? "
     f"competitor opened nearby? dealer disengaged?\n"
     f"  Potential recovery at 75% of LY: +{underperf_gap:.1f}ML."))

insight_box(sl,Inches(6.7),Inches(5.22),Inches(6.45),Inches(1.65),
    "COMBINED RECOVERY TARGET",
    (f"  Inactive stations:      +{recovery_inact:.1f}ML\n"
     f"  Under-performers:       +{underperf_gap:.1f}ML\n"
     f"  TOTAL RECOVERY TARGET:  +{recovery_inact+underperf_gap:.1f}ML\n\n"
     f"  This is {(recovery_inact+underperf_gap)/tot_cy*100:.0f}% of Karachi's current CY volume. "
     f"All of it from existing outlets — no new stations needed."))

footer(sl,"Under-performing = active stations with CY vol < 50% of their own LY baseline  |  Recovery = directional estimate")
notes(sl,
    f"SLIDE 6 — LOST VOLUME\n"
    f"Inactive: {n_inact} stations, LY vol {inact_ly_vol:.2f}ML, recovery {recovery_inact:.2f}ML\n"
    f"Under-performers: {len(undp)}, vol gap {underperf_gap:.2f}ML\n"
    f"Combined: {recovery_inact+underperf_gap:.2f}ML recoverable\n"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Volume Growth Roadmap
# ─────────────────────────────────────────────────────────────────────────────
sl=prs.slides.add_slide(BL)
total_pot=recovery_inact+underperf_gap+lube_potential+petrol_addon_potential
hdr(sl,f"KARACHI VOLUME GROWTH ROADMAP — +{total_pot:.0f}ML IDENTIFIABLE OPPORTUNITY",
    "Four data-driven levers. Prioritized by speed and certainty. All from existing outlets — no new stations needed.")

levers=[
    ("LEVER 1","Reactivate Inactive Stations",
     f"{n_inact} stations","Fastest win",
     f"+{recovery_inact:.1f}ML",
     [f"Assign field rep to each of the {n_inact} inactive stations within 2 weeks",
      f"Classify each: temp closure vs permanent — only reactivate viable outlets",
      f"Stations with > 1ML LY vol are highest priority (direct revenue)",
      "Root causes: equipment failure, credit freeze, management change, dealer exit",
      "Target: {:.0f}ML recovered within 60 days".format(recovery_inact*0.7)],
     RED,LRED,"30–60 days"),
    ("LEVER 2","Recover Under-Performing Stations",
     f"{len(undp)} stations","Diagnostic required",
     f"+{underperf_gap:.1f}ML",
     ["Station-by-station field diagnosis — what happened since LY?",
      "Compare each station's CY vs LY by product (which product dropped?)",
      "Check: competitor proximity, supply chain, credit limit, pricing",
      "Set 90-day volume recovery targets per station",
      "Weekly tracking vs target — escalate if no recovery in 3 weeks"],
     ORANGE,LYELL,"60–90 days"),
    ("LEVER 3","Add Lubes to Fuel-Only Stations",
     f"{n_fuel_no_lub} eligible","Structured rollout",
     f"+{lube_potential:.1f}ML",
     [f"Priority: stations within top-80% GRS that have zero lube sales",
      f"Current lube sellers average {avg_lube_per_stn*1000:.0f} litres/station in 10M",
      "Provide: lubricant display units, starter inventory, staff training",
      "Product focus: MCO first (easiest consumer pull), then DEO, then PCMO",
      "Trial: 10 stations in Month 1, full rollout if conversion > 70%"],
     DBLUE,LBLUE,"90–120 days"),
    ("LEVER 4","Add Petrol to Diesel-Only Stations",
     f"{n_d_no_p} eligible","Infrastructure dependent",
     f"+{petrol_addon_potential:.1f}ML",
     [f"{n_d_no_p} stations sell diesel but not petrol — potential PMG/R95 uplift",
      "Screen for: petrol storage capacity, location type (highway vs urban)",
      "Urban highway diesel stations are highest-value candidates",
      "Volume assumption: 25% of their current diesel vol could convert to petrol",
      "Highest ticket-size lever but requires physical investment — plan for H2"],
     GREEN,LGREEN,"120–180 days"),
]
lh=Inches(5.8)/len(levers)
for li,(lever,title,scope,tag,potential,actions,hbg,bbg,timeline) in enumerate(levers):
    y=Inches(0.85)+li*lh
    # Header bar
    rect(sl,Inches(0.15),y,Inches(12.97),Inches(0.3),hbg)
    txt(sl,Inches(0.25),y+Inches(0.03),Inches(1.2),Inches(0.25),lever,size=9,bold=True,fg=WHITE)
    txt(sl,Inches(1.5),y+Inches(0.03),Inches(6),Inches(0.25),title,size=9.5,bold=True,fg=WHITE)
    txt(sl,Inches(8),y+Inches(0.03),Inches(1.8),Inches(0.25),scope,size=8,fg=LGREY,align=PP_ALIGN.RIGHT)
    txt(sl,Inches(10),y+Inches(0.03),Inches(1.5),Inches(0.25),timeline,size=8,fg=GOLD,align=PP_ALIGN.CENTER)
    rect(sl,Inches(11.6),y,Inches(1.52),Inches(0.3),GOLD)
    txt(sl,Inches(11.62),y+Inches(0.03),Inches(1.48),Inches(0.25),potential,size=10,bold=True,fg=NAVY,align=PP_ALIGN.CENTER)
    # Body
    rect(sl,Inches(0.15),y+Inches(0.3),Inches(12.97),lh-Inches(0.34),bbg)
    for ai,action in enumerate(actions):
        ay=y+Inches(0.33)+ai*Inches((lh-Inches(0.37))/len(actions))
        txt(sl,Inches(0.35),ay,Inches(12.55),Inches((lh-Inches(0.37))/len(actions)),
            f"• {action}",size=8,fg=NAVY)

# Summary strip
rect(sl,Inches(0),Inches(6.82),SW,Inches(0.35),NAVY)
summary=(
    f"TOTAL OPPORTUNITY: +{total_pot:.0f}ML  =  "
    f"Inactive {recovery_inact:.0f}ML  +  Under-performers {underperf_gap:.0f}ML  +  "
    f"Lubes {lube_potential:.0f}ML  +  Petrol add-on {petrol_addon_potential:.0f}ML  "
    f"|  Current Karachi CY: {tot_cy:.0f}ML  |  Potential: {tot_cy+total_pot:.0f}ML  "
    f"({total_pot/tot_cy*100:.0f}% uplift)"
)
txt(sl,Inches(0.15),Inches(6.83),SW-Inches(0.3),Inches(0.3),summary,
    size=8.5,bold=True,fg=GOLD,align=PP_ALIGN.CENTER)

notes(sl,
    f"SLIDE 7 — ROADMAP\n"
    f"Lever 1 Inactive: +{recovery_inact:.1f}ML (65% of {inact_ly_vol:.1f}ML LY)\n"
    f"Lever 2 Under-perf: +{underperf_gap:.1f}ML\n"
    f"Lever 3 Lubes: +{lube_potential:.1f}ML ({n_fuel_no_lub} stns x {avg_lube_per_stn:.4f}ML)\n"
    f"Lever 4 Petrol: +{petrol_addon_potential:.1f}ML ({n_d_no_p} stns x 25% diesel vol)\n"
    f"Total: +{total_pot:.1f}ML = +{total_pot/tot_cy*100:.0f}% on {tot_cy:.1f}ML base\n"
)

# ── save ──────────────────────────────────────────────────────────────────────
out=Path("reports/PSO_Karachi_Volume_Focus.pptx")
out.parent.mkdir(exist_ok=True)
prs.save(str(out))
print(f"\nSaved → {out.resolve()}")
print(f"Slides: {len(prs.slides)}")
