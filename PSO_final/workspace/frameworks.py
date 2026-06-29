"""
PSO Lubricants — Strategic Framework Analysis
Hardcoded: BCG Matrix, Volume-Margin Matrix, Distribution Efficiency Matrix,
           Product Life Cycle, Portfolio Pareto
AI layer:  Claude / OpenAI / Gemini (whichever key is configured) synthesises
           all framework positions into strategic recommendations.
Output:    reports/PSO_Lubes_Framework_Analysis_<period>.pptx
"""
import sys, io, os
sys.path.insert(0, 'src')
sys.stdout.reconfigure(encoding='utf-8')

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pso import ingest
from _pso_common import INPUT_PATH, get_period_label, out_path

# ── palette ───────────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x00, 0x47, 0x9D)
C_DNAVY  = RGBColor(0x1F, 0x38, 0x64)
C_GREEN  = RGBColor(0x00, 0x8C, 0x4A)
C_ORANGE = RGBColor(0xE4, 0x6C, 0x0A)
C_RED    = RGBColor(0xC0, 0x00, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGREY  = RGBColor(0xF2, 0xF2, 0xF2)
C_GOLD   = RGBColor(0xFF, 0xD7, 0x00)

H_NAVY   = '#00479D'; H_DNAVY = '#1F3864'
H_GREEN  = '#008C4A'; H_ORANGE= '#E46C0A'
H_RED    = '#C00000'; H_GREY  = '#BFBFBF'
H_GOLD   = '#E8B000'; H_LGREY = '#F2F2F2'

CAT_COLORS = {'LOW GRADE': H_NAVY, 'DEO': H_GREEN, 'MCO': H_ORANGE, 'PCMO': H_RED}

# ── PPTX helpers ─────────────────────────────────────────────────────────────
def _prs():
    p = Presentation()
    p.slide_width  = Inches(10)
    p.slide_height = Inches(5.63)
    return p

def _blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(sl, x, y, w, h, rgb):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = rgb
    return s

def _txt(sl, x, y, w, h, text, size=9, bold=False, italic=False,
         color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def _hdr(sl, title, sub=''):
    _rect(sl, 0, 0, 10, 0.60, C_DNAVY)
    _txt(sl, 0.25, 0.10, 7.5, 0.34, title, size=16, bold=True)
    if sub:
        _txt(sl, 0.25, 0.40, 9.5, 0.18, sub, size=7.5,
             color=RGBColor(0xBB, 0xCC, 0xFF))

def _embed(sl, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0); plt.close(fig)
    sl.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

def _sidebar_card(sl, rx, by, cat_name, quad_text, detail, col_rgb):
    _rect(sl, rx, by, 3.30, 0.90, C_LGREY)
    _rect(sl, rx, by, 0.05, 0.90, col_rgb)
    _txt(sl, rx+0.12, by+0.06, 2.2, 0.22, cat_name,
         size=8, bold=True, color=C_DNAVY)
    _txt(sl, rx+0.12, by+0.30, 2.5, 0.20, quad_text,
         size=8, bold=True, color=col_rgb)
    _txt(sl, rx+0.12, by+0.52, 3.10, 0.32, detail,
         size=7, color=RGBColor(0x44, 0x44, 0x44), wrap=True)

# ── Load & aggregate ──────────────────────────────────────────────────────────
print("Loading data…")
df, _ = ingest.load(INPUT_PATH)
PERIOD = get_period_label(df)
retail = df[df['IsRetail'] & ~df['IsInternational']].copy()
lubes  = retail[retail['FuelSegment'] == 'Lubricants'].copy()

MAIN = ['LOW GRADE', 'DEO', 'MCO', 'PCMO']

cat = (lubes.groupby('LubeCategory')
       .agg(vol_cy   =('SalesLtr_CY',  'sum'),
            vol_sply =('SalesLtr_SPLY','sum'),
            mgn_cy   =('NetMargin_CY', 'sum'),
            stns     =('Customer Number','nunique'))
       .assign(growth  =lambda d: (d.vol_cy-d.vol_sply)/d.vol_sply.replace(0,np.nan)*100,
               share   =lambda d: d.vol_cy/d.vol_cy.sum()*100,
               mgn_pl  =lambda d: d.mgn_cy/d.vol_cy.replace(0,np.nan),
               vol_kl  =lambda d: d.vol_cy/1000)
       .loc[MAIN])

total_stns      = lubes['Customer Number'].nunique()
cat['stn_pct']  = cat['stns'] / total_stns * 100
cat['kl_per_stn'] = cat['vol_kl'] / cat['stns']

# ── ① BCG Growth-Share Matrix ─────────────────────────────────────────────────
# "Market share" = % of PSO's own lubricant portfolio volume (internal application)
BCG_SHARE_THRESH  = 20.0   # % of portfolio — natural break DEO(29%) / MCO(9%)
BCG_GROWTH_THRESH = 10.0   # % growth vs SPLY — approximately portfolio median

def _bcg(share, growth):
    hi_s = share  >= BCG_SHARE_THRESH
    hi_g = growth >= BCG_GROWTH_THRESH
    if   hi_s and hi_g:  return 'Star'
    elif hi_s:           return 'Cash Cow'
    elif hi_g:           return 'Question Mark'
    else:                return 'Dog'

cat['bcg'] = cat.apply(lambda r: _bcg(r.share, r.growth), axis=1)

BCG_RGB = {'Star': C_GOLD, 'Cash Cow': C_GREEN,
           'Question Mark': C_NAVY, 'Dog': C_RED}
BCG_HEX = {'Star': H_GOLD, 'Cash Cow': H_GREEN,
           'Question Mark': H_NAVY, 'Dog': H_RED}

# ── ② Volume-Margin Matrix ────────────────────────────────────────────────────
vol_med = cat['vol_kl'].median()
mgn_med = cat['mgn_pl'].median()

def _vm(vol, mgn):
    if   vol >= vol_med and mgn >= mgn_med: return 'Champion'
    elif vol >= vol_med:                    return 'Volume Driver'
    elif mgn >= mgn_med:                    return 'Premium Star'
    else:                                   return 'Restructure'

cat['vm'] = cat.apply(lambda r: _vm(r.vol_kl, r.mgn_pl), axis=1)

VM_RGB = {'Champion': C_GOLD, 'Volume Driver': C_NAVY,
          'Premium Star': C_GREEN, 'Restructure': C_RED}

# ── ③ Distribution Efficiency Matrix ─────────────────────────────────────────
pen_med = cat['stn_pct'].median()
prd_med = cat['kl_per_stn'].median()

def _de(pen, prd):
    if   pen >= pen_med and prd >= prd_med: return 'Optimise'
    elif pen >= pen_med:                    return 'Push Productivity'
    elif prd >= prd_med:                    return 'Expand Reach'
    else:                                   return 'Rethink'

cat['de'] = cat.apply(lambda r: _de(r.stn_pct, r.kl_per_stn), axis=1)

DE_RGB = {'Optimise': C_GREEN, 'Push Productivity': C_ORANGE,
          'Expand Reach': C_NAVY, 'Rethink': C_RED}
DE_HINT = {
    'Optimise':          'Wide reach + high output → maintain & enhance',
    'Push Productivity': 'Wide reach but low output → training & incentives',
    'Expand Reach':      'High output but few stations → add distribution',
    'Rethink':           'Low reach + low output → reposition or phase out',
}

# ── ④ Product Life Cycle ──────────────────────────────────────────────────────
def _plc(growth):
    if   growth > 20: return 'Growth'
    elif growth > 10: return 'Late Growth'
    elif growth > 3:  return 'Maturity'
    elif growth > 0:  return 'Late Maturity'
    else:             return 'Decline'

cat['plc'] = cat['growth'].apply(_plc)

PLC_HEX = {'Growth': H_GREEN, 'Late Growth': '#6EC048',
           'Maturity': H_ORANGE, 'Late Maturity': '#E4A030', 'Decline': H_RED}

# ── Print summary ─────────────────────────────────────────────────────────────
print("\n=== FRAMEWORK POSITIONS ===")
hdr = f"{'Category':<16}  {'Share%':>7}  {'Growth%':>8}  {'BCG':>14}  {'Vol-Margin':>14}  {'Dist.Eff.':>16}  {'PLC':>12}"
print(hdr); print("-" * len(hdr))
for c, r in cat.iterrows():
    print(f"{c:<16}  {r.share:>7.1f}  {r.growth:>+8.1f}  "
          f"{r.bcg:>14}  {r.vm:>14}  {r.de:>16}  {r.plc:>12}")

# ── Build charts ─────────────────────────────────────────────────────────────
print("\nBuilding charts…")
max_vol  = cat['vol_kl'].max()
max_stns = cat['stns'].max()

def _bubble(ax, x_col, y_col, size_col, size_max, size_scale=2500):
    for c_name, row in cat.iterrows():
        sz  = (row[size_col] / size_max) * size_scale + 200
        col = CAT_COLORS.get(c_name, H_GREY)
        ax.scatter(row[x_col], row[y_col], s=sz, color=col,
                   alpha=0.82, edgecolors='white', linewidth=2, zorder=5)
        ax.annotate(c_name.replace(' ', '\n'), (row[x_col], row[y_col]),
                    textcoords='offset points', xytext=(10, 4),
                    fontsize=7.5, fontweight='bold', color=col)

def _quad_bg(ax, thresh_x, thresh_y, xmin, xmax, ymin, ymax):
    # normalized y position of threshold
    yn = (thresh_y - ymin) / (ymax - ymin)
    xn = (thresh_x - xmin) / (xmax - xmin)
    ax.axvspan(thresh_x, xmax, ymin=yn, ymax=1,   alpha=0.07, color=H_GOLD)    # top-right
    ax.axvspan(thresh_x, xmax, ymin=0,  ymax=yn,  alpha=0.07, color=H_GREEN)   # bot-right
    ax.axvspan(xmin, thresh_x, ymin=yn, ymax=1,   alpha=0.07, color=H_NAVY)    # top-left
    ax.axvspan(xmin, thresh_x, ymin=0,  ymax=yn,  alpha=0.07, color=H_RED)     # bot-left
    ax.axvline(thresh_x, color='#999', linestyle='--', linewidth=1.2, zorder=1)
    ax.axhline(thresh_y, color='#999', linestyle='--', linewidth=1.2, zorder=1)

# ── Chart 1: BCG ──────────────────────────────────────────────────────────────
fig_bcg, ax = plt.subplots(figsize=(6.8, 4.8))
fig_bcg.patch.set_facecolor('white')
x0, x1 = -3, cat['share'].max() * 1.25
y0, y1 = cat['growth'].min() - 1.5, cat['growth'].max() + 2.0
_quad_bg(ax, BCG_SHARE_THRESH, BCG_GROWTH_THRESH, x0, x1, y0, y1)

for lbl, lx, ly in [('STAR', (BCG_SHARE_THRESH+x1)/2, (BCG_GROWTH_THRESH+y1)/2),
                     ('CASH COW', (BCG_SHARE_THRESH+x1)/2, (y0+BCG_GROWTH_THRESH)/2),
                     ('QUESTION MARK', (x0+BCG_SHARE_THRESH)/2, (BCG_GROWTH_THRESH+y1)/2),
                     ('DOG', (x0+BCG_SHARE_THRESH)/2, (y0+BCG_GROWTH_THRESH)/2)]:
    ax.text(lx, ly, lbl, ha='center', va='center', fontsize=8.5,
            color='#CCCCCC', fontweight='bold', alpha=0.9)

_bubble(ax, 'share', 'growth', 'vol_kl', max_vol)
ax.set_xlim(x0, x1); ax.set_ylim(y0, y1)
ax.set_xlabel('Portfolio Share — % of Total PSO Lube Volume', fontsize=9, labelpad=5)
ax.set_ylabel('Volume Growth vs SPLY (%)', fontsize=9, labelpad=5)
ax.set_title('BCG Growth-Share Matrix (Internal Portfolio)', fontsize=10,
             fontweight='bold', color=H_NAVY, pad=9)
ax.annotate(f'Share threshold {BCG_SHARE_THRESH}%', xy=(BCG_SHARE_THRESH, y0+0.3),
            fontsize=7, color='#888', ha='center')
ax.annotate(f'Growth threshold {BCG_GROWTH_THRESH}%', xy=(x0+1, BCG_GROWTH_THRESH+0.15),
            fontsize=7, color='#888')
# bubble legend
for vl, lbl in [(7000,'7,000 KL'),(2000,'2,000 KL'),(500,'500 KL')]:
    ax.scatter([], [], s=(vl/max_vol)*2500+200, color='#AAA', alpha=0.7, label=lbl)
ax.legend(title='Bubble = Vol CY', fontsize=7, title_fontsize=7, loc='lower right')
ax.set_axisbelow(True); ax.grid(True, linestyle=':', alpha=0.3)
ax.spines[['top','right']].set_visible(False)
plt.tight_layout()

# ── Chart 2: Volume-Margin ────────────────────────────────────────────────────
fig_vm, ax2 = plt.subplots(figsize=(6.8, 4.8))
fig_vm.patch.set_facecolor('white')
vx0, vx1 = 0, cat['vol_kl'].max() * 1.22
vy0, vy1 = 0, cat['mgn_pl'].max() * 1.22
_quad_bg(ax2, vol_med, mgn_med, vx0, vx1, vy0, vy1)

for lbl, lx, ly in [('CHAMPION',     (vol_med+vx1)/2, (mgn_med+vy1)/2),
                     ('VOLUME DRIVER',(vol_med+vx1)/2, (vy0+mgn_med)/2),
                     ('PREMIUM STAR', (vx0+vol_med)/2, (mgn_med+vy1)/2),
                     ('RESTRUCTURE',  (vx0+vol_med)/2, (vy0+mgn_med)/2)]:
    ax2.text(lx, ly, lbl, ha='center', va='center', fontsize=8.5,
             color='#CCCCCC', fontweight='bold', alpha=0.9)

_bubble(ax2, 'vol_kl', 'mgn_pl', 'stns', max_stns)
ax2.set_xlim(vx0, vx1); ax2.set_ylim(vy0, vy1)
ax2.set_xlabel('Volume CY (KL)', fontsize=9, labelpad=5)
ax2.set_ylabel('Net Margin / Litre (PKR)', fontsize=9, labelpad=5)
ax2.set_title('Volume-Margin Matrix', fontsize=10, fontweight='bold', color=H_NAVY, pad=9)
ax2.xaxis.set_major_formatter(FuncFormatter(lambda v,_: f'{v:,.0f}'))
ax2.yaxis.set_major_formatter(FuncFormatter(lambda v,_: f'PKR {v:,.0f}'))
for vl, lbl in [(3500,'3,500 stns'),(2000,'2,000 stns'),(800,'800 stns')]:
    ax2.scatter([], [], s=(vl/max_stns)*2500+200, color='#AAA', alpha=0.7, label=lbl)
ax2.legend(title='Bubble = Stations', fontsize=7, title_fontsize=7, loc='upper left')
ax2.set_axisbelow(True); ax2.grid(True, linestyle=':', alpha=0.3)
ax2.spines[['top','right']].set_visible(False)
plt.tight_layout()

# ── Chart 3: Distribution Efficiency ─────────────────────────────────────────
fig_de, ax3 = plt.subplots(figsize=(6.8, 4.8))
fig_de.patch.set_facecolor('white')
dx0 = max(0, cat['stn_pct'].min() - 5)
dx1 = min(100, cat['stn_pct'].max() + 5)
dy0, dy1 = 0, cat['kl_per_stn'].max() * 1.25
_quad_bg(ax3, pen_med, prd_med, dx0, dx1, dy0, dy1)

for lbl, lx, ly in [('OPTIMISE',          (pen_med+dx1)/2, (prd_med+dy1)/2),
                     ('PUSH PRODUCTIVITY', (pen_med+dx1)/2, (dy0+prd_med)/2),
                     ('EXPAND REACH',      (dx0+pen_med)/2, (prd_med+dy1)/2),
                     ('RETHINK',           (dx0+pen_med)/2, (dy0+prd_med)/2)]:
    ax3.text(lx, ly, lbl, ha='center', va='center', fontsize=8.5,
             color='#CCCCCC', fontweight='bold', alpha=0.9)

_bubble(ax3, 'stn_pct', 'kl_per_stn', 'vol_kl', max_vol)
ax3.set_xlim(dx0, dx1); ax3.set_ylim(dy0, dy1)
ax3.set_xlabel('Station Penetration (% of PSO Lube Stations Stocking Category)', fontsize=9, labelpad=5)
ax3.set_ylabel('Volume per Selling Station (KL)', fontsize=9, labelpad=5)
ax3.set_title('Distribution Efficiency Matrix', fontsize=10, fontweight='bold', color=H_NAVY, pad=9)
for vl, lbl in [(7000,'7,000 KL'),(2000,'2,000 KL'),(500,'500 KL')]:
    ax3.scatter([], [], s=(vl/max_vol)*2500+200, color='#AAA', alpha=0.7, label=lbl)
ax3.legend(title='Bubble = Vol CY', fontsize=7, title_fontsize=7, loc='upper right')
ax3.set_axisbelow(True); ax3.grid(True, linestyle=':', alpha=0.3)
ax3.spines[['top','right']].set_visible(False)
plt.tight_layout()

# ── Chart 4: PLC + Pareto ─────────────────────────────────────────────────────
fig_pp, (axL, axR) = plt.subplots(1, 2, figsize=(10, 4.5))
fig_pp.patch.set_facecolor('white')

# PLC bar chart
cat_s = cat.sort_values('share', ascending=False)
ypos  = range(len(cat_s))
bars  = axL.barh(list(ypos),
                 cat_s['growth'],
                 color=[PLC_HEX.get(cat_s.loc[c,'plc'], H_GREY) for c in cat_s.index],
                 alpha=0.85, zorder=3)
axL.axvline(BCG_GROWTH_THRESH, color='#888', linewidth=1.2, linestyle='--')
axL.set_yticks(list(ypos))
axL.set_yticklabels([c[:11] for c in cat_s.index], fontsize=9)
axL.set_xlabel('Growth % vs SPLY', fontsize=9)
axL.set_title('Product Life Cycle Stage', fontsize=10, fontweight='bold', color=H_NAVY, pad=8)
axL.set_axisbelow(True); axL.xaxis.grid(True, linestyle='--', alpha=0.4)
axL.spines[['top','right']].set_visible(False)
for bar, (c_nm, row) in zip(bars, cat_s.iterrows()):
    axL.annotate(f" {row.growth:+.1f}% — {row.plc}",
                 xy=(bar.get_width(), bar.get_y()+bar.get_height()/2),
                 va='center', fontsize=7.5, color='#333', fontweight='bold')

# Pareto
pareto = cat_s[['vol_kl','mgn_cy']].copy()
pareto['vol_cum'] = pareto['vol_kl'].cumsum() / pareto['vol_kl'].sum() * 100
pareto['mgn_cum'] = pareto['mgn_cy'].cumsum() / pareto['mgn_cy'].sum() * 100
xp = list(range(1, len(pareto)+1))
axR.bar(xp, cat_s['share'],
        color=[CAT_COLORS.get(c, H_GREY) for c in pareto.index], alpha=0.82, zorder=3)
axR2 = axR.twinx()
axR2.plot(xp, pareto['vol_cum'], 'o-', color=H_NAVY, lw=2,
          markersize=7, markerfacecolor='white', markeredgewidth=2, label='Cum. Vol %')
axR2.plot(xp, pareto['mgn_cum'], 's--', color=H_ORANGE, lw=1.5,
          markersize=6, markerfacecolor='white', markeredgewidth=2, label='Cum. Mgn %')
axR.set_xticks(xp)
axR.set_xticklabels([c[:11] for c in pareto.index], rotation=20, ha='right', fontsize=8)
axR.set_ylabel('Volume Share (%)', fontsize=9)
axR2.set_ylabel('Cumulative %', fontsize=9)
axR.set_title('Portfolio Concentration (Pareto)', fontsize=10, fontweight='bold', color=H_NAVY, pad=8)
axR.set_axisbelow(True); axR.yaxis.grid(True, linestyle='--', alpha=0.4)
axR.spines[['top','right']].set_visible(False)
lines, labels = axR2.get_legend_handles_labels()
axR2.legend(lines, labels, fontsize=7.5, loc='center right')
plt.tight_layout()

# ── Conflict detection ────────────────────────────────────────────────────────
def _detect_conflicts(row):
    """Return list of conflict strings between frameworks for one category."""
    conflicts = []
    bcg_pos = row.bcg in ('Star', 'Cash Cow')
    vm_pos  = row.vm  in ('Champion', 'Premium Star')
    de_pos  = row.de  in ('Optimise', 'Expand Reach')
    if bcg_pos and not vm_pos:
        conflicts.append(
            f"BCG={row.bcg} (high volume/share) conflicts with Vol-Margin={row.vm} "
            f"(margin PKR {row.mgn_pl:.0f}/L below portfolio median PKR {mgn_med:.0f}/L) "
            f"— volume growth not margin-accretive"
        )
    if not bcg_pos and vm_pos:
        conflicts.append(
            f"BCG={row.bcg} (share {row.share:.0f}%, growth {row.growth:+.1f}%) conflicts with "
            f"Vol-Margin={row.vm} (margin PKR {row.mgn_pl:.0f}/L above median) "
            f"— strong unit economics, insufficient scale"
        )
    if bcg_pos and not de_pos:
        conflicts.append(
            f"BCG={row.bcg} but DE={row.de} ({row.stn_pct:.0f}% penetration, "
            f"{row.kl_per_stn:.2f} KL/stn) — market position at risk from distribution weakness"
        )
    if de_pos and not bcg_pos:
        conflicts.append(
            f"DE={row.de} ({row.stn_pct:.0f}% penetration) despite BCG={row.bcg} "
            f"— distribution asset underperforming its potential"
        )
    return conflicts

# ── AI calls (two-call Option B) ──────────────────────────────────────────────
print("\nCalling AI for strategic recommendations…")
from pso.ai_insights import _build_caller
caller, provider = _build_caller()

def _prompt_per_category():
    sections = []
    for c_name, row in cat.iterrows():
        conflicts = _detect_conflicts(row)
        conflict_txt = ("\n  — " + "\n  — ".join(conflicts)) if conflicts \
                       else "  All 4 frameworks aligned — no signal conflicts."
        sections.append(
            f"### {c_name}\n"
            f"- BCG: {row.bcg}  |  Vol-Margin: {row.vm}  |  Dist. Eff.: {row.de}  |  PLC: {row.plc}\n"
            f"- Volume: {row.vol_kl:,.0f} KL | Share: {row.share:.1f}% | Growth: {row.growth:+.1f}% vs SPLY\n"
            f"- Margin: PKR {row.mgn_pl:.0f}/L | Stations: {int(row.stns):,} ({row.stn_pct:.0f}% penetration) | "
            f"Productivity: {row.kl_per_stn:.2f} KL/stn\n"
            f"- Framework conflicts:\n{conflict_txt}\n"
        )
    cat_block = "\n".join(sections)
    return f"""Period: {PERIOD}
Task: Per-category strategic signal reconciliation for PSO Lubricants portfolio.

Where framework signals conflict for a category, resolve the contradiction into a single coherent strategic directive.

Portfolio context:
- Total volume: {cat['vol_kl'].sum():,.0f} KL | BCG share threshold: ≥{BCG_SHARE_THRESH}% | Growth threshold: ≥{BCG_GROWTH_THRESH}%
- Vol-Margin medians: vol {vol_med:,.0f} KL, margin PKR {mgn_med:.0f}/L
- Distribution medians: penetration {pen_med:.0f}%, productivity {prd_med:.2f} KL/stn

{cat_block}
For EACH of the 4 categories provide exactly:
a) SIGNALS: Agree or conflict? Nature of tension if any.
b) DIRECTIVE (1 sentence): Single clear instruction — e.g. "Harvest margin while defending volume share."
c) RATIONALE (2-3 sentences): Which framework signals dominate, and why conflicting ones are overridden.
d) 2 SPECIFIC ACTIONS: Concrete 3-month actions for the sales/marketing team with expected metric impact.

Use exact numbers. No filler language."""

def _prompt_portfolio(per_cat_text):
    top2_share = cat.loc['LOW GRADE', 'share'] + cat.loc['DEO', 'share']
    top2_mgn   = (cat.loc['LOW GRADE', 'mgn_cy'] + cat.loc['DEO', 'mgn_cy']) / cat['mgn_cy'].sum() * 100
    blended_pl = cat['mgn_cy'].sum() / cat['vol_kl'].sum() / 1000
    return f"""Period: {PERIOD}
Task: Portfolio-level strategy synthesis for PSO Lubricants.

## Per-Category Analysis (derived in prior step)
{per_cat_text}

## Portfolio Metrics
- 4 categories | Total: {cat['vol_kl'].sum():,.0f} KL
- LOW GRADE + DEO = {top2_share:.0f}% of volume and {top2_mgn:.0f}% of margin pool
- All categories growing: range {cat['growth'].min():+.1f}% to {cat['growth'].max():+.1f}%
- Blended portfolio margin: PKR {blended_pl:.0f}/L (LOW GRADE volume diluting average)

Provide the following — be specific, cite exact categories and numbers:

1. PORTFOLIO HEALTH SCORE (1-10 + one sentence justification).

2. DOMINANT RISK (1 sentence): The single biggest strategic risk specific to this PSO lubricants portfolio.

3. TOP 5 PRIORITY ACTIONS (numbered, each ≤2 sentences): Ordered by impact. Name the category/cities and expected metric movement.

4. 6-MONTH OUTLOOK (3-4 sentences): If the per-category directives are executed, what moves and in which direction?

5. CONTRARIAN OBSERVATION (1-2 sentences): Something the frameworks reveal that contradicts conventional lubricant portfolio wisdom."""

AI_CAT_TEXT  = ""
AI_PORT_TEXT = ""
if caller:
    try:
        print("  Call 1: per-category reconciliation…")
        AI_CAT_TEXT  = caller(_prompt_per_category())
        print(f"  ✓ {len(AI_CAT_TEXT)} chars via {provider}")
        print("  Call 2: portfolio synthesis…")
        AI_PORT_TEXT = caller(_prompt_portfolio(AI_CAT_TEXT))
        print(f"  ✓ {len(AI_PORT_TEXT)} chars via {provider}")
    except Exception as e:
        msg = f"[AI call failed: {e}]"
        AI_CAT_TEXT  = AI_CAT_TEXT  or msg
        AI_PORT_TEXT = AI_PORT_TEXT or msg
        print(f"  !! {msg}")
else:
    AI_CAT_TEXT  = "[No AI provider configured — set ANTHROPIC_API_KEY, OPENAI_API_KEY, or GEMINI_API_KEY]"
    AI_PORT_TEXT = AI_CAT_TEXT
    print(f"  !! {AI_CAT_TEXT}")

# ── Build PPTX ────────────────────────────────────────────────────────────────
print("Building PPTX…")
prs = _prs()

# ── Slide 1: Cover ────────────────────────────────────────────────────────────
sl1 = _blank(prs)
_rect(sl1, 0, 0, 10, 5.63, C_DNAVY)
_rect(sl1, 0, 2.05, 10, 0.04, C_ORANGE)

_txt(sl1, 1.0, 0.55, 8.0, 0.40, "PSO RETAIL LUBRICANTS",
     size=12, bold=True, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
_txt(sl1, 1.0, 1.00, 8.0, 0.75, "Strategic Framework Analysis",
     size=30, bold=True, align=PP_ALIGN.CENTER)
_txt(sl1, 1.0, 1.85, 8.0, 0.22,
     "BCG Matrix  ·  Volume-Margin  ·  Distribution Efficiency  ·  Product Life Cycle  ·  Pareto",
     size=9, italic=True, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
_txt(sl1, 1.0, 2.18, 8.0, 0.24, PERIOD,
     size=10, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# 4-column summary grid
cols = [
    ('BCG',        [(c, cat.loc[c,'bcg'])  for c in MAIN]),
    ('Vol-Margin', [(c, cat.loc[c,'vm'])   for c in MAIN]),
    ('Dist. Eff.', [(c, cat.loc[c,'de'])   for c in MAIN]),
    ('PLC Stage',  [(c, cat.loc[c,'plc'])  for c in MAIN]),
]
for i, (title, items) in enumerate(cols):
    bx = 0.35 + i * 2.35
    _rect(sl1, bx, 2.55, 2.25, 2.78, RGBColor(0x10, 0x22, 0x44))
    _txt(sl1, bx+0.08, 2.62, 2.10, 0.22, title,
         size=8, bold=True, color=C_ORANGE)
    for j, (cat_name, pos) in enumerate(items):
        _txt(sl1, bx+0.08, 2.90 + j*0.60, 2.10, 0.22,
             cat_name, size=7.5, bold=True, color=RGBColor(0xBB,0xCC,0xFF))
        _txt(sl1, bx+0.08, 3.10 + j*0.60, 2.10, 0.20,
             pos, size=7, color=RGBColor(0xDD,0xDD,0xDD))

_txt(sl1, 0, 5.40, 10, 0.18,
     f"Hardcoded framework logic  |  AI: {provider or 'N/A'}  |  {PERIOD}",
     size=6.5, color=RGBColor(0x77, 0x88, 0xAA), align=PP_ALIGN.CENTER)

# ── Slide 2: BCG ──────────────────────────────────────────────────────────────
sl2 = _blank(prs)
_hdr(sl2, "BCG Growth-Share Matrix",
     f"Portfolio share threshold {BCG_SHARE_THRESH}%  |  Growth threshold {BCG_GROWTH_THRESH}%  |  Bubble = Volume CY  |  Applied to PSO internal portfolio")
_embed(sl2, fig_bcg, 0.18, 0.68, 6.10, 4.72)

rx = 6.40
_txt(sl2, rx, 0.68, 3.40, 0.20, "QUADRANT POSITIONS",
     size=7.5, bold=True, color=C_DNAVY)
for i, c_name in enumerate(MAIN):
    row  = cat.loc[c_name]
    quad = row.bcg
    rgb  = BCG_RGB.get(quad, C_LGREY)
    _sidebar_card(sl2, rx, 0.92 + i*1.15, c_name, quad,
                  f"Share {row.share:.1f}%  Growth {row.growth:+.1f}%  Vol {row.vol_kl:,.0f} KL", rgb)

# ── Slide 3: Volume-Margin ────────────────────────────────────────────────────
sl3 = _blank(prs)
_hdr(sl3, "Volume-Margin Matrix",
     f"Median split: Vol {vol_med:,.0f} KL  |  Margin PKR {mgn_med:,.0f}/L  |  Bubble = Stations")
_embed(sl3, fig_vm, 0.18, 0.68, 6.10, 4.72)

rx = 6.40
_txt(sl3, rx, 0.68, 3.40, 0.20, "QUADRANT POSITIONS",
     size=7.5, bold=True, color=C_DNAVY)
for i, c_name in enumerate(MAIN):
    row  = cat.loc[c_name]
    quad = row.vm
    rgb  = VM_RGB.get(quad, C_LGREY)
    _sidebar_card(sl3, rx, 0.92 + i*1.15, c_name, quad,
                  f"Vol {row.vol_kl:,.0f} KL  Mgn PKR {row.mgn_pl:.0f}/L  Stns {int(row.stns):,}", rgb)

# ── Slide 4: Distribution Efficiency ─────────────────────────────────────────
sl4 = _blank(prs)
_hdr(sl4, "Distribution Efficiency Matrix",
     f"Median split: Penetration {pen_med:.0f}%  |  Productivity {prd_med:.2f} KL/stn  |  Bubble = Vol CY")
_embed(sl4, fig_de, 0.18, 0.68, 6.10, 4.72)

rx = 6.40
_txt(sl4, rx, 0.68, 3.40, 0.20, "QUADRANT POSITIONS",
     size=7.5, bold=True, color=C_DNAVY)
for i, c_name in enumerate(MAIN):
    row  = cat.loc[c_name]
    quad = row.de
    rgb  = DE_RGB.get(quad, C_LGREY)
    _sidebar_card(sl4, rx, 0.92 + i*1.15, c_name, quad,
                  f"Pen {row.stn_pct:.0f}%  Prod {row.kl_per_stn:.2f} KL/stn\n{DE_HINT[quad]}", rgb)

# ── Slide 5: PLC + Pareto ─────────────────────────────────────────────────────
sl5 = _blank(prs)
_hdr(sl5, "Product Life Cycle & Portfolio Concentration",
     "Growth >20%=Growth  >10%=Late Growth  >3%=Maturity  >0%=Late Maturity  ≤0%=Decline  |  Pareto: category concentration")
_embed(sl5, fig_pp, 0.15, 0.68, 9.70, 4.75)

# ── Shared AI text renderer ───────────────────────────────────────────────────
_AI_HDRS = ('###', 'a)', 'b)', 'c)', 'd)',
            '1.', '2.', '3.', '4.', '5.',
            '##', '**', 'DIRECTIVE', 'SIGNALS', 'RATIONALE', 'ACTIONS',
            'PORTFOLIO', 'TOP 5', 'DOMINANT', 'OUTLOOK', 'SCORE', 'CONTRARIAN')

def _render_ai_text(sl, text):
    tb = sl.shapes.add_textbox(Inches(0.22), Inches(0.70), Inches(9.56), Inches(4.78))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for raw_line in text.strip().split('\n'):
        line = raw_line.strip()
        if not line:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(1)
        r = p.add_run()
        is_hdr = any(line.startswith(h) for h in _AI_HDRS)
        display = line.lstrip('#').strip().strip('*').strip()
        r.text = display
        r.font.size  = Pt(8.5) if is_hdr else Pt(7.5)
        r.font.bold  = is_hdr
        r.font.color.rgb = C_DNAVY if is_hdr else RGBColor(0x22, 0x22, 0x33)

# ── Slide 6: AI Per-Category Signal Reconciliation ───────────────────────────
sl6 = _blank(prs)
_hdr(sl6, "AI: Per-Category Signal Reconciliation",
     f"Framework conflicts resolved per category  |  Provider: {provider or 'not configured'}  |  {PERIOD}")
_render_ai_text(sl6, AI_CAT_TEXT)

# ── Slide 7: AI Portfolio Strategy ───────────────────────────────────────────
sl7 = _blank(prs)
_hdr(sl7, "AI: Portfolio Strategy & Priority Actions",
     f"Synthesised across all categories  |  Provider: {provider or 'not configured'}  |  {PERIOD}")
_render_ai_text(sl7, AI_PORT_TEXT)

# ── Save ──────────────────────────────────────────────────────────────────────
out = out_path('PSO_Lubes_Framework_Analysis', 'pptx', df)
prs.save(out)
print(f"\nSaved → {out}")
