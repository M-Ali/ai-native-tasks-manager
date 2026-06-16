"""Inspect karachi_profile.pptx — slides, shapes, text, layout."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

path = 'karachi_profile.pptx'
prs  = Presentation(path)

SW = prs.slide_width;  SH = prs.slide_height
print(f"Slide size: {SW.inches:.2f}\" x {SH.inches:.2f}\"")
print(f"Slides: {len(prs.slides)}")
print()

for si, slide in enumerate(prs.slides):
    print(f"{'='*70}")
    print(f"SLIDE {si+1}  (layout: {slide.slide_layout.name})")
    print(f"{'='*70}")
    for shape in slide.shapes:
        st = shape.shape_type
        nm = shape.name
        x  = shape.left/914400 if shape.left else 0
        y  = shape.top/914400  if shape.top  else 0
        w  = shape.width/914400 if shape.width else 0
        h  = shape.height/914400 if shape.height else 0
        print(f"  [{st}] '{nm}'  pos=({x:.2f}\",{y:.2f}\")  size=({w:.2f}\"x{h:.2f}\")")

        if hasattr(shape,'text') and shape.text.strip():
            for line in shape.text.strip().split('\n'):
                if line.strip():
                    print(f"       TEXT: {line.strip()[:120]}")

        if st == 19:  # TABLE
            tbl = shape.table
            print(f"       TABLE: {tbl.rows.__len__()} rows x {len(tbl.columns)} cols")
            for ri in range(min(len(tbl.rows),30)):
                row_texts = []
                for ci in range(len(tbl.columns)):
                    t = tbl.cell(ri,ci).text.strip()
                    row_texts.append(t[:20] if t else '')
                print(f"       row{ri}: {' | '.join(row_texts)}")

        if hasattr(shape,'chart'):
            cht = shape.chart
            print(f"       CHART type={cht.chart_type}  title={cht.has_title}")
            if cht.has_title and cht.chart_title.has_text_frame:
                print(f"       CHART TITLE: {cht.chart_title.text_frame.text}")
            for series in cht.series:
                print(f"       SERIES: '{series.name}'  vals={list(series.values)[:8]}")

        if hasattr(shape, 'image'):
            print(f"       IMAGE")
    print()
