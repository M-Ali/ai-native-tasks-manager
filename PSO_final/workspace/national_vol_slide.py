"""
PSO National Lubricants Volume Station Profile — Single Slide
Aggregates all retail lubes stations nationally (not just top cities).
Output: reports/PSO_National_Lubes_Vol_Station_Profile_<period>.pptx
"""
import sys, os
sys.path.insert(0, 'src')
sys.stdout.reconfigure(encoding='utf-8')

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pso import ingest
from _pso_common import INPUT_PATH, get_period_label, out_path

# ── constants ─────────────────────────────────────────────────────────────────
SW = Inches(10); SH = Inches(5.625)

def rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = dict(
    NAVY      ='0F2035', DNAVY='0A1628', DNAV2='1B2A4A',
    ORANGE    ='E67E22', ORANGE_L='E59866',
    BLUE      ='2E86C1', BLUE_L='7FB3D3', BLUE2='1A5276',
    GREEN     ='1E8449', GREEN_L='EAFAF1',
    RED       ='E74C3C', RED_D='922B21', RED_BG='FDEDEC',
    PURPLE    ='8E44AD',
    AMBER     ='D35400',
    GOLD      ='D4AC0D',
    GREY      ='7F8C8D', GREY2='5D6D7E', GREY3='95A5A6',
    WHITE     ='FFFFFF',
    BG_BLUE   ='EAF2FF', BG_GREEN='EAFAF1',
    BG_YELLOW ='FEF9E7', BG_RED='FDEDEC', BG_LIGHT='EBF5FB',
    KPI_BG    ='F4F6F8',
    BROWN     ='784212',
    TRACK     ='EAF2FF',
    TEAL      ='148F77', TEAL_L='A2D9CE',
)

CAT_COLORS = {
    'DEO':       C['BLUE'],
    'PCMO':      C['PURPLE'],
    'MCO':       C['ORANGE'],
    'LOW GRADE': C['GREEN'],
    'OTHERS':    C['GREY3'],
}
CAT_ORDER = ['DEO', 'PCMO', 'MCO', 'LOW GRADE', 'OTHERS']

# ── helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill_hex):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    fill = shape.fill; fill.solid()
    fill.fore_color.rgb = fill_hex if isinstance(fill_hex, RGBColor) else rgb(fill_hex)
    return shape

def txt(slide, x, y, w, h, text, size=9, bold=False, color='FFFFFF',
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = wrap; tf.auto_size = None
    para = tf.paragraphs[0]; para.alignment = align
    run  = para.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = rgb(color); run.font.name = 'Arial'
    return txb

def set_bg(slide, fill_hex):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = rgb(fill_hex)

def fmt_kl(v):
    kl = v / 1000
    if kl >= 1000: return f"{kl/1000:.3f} ML"
    if kl >= 100:  return f"{kl:.1f} KL"
    if kl >= 10:   return f"{kl:.2f} KL"
    if kl >= 1:    return f"{kl:.3f} KL"
    return f"{v:.0f} L"

def fmt_kl_total(v):
    kl = v / 1000
    if kl >= 1000: return f"{kl/1000:.2f} ML"
    return f"{kl:.1f} KL"

def chg_str(v):
    return f"{v:+.1f}%"

# ── data ──────────────────────────────────────────────────────────────────────
print("Loading data…")
df, _ = ingest.load(INPUT_PATH)
REPORT_PERIOD = get_period_label(df)
retail    = df[df['IsRetail'] & ~df['IsInternational']].copy()
lubes_all = retail[retail['FuelSegment'] == 'Lubricants'].copy()

# ── national aggregation ──────────────────────────────────────────────────────
print("Computing national station metrics…")

lubes_vol    = lubes_all['SalesLtr_CY'].sum()
lubes_vol_ly = lubes_all['SalesLtr_LY'].sum()
vol_chg      = (lubes_vol - lubes_vol_ly) / lubes_vol_ly * 100 if lubes_vol_ly else 0

stn_total = (lubes_all.groupby('Customer Number')['SalesLtr_CY']
             .sum().sort_values(ascending=False))
n_stns    = retail['Customer Number'].nunique()
n_active  = (stn_total > 0).sum()
n_zero    = n_stns - n_active

avg_per_stn = lubes_vol / n_stns if n_stns else 0
med_per_stn = stn_total[stn_total > 0].median() if n_active else 0
p75         = stn_total[stn_total > 0].quantile(0.75)

n_high   = (stn_total  > p75).sum()
n_mid    = ((stn_total > med_per_stn) & (stn_total <= p75)).sum()
n_low    = ((stn_total > 0) & (stn_total <= med_per_stn)).sum()
vol_high = stn_total[stn_total  > p75].sum()
vol_mid  = stn_total[(stn_total > med_per_stn) & (stn_total <= p75)].sum()
vol_low  = stn_total[(stn_total > 0) & (stn_total <= med_per_stn)].sum()

top10_vol = stn_total.head(10).sum()
top10_pct = top10_vol / lubes_vol * 100 if lubes_vol else 0

# top 10 with category breakdown
top10_codes  = stn_total.head(10).index.tolist()
top10_detail = []
for code in top10_codes:
    stn_data = lubes_all[lubes_all['Customer Number'] == code]
    name  = stn_data['Name 1'].iloc[0] if len(stn_data) else str(code)
    city  = stn_data['CityNorm'].iloc[0] if len(stn_data) else ''
    total = stn_total[code]
    by_cat = {cat: stn_data[stn_data['LubeCategory'] == cat]['SalesLtr_CY'].sum()
              for cat in CAT_ORDER}
    top10_detail.append({'code': code, 'name': name, 'city': city,
                         'total': total, 'by_cat': by_cat})

# category summary
cat_stats = {}
for cat in CAT_ORDER:
    df_cat     = lubes_all[lubes_all['LubeCategory'] == cat]
    stn_cat    = df_cat.groupby('Customer Number')['SalesLtr_CY'].sum()
    vol        = stn_cat.sum()
    vol_ly     = df_cat['SalesLtr_LY'].sum()
    n_sell     = (stn_cat > 0).sum()
    cat_stats[cat] = dict(
        vol     = vol,
        vol_pct = vol / lubes_vol * 100 if lubes_vol else 0,
        vol_chg = (vol - vol_ly) / vol_ly * 100 if vol_ly else 0,
        n_sell  = n_sell,
    )

d = dict(
    n_stns=n_stns, n_active=n_active, n_zero=n_zero,
    lubes_vol=lubes_vol, lubes_vol_ly=lubes_vol_ly, vol_chg=vol_chg,
    avg_per_stn=avg_per_stn, med_per_stn=med_per_stn, p75=p75,
    n_high=n_high, n_mid=n_mid, n_low=n_low, n_zero_stn=n_zero,
    vol_high=vol_high, vol_mid=vol_mid, vol_low=vol_low,
    top10_vol=top10_vol, top10_pct=top10_pct, top10=top10_detail,
    cat=cat_stats,
)

# ── build slide ───────────────────────────────────────────────────────────────
def build_national_slide(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # header
    rect(sl, 0, 0, 10, 0.58, C['DNAVY'])
    txt(sl, 0.25, 0.10, 6.5, 0.20,
        "PSO RETAIL — NATIONAL LUBRICANTS",
        size=8, color=C['TEAL_L'])
    txt(sl, 0.25, 0.27, 6.5, 0.26,
        "Station Volume Performance — Who Is Delivering",
        size=17, bold=True, color=C['WHITE'])
    txt(sl, 6.5, 0.20, 3.25, 0.18,
        f"Tiers based on Total Lubes Volume  |  Median {fmt_kl(d['med_per_stn'])}  |  {REPORT_PERIOD}",
        size=8, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    # ── LEFT: KPI strip + station tiers ───────────────────────────────────────
    # KPI strip
    kpis = [
        (fmt_kl_total(d['lubes_vol']), "Total Volume",    C['ORANGE'],  chg_str(d['vol_chg'])),
        (str(d['n_stns']),             "Total Stations",  C['BLUE'],    f"{d['n_active']} active"),
        (fmt_kl(d['avg_per_stn']),     "Avg / Station",   C['TEAL'],    f"Median {fmt_kl(d['med_per_stn'])}"),
        (f"{d['top10_pct']:.1f}%",     "Top 10 Conc.",    C['PURPLE'],  fmt_kl_total(d['top10_vol'])),
    ]
    kpi_w = 0.78
    for ki, (val, lbl, col, sub) in enumerate(kpis):
        kx = 0.19 + ki * (kpi_w + 0.04)
        rect(sl, kx, 0.65, kpi_w, 0.58, C['KPI_BG'])
        rect(sl, kx, 0.65, kpi_w, 0.03, col)
        txt(sl, kx+0.06, 0.71, kpi_w-0.08, 0.24, val,
            size=13, bold=True, color=col)
        txt(sl, kx+0.06, 0.94, kpi_w-0.08, 0.13, lbl,
            size=6.5, color=C['GREY'])
        txt(sl, kx+0.06, 1.09, kpi_w-0.08, 0.11, sub,
            size=6, color=C['GREY2'])

    txt(sl, 0.19, 1.32, 3.26, 0.14, "STATION VOLUME TIERS",
        size=7.5, bold=True, color=C['GREY'])

    tiers = [
        (C['GREEN'],  f"HIGH PERFORMERS    {d['n_high']} stations",
         fmt_kl_total(d['vol_high']),
         f"{d['vol_high']/d['lubes_vol']*100:.1f}% of National Lubes  |  Above {fmt_kl(d['p75'])} each"),
        (C['BLUE'],   f"MID PERFORMERS    {d['n_mid']} stations",
         fmt_kl_total(d['vol_mid']),
         f"{d['vol_mid']/d['lubes_vol']*100:.1f}% of National Lubes  |  {fmt_kl(d['med_per_stn'])} – {fmt_kl(d['p75'])}"),
        (C['AMBER'],  f"LOW PERFORMERS    {d['n_low']} stations",
         fmt_kl_total(d['vol_low']),
         f"{d['vol_low']/d['lubes_vol']*100:.1f}% of National Lubes  |  Below {fmt_kl(d['med_per_stn'])}"),
        (C['RED'],    f"ZERO LUBES    {d['n_zero_stn']} stations",
         "0 KL",
         "Not selling any lubricant category"),
    ]
    for ti, (col, title, val, sub) in enumerate(tiers):
        by = 1.50 + ti * 0.97
        rect(sl, 0.19, by, 3.19, 0.83, col)
        txt(sl, 0.33, by+0.08, 2.98, 0.18, title,
            size=9, bold=True, color=C['WHITE'])
        txt(sl, 0.33, by+0.28, 2.98, 0.30, val,
            size=18, bold=True, color=C['WHITE'])
        txt(sl, 0.33, by+0.62, 2.98, 0.16, sub,
            size=8, color=C['WHITE'])

    # ── RIGHT: top 10 stacked bars ────────────────────────────────────────────
    txt(sl, 3.56, 0.68, 6.35, 0.14,
        "TOP 10 STATIONS NATIONALLY BY TOTAL LUBES VOLUME (KL)",
        size=7.5, bold=True, color=C['GREY'])

    top10    = d['top10']
    chart_x  = 3.56; chart_y = 0.86
    max_vol  = top10[0]['total'] if top10 else 1
    bar_max_w = 5.20
    row_h    = 0.24; row_gap = 0.03

    for si, stn in enumerate(top10[:10]):
        ry = chart_y + si * (row_h + row_gap)
        label = f"{stn['name'][:18]}  ({stn['city']})"
        txt(sl, chart_x, ry+0.02, 1.90, row_h-0.04,
            label[:28], size=7, color=C['DNAV2'])
        bx_start = chart_x + 1.93
        for cat in CAT_ORDER:
            vol = stn['by_cat'].get(cat, 0)
            if vol > 0:
                seg_w = vol / max_vol * bar_max_w
                rect(sl, bx_start, ry+0.02, seg_w, row_h-0.04,
                     CAT_COLORS.get(cat, C['GREY3']))
                bx_start += seg_w
        txt(sl, bx_start+0.06, ry+0.02, 0.90, row_h-0.04,
            fmt_kl(stn['total']), size=7, color=C['GREY2'])

    # legend
    leg_y = chart_y + 10*(row_h+row_gap) + 0.04
    for li, cat in enumerate(CAT_ORDER):
        lx = chart_x + li * 1.25
        rect(sl, lx, leg_y, 0.12, 0.10, CAT_COLORS.get(cat, C['GREY3']))
        txt(sl, lx+0.15, leg_y, 1.08, 0.12, cat, size=6.5, color=C['GREY'])

    # insight box
    ins_y = leg_y + 0.18
    rect(sl, 3.56, ins_y, 6.25, 0.52, C['BG_YELLOW'])
    rect(sl, 3.58, ins_y, 0.02, 0.52, C['BROWN'])
    top_name = top10[0]['name'][:22] if top10 else ''
    top_city = top10[0]['city'] if top10 else ''
    ins_txt  = (
        f"Top 10 stations = {fmt_kl_total(d['top10_vol'])} ({d['top10_pct']:.1f}% of national lubes volume). "
        f"{top_name} ({top_city}) leads at {fmt_kl(top10[0]['total'])}. "
        f"{d['n_zero_stn']} stations across the network sell zero lubricants — "
        f"each conversion at national avg of {fmt_kl(d['avg_per_stn'])} adds immediate volume."
    )
    txt(sl, 3.72, ins_y+0.06, 6.06, 0.40, ins_txt,
        size=8, color=C['DNAV2'], wrap=True)


# ── run ───────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

build_national_slide(prs, d)

out = out_path('PSO_National_Lubes_Vol_Station_Profile', 'pptx', df)
prs.save(out)
print(f"Saved → {out}")
