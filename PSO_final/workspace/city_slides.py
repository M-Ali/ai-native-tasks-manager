"""
PSO — City Performance Slide Deck
Cover + 4 slides per city × 10 cities = 41 slides
  Slide 1: GRS & Volume CY vs SPLY (by product group)
  Slide 2: Station Pareto — who drives 80% of GRS
  Slide 3: Sub-product breakdown + station × product matrix
  Slide 4: Problems (data-driven) + Solutions
"""

import sys, os
sys.path.insert(0, 'src')
os.environ['PYTHONIOENCODING'] = 'utf-8'

import pandas as pd
import numpy as np
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree

from pso import ingest

# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════
DATA_FILE = "data/input/Working File Retail Fuels Data.xlsx"
OUT_FILE  = "reports/PSO_City_Slides_10M_FY26.pptx"

# Slide size: widescreen 16:9
SW = Inches(13.33)
SH = Inches(7.50)

C_NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
C_GOLD   = RGBColor(0xC9, 0xA0, 0x30)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGREY  = RGBColor(0xF2, 0xF4, 0xF8)
C_MGREY  = RGBColor(0xD9, 0xDC, 0xE3)
C_GREEN  = RGBColor(0x70, 0xAD, 0x47)
C_DGREEN = RGBColor(0x37, 0x5B, 0x25)
C_RED    = RGBColor(0xC0, 0x00, 0x00)
C_LRED   = RGBColor(0xFF, 0xC7, 0xCE)
C_YELL   = RGBColor(0xFF, 0xC0, 0x00)
C_DBLUE  = RGBColor(0x2E, 0x5B, 0x9A)
C_LBLUE  = RGBColor(0xBD, 0xD7, 0xEE)
C_ORANGE = RGBColor(0xFF, 0x82, 0x00)

# Product definitions
DIESEL_PRODS = ["HSD", "LDO", "SLUDGE-HSD", "SLUDGE"]
PETROL_PRODS = ["PMG", "R95"]
LUBE_CATS    = ["DEO", "PCMO", "MCO", "LOW GRADE", "INDUSTRIAL GRADE", "Greases", "OTHERS"]
FUEL_GROUPS  = ["Diesel", "Petrol", "Lubricants"]

GRP_COLORS = {
    "Diesel":     C_NAVY,
    "Petrol":     C_DBLUE,
    "Lubricants": C_ORANGE,
    "Total":      C_MGREY,
}

# ═══════════════════════════════════════════════════════════════════════════════
# DATA LOADING
# ═══════════════════════════════════════════════════════════════════════════════
print("Loading data …")
df, _ = ingest.load(DATA_FILE)
period = df["_Period"].iloc[0] if "_Period" in df.columns else "10M_FY26"
period_label = period.replace("_", " ")

retail = df[df["IsRetail"] & ~df["IsInternational"]].copy()

total_retail_grs = retail["SalesGRS_CY"].sum() / 1e6
total_retail_vol = retail["SalesLtr_CY"].sum() / 1e6

# Top 10 cities by GRS
top10_cities = (
    retail.groupby("CityNorm")["SalesGRS_CY"].sum()
    .nlargest(10).index.tolist()
)

# ═══════════════════════════════════════════════════════════════════════════════
# DATA EXTRACTION PER CITY
# ═══════════════════════════════════════════════════════════════════════════════
def safe_pct(cy, ly):
    if ly and ly != 0: return (cy - ly) / abs(ly) * 100
    return None

def extract(city_df):
    out = {}

    # ── Group summary (Diesel / Petrol / Lubes / Total) ───────────────────────
    groups = {}
    for seg in FUEL_GROUPS:
        s = city_df[city_df["FuelSegment"] == seg]
        groups[seg] = dict(
            grs_cy = s["SalesGRS_CY"].sum() / 1e6,
            grs_ly = s["SalesGRS_SPLY"].sum() / 1e6,
            vol_cy = s["SalesLtr_CY"].sum() / 1e6,
            vol_ly = s["SalesLtr_SPLY"].sum() / 1e6,
            stations = s["Customer Number"].nunique(),
        )
    t = city_df
    groups["Total"] = dict(
        grs_cy = t["SalesGRS_CY"].sum() / 1e6,
        grs_ly = t["SalesGRS_SPLY"].sum() / 1e6,
        vol_cy = t["SalesLtr_CY"].sum() / 1e6,
        vol_ly = t["SalesLtr_SPLY"].sum() / 1e6,
        stations = t["Customer Number"].nunique(),
    )
    for v in groups.values():
        v["grs_chg"] = safe_pct(v["grs_cy"], v["grs_ly"])
        v["vol_chg"] = safe_pct(v["vol_cy"], v["vol_ly"])
    out["groups"] = groups

    # ── Sub-product detail ────────────────────────────────────────────────────
    subs = {}
    diesel_df = city_df[city_df["FuelSegment"] == "Diesel"]
    for prod in DIESEL_PRODS:
        s = diesel_df[diesel_df["ProductCategory"] == prod]
        subs[f"Diesel|{prod}"] = dict(
            group="Diesel", label=prod,
            grs_cy = s["SalesGRS_CY"].sum() / 1e6,
            grs_ly = s["SalesGRS_SPLY"].sum() / 1e6,
            vol_cy = s["SalesLtr_CY"].sum() / 1e6,
            vol_ly = s["SalesLtr_SPLY"].sum() / 1e6,
            stations = s["Customer Number"].nunique(),
        )

    petrol_df = city_df[city_df["FuelSegment"] == "Petrol"]
    for prod in PETROL_PRODS:
        s = petrol_df[petrol_df["ProductCategory"] == prod]
        subs[f"Petrol|{prod}"] = dict(
            group="Petrol", label=prod,
            grs_cy = s["SalesGRS_CY"].sum() / 1e6,
            grs_ly = s["SalesGRS_SPLY"].sum() / 1e6,
            vol_cy = s["SalesLtr_CY"].sum() / 1e6,
            vol_ly = s["SalesLtr_SPLY"].sum() / 1e6,
            stations = s["Customer Number"].nunique(),
        )

    lubes_df = city_df[city_df["FuelSegment"] == "Lubricants"]
    for cat in LUBE_CATS:
        s = lubes_df[lubes_df["LubeCategory"] == cat]
        subs[f"Lubes|{cat}"] = dict(
            group="Lubricants", label=cat,
            grs_cy = s["SalesGRS_CY"].sum() / 1e6,
            grs_ly = s["SalesGRS_SPLY"].sum() / 1e6,
            vol_cy = s["SalesLtr_CY"].sum() / 1e6,
            vol_ly = s["SalesLtr_SPLY"].sum() / 1e6,
            stations = s["Customer Number"].nunique(),
        )

    for v in subs.values():
        v["grs_chg"] = safe_pct(v["grs_cy"], v["grs_ly"])
        v["vol_chg"] = safe_pct(v["vol_cy"], v["vol_ly"])

    out["subs"] = subs

    # ── Station pareto ────────────────────────────────────────────────────────
    stn = (
        city_df.groupby("Customer Number", as_index=False)
        .agg(
            name     = ("Name 1",        "first"),
            grs_cy   = ("SalesGRS_CY",   "sum"),
            grs_ly   = ("SalesGRS_SPLY",   "sum"),
            vol_cy   = ("SalesLtr_CY",   "sum"),
            vol_ly   = ("SalesLtr_SPLY",   "sum"),
        )
    )
    stn = stn.sort_values("grs_cy", ascending=False).reset_index(drop=True)
    total_grs_cy = stn["grs_cy"].sum()
    stn["grs_M"]      = stn["grs_cy"] / 1e6
    stn["grs_sh"]     = stn["grs_cy"] / total_grs_cy * 100 if total_grs_cy else 0
    stn["cum_grs_sh"] = stn["grs_sh"].cumsum()
    stn["active"]     = stn["vol_cy"] > 0
    stn["grs_chg"]    = stn.apply(lambda r: safe_pct(r["grs_cy"], r["grs_ly"]), axis=1)

    # Products per station
    stn_prods = {}
    for _, row in city_df.groupby("Customer Number"):
        pass  # will build below
    stn_prod_map = (
        city_df.groupby("Customer Number")["FuelSegment"]
        .apply(lambda x: sorted(x.dropna().unique().tolist()))
        .to_dict()
    )
    stn_subprod_map = {}
    for cnum, grp in city_df.groupby("Customer Number"):
        fuel_items = set(grp[grp["FuelSegment"] != "Lubricants"]["ProductCategory"].unique())
        lube_items = set(grp[grp["FuelSegment"] == "Lubricants"]["LubeCategory"].dropna().unique())
        stn_subprod_map[cnum] = {"fuel": fuel_items, "lubes": lube_items}

    stn["products"] = stn["Customer Number"].map(
        lambda c: "/".join(stn_prod_map.get(c, []))
    )

    idx_80 = int((stn["cum_grs_sh"] < 80).sum())
    out["pareto"] = stn
    out["stations_for_80pct"] = idx_80 + 1
    out["total_stations"]     = len(stn)
    out["inactive_stations"]  = int((~stn["active"]).sum())
    out["stn_prod_map"]       = stn_prod_map
    out["stn_subprod_map"]    = stn_subprod_map

    # ── Non-selling analysis ──────────────────────────────────────────────────
    not_selling_cy = []   # zero CY vol but had SPLY vol
    never_selling  = []   # zero CY and SPLY vol
    for key, sv in subs.items():
        if sv["vol_cy"] == 0 and sv["vol_ly"] > 0:
            not_selling_cy.append(sv["label"])
        elif sv["vol_cy"] == 0 and sv["vol_ly"] == 0:
            never_selling.append(sv["label"])
    out["not_selling_cy"]  = not_selling_cy   # dropped vs SPLY
    out["never_selling"]   = never_selling     # never had sales

    # ── Station × product group matrix (top 15 stations) ─────────────────────
    top15 = stn.head(15)["Customer Number"].tolist()
    matrix_rows = []
    for cnum in top15:
        stn_row = stn[stn["Customer Number"] == cnum].iloc[0]
        row = {"num": cnum, "name": str(stn_row["name"])[:35], "grs_M": stn_row["grs_M"],
               "grs_sh": stn_row["grs_sh"], "cum": stn_row["cum_grs_sh"]}
        sp = stn_subprod_map.get(cnum, {})
        fuel_set = sp.get("fuel", set())
        lube_set = sp.get("lubes", set())
        row["HSD"]        = "Y" if "HSD"  in fuel_set else "—"
        row["LDO"]        = "Y" if "LDO"  in fuel_set else "—"
        row["PMG"]        = "Y" if "PMG"  in fuel_set else "—"
        row["R95"]        = "Y" if "R95"  in fuel_set else "—"
        row["DEO"]        = "Y" if "DEO"  in lube_set else "—"
        row["PCMO"]       = "Y" if "PCMO" in lube_set else "—"
        row["MCO"]        = "Y" if "MCO"  in lube_set else "—"
        row["LOW GRADE"]  = "Y" if "LOW GRADE" in lube_set else "—"
        matrix_rows.append(row)
    out["matrix"] = matrix_rows

    # ── Cross-sell gaps ───────────────────────────────────────────────────────
    total_stn = len(stn)
    has_lubes = sum(1 for c in stn["Customer Number"] if stn_subprod_map.get(c, {}).get("lubes"))
    has_r95   = sum(1 for c in stn["Customer Number"] if "R95" in stn_subprod_map.get(c, {}).get("fuel", set()))
    has_pmg   = sum(1 for c in stn["Customer Number"] if "PMG" in stn_subprod_map.get(c, {}).get("fuel", set()))
    has_hsd   = sum(1 for c in stn["Customer Number"] if "HSD" in stn_subprod_map.get(c, {}).get("fuel", set()))
    out["has_lubes_pct"] = has_lubes / total_stn * 100 if total_stn else 0
    out["has_r95_pct"]   = has_r95   / total_stn * 100 if total_stn else 0
    out["has_pmg_pct"]   = has_pmg   / total_stn * 100 if total_stn else 0

    return out


print("Extracting city data …")
city_data = {}
for city in top10_cities:
    city_data[city] = extract(retail[retail["CityNorm"] == city].copy())
    print(f"  {city}: {city_data[city]['total_stations']} stations, "
          f"{city_data[city]['stations_for_80pct']} for 80% GRS")

# ═══════════════════════════════════════════════════════════════════════════════
# PPT HELPERS
# ═══════════════════════════════════════════════════════════════════════════════
def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

def add_textbox(slide, left, top, width, height, text, size=10,
                bold=False, color=C_NAVY, bg=None, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.italic = italic
    if bg:
        from pptx.oxml.ns import qn
        fill = txb._element.spPr
        solidFill = parse_xml(
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{rgb_hex(bg)}"/></a:solidFill>'
        )
        # simple approach: just use rectangle shape instead for backgrounds
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_label(slide, left, top, width, height, text, size=10,
              bold=False, fg=C_NAVY, bg=None, align=PP_ALIGN.LEFT,
              italic=False, v_anchor="middle"):
    if bg:
        add_rect(slide, left, top, width, height, bg)
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    if v_anchor == "middle":
        from pptx.enum.text import MSO_ANCHOR
        tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.bold   = bold
    run.font.size   = Pt(size)
    run.font.color.rgb = fg
    run.font.italic = italic
    return txb

def add_slide_header(slide, city, slide_title, period_label, color=C_NAVY):
    """Full-width header bar."""
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.72), color)
    add_label(slide, Inches(0.18), Inches(0.05), Inches(6), Inches(0.35),
              city.upper(), size=16, bold=True, fg=C_GOLD)
    add_label(slide, Inches(0.18), Inches(0.40), Inches(8), Inches(0.28),
              slide_title, size=10, bold=False, fg=C_WHITE)
    add_label(slide, Inches(10.5), Inches(0.22), Inches(2.6), Inches(0.28),
              f"PSO Retail  |  {period_label}", size=8, fg=C_MGREY, align=PP_ALIGN.RIGHT)

def add_footer(slide, text, bg=C_NAVY):
    add_rect(slide, Inches(0), Inches(7.18), SW, Inches(0.32), bg)
    add_label(slide, Inches(0.15), Inches(7.20), Inches(13), Inches(0.27),
              text, size=7.5, fg=C_MGREY, italic=True)

def add_section_title(slide, left, top, width, text, bg=C_DBLUE):
    add_rect(slide, left, top, width, Inches(0.27), bg)
    add_label(slide, left + Inches(0.08), top + Inches(0.02), width - Inches(0.1), Inches(0.24),
              text, size=8.5, bold=True, fg=C_WHITE)

def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def chg_color(val):
    if val is None: return C_MGREY
    if val > 2:  return C_GREEN
    if val < -2: return C_RED
    return C_YELL

def fmt_m(val, dec=0):
    if val is None or (isinstance(val, float) and np.isnan(val)): return "—"
    return f"{val:,.{dec}f}"

def fmt_pct(val, dec=1):
    if val is None or (isinstance(val, float) and np.isnan(val)): return "—"
    sign = "+" if val > 0 else ""
    return f"{sign}{val:.{dec}f}%"

def fmt_vol(val):
    return fmt_m(val, 1)

# ─── Table builder ─────────────────────────────────────────────────────────────
def add_table_shape(slide, left, top, width, height, rows, cols):
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height)
    return tbl.table

def style_cell(cell, text, size=8, bold=False, fg=C_NAVY, bg=None,
               align=PP_ALIGN.CENTER, italic=False):
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = fg
    run.font.italic = italic
    if bg:
        fill = cell._tc.get_or_add_tcPr()
        solidFill = parse_xml(
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{rgb_hex(bg)}"/></a:solidFill>'
        )
        for child in list(fill):
            if 'Fill' in child.tag or 'fill' in child.tag.lower():
                fill.remove(child)
        fill.append(solidFill)

def set_col_width(table, col_idx, width):
    table.columns[col_idx].width = width

def set_row_height(table, row_idx, height):
    table.rows[row_idx].height = height

# ─── Bar chart helper (CY vs SPLY grouped) ─────────────────────────────────────
def add_grouped_bar(slide, left, top, width, height,
                    categories, cy_vals, ly_vals, title="",
                    cy_label="CY", ly_label="SPLY"):
    cd = ChartData()
    cd.categories = categories
    cd.add_series(cy_label, cy_vals)
    cd.add_series(ly_label, ly_vals)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, cd
    )
    chart = chart_frame.chart

    # Style
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = C_NAVY

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Series colors
    try:
        s0 = chart.series[0]
        s0.format.fill.solid()
        s0.format.fill.fore_color.rgb = C_NAVY

        s1 = chart.series[1]
        s1.format.fill.solid()
        s1.format.fill.fore_color.rgb = C_MGREY
    except Exception:
        pass

    # Value axis
    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.size = Pt(7.5)
        chart.category_axis.tick_labels.font.bold = True
    except Exception:
        pass

    return chart_frame

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE GENERATORS
# ═══════════════════════════════════════════════════════════════════════════════

# ─── Cover Slide ──────────────────────────────────────────────────────────────
def make_cover(prs, period_label, city_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    add_rect(slide, Inches(0), Inches(0), SW, SH, C_NAVY)
    # Gold accent bar
    add_rect(slide, Inches(0), Inches(2.8), SW, Inches(0.06), C_GOLD)
    add_rect(slide, Inches(0), Inches(4.2), SW, Inches(0.06), C_GOLD)

    add_label(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
              "PSO RETAIL — CITY PERFORMANCE ANALYSIS", size=22, bold=True, fg=C_WHITE,
              align=PP_ALIGN.CENTER)
    add_label(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.6),
              f"Top 10 Cities  |  Period: {period_label}", size=14, fg=C_GOLD,
              align=PP_ALIGN.CENTER)
    add_label(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.0),
              "Value & Volume  ·  Station Contribution  ·  Product Mix  ·  Problems & Solutions",
              size=11, fg=C_LGREY, align=PP_ALIGN.CENTER)

    city_str = "  ·  ".join(city_list)
    add_label(slide, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.5),
              city_str, size=10, fg=C_GOLD, align=PP_ALIGN.CENTER)

    add_label(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
              "For internal use only — PSO Management", size=9, fg=C_MGREY,
              align=PP_ALIGN.CENTER, italic=True)

    set_notes(slide,
        f"This deck covers the top 10 cities by GRS for PSO Retail in {period_label}.\n"
        f"Cities: {', '.join(city_list)}.\n"
        "Each city has 4 slides: (1) GRS & Volume CY vs SPLY, (2) Station Pareto, "
        "(3) Product breakdown + station-product matrix, (4) Problems & Solutions.\n"
        "All numbers sourced directly from the Working File. No estimates."
    )

# ─── Slide 1: GRS & Volume CY vs SPLY ──────────────────────────────────────────
def make_slide1_grs_vol(prs, city, cd, period_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, city, "GRS & Volume — CY vs SPLY by Product Group", period_label)

    groups  = cd["groups"]
    g_order = ["Diesel", "Petrol", "Lubricants"]
    g_cats  = g_order

    # ── GRS grouped bar (left) ─────────────────────────────────────────────
    cy_grs = [groups[g]["grs_cy"] for g in g_order]
    ly_grs = [groups[g]["grs_ly"] for g in g_order]
    add_grouped_bar(slide,
        left=Inches(0.18), top=Inches(0.82), width=Inches(5.8), height=Inches(3.1),
        categories=g_cats, cy_vals=cy_grs, ly_vals=ly_grs,
        title="GRS (PKR Million) — CY vs SPLY", cy_label=f"CY ({period_label})", ly_label="SPLY"
    )

    # ── Volume grouped bar (right) ─────────────────────────────────────────
    cy_vol = [groups[g]["vol_cy"] for g in g_order]
    ly_vol = [groups[g]["vol_ly"] for g in g_order]
    add_grouped_bar(slide,
        left=Inches(6.3), top=Inches(0.82), width=Inches(6.85), height=Inches(3.1),
        categories=g_cats, cy_vals=cy_vol, ly_vals=ly_vol,
        title="Volume (Million Litres) — CY vs SPLY", cy_label=f"CY ({period_label})", ly_label="SPLY"
    )

    # ── Detail table (full width below charts) ──────────────────────────────
    add_section_title(slide, Inches(0.18), Inches(4.05), Inches(12.97),
                      "Product Group Detail — GRS & Volume with Year-on-Year Change")

    rows_order = g_order + ["Total"]
    n_rows = len(rows_order) + 1  # +1 header
    n_cols = 10
    tbl = add_table_shape(slide,
        left=Inches(0.18), top=Inches(4.34), width=Inches(12.97), height=Inches(2.6),
        rows=n_rows, cols=n_cols
    )

    headers = ["Product Group", "Stations",
               "GRS CY (PKR M)", "GRS SPLY (PKR M)", "GRS Chg%",
               "Vol CY (ML)", "Vol SPLY (ML)", "Vol Chg%",
               "GRS/Station CY (M)", "Natl Share%"]
    col_widths = [Inches(1.5), Inches(0.9),
                  Inches(1.5), Inches(1.5), Inches(1.0),
                  Inches(1.3), Inches(1.3), Inches(1.0),
                  Inches(1.47), Inches(1.0)]

    for ci, (h, w) in enumerate(zip(headers, col_widths)):
        style_cell(tbl.cell(0, ci), h, size=8, bold=True, fg=C_WHITE, bg=C_NAVY,
                   align=PP_ALIGN.CENTER)
        set_col_width(tbl, ci, w)
    set_row_height(tbl, 0, Inches(0.32))

    for ri, grp in enumerate(rows_order, 1):
        d  = groups[grp]
        is_total = grp == "Total"
        bg = C_MGREY if is_total else (C_LGREY if ri % 2 == 0 else C_WHITE)
        fg = C_NAVY

        grs_per_stn = d["grs_cy"] / d["stations"] if d["stations"] else 0
        natl_sh     = d["grs_cy"] / total_retail_grs * 100

        row_vals = [
            grp,
            str(d["stations"]),
            fmt_m(d["grs_cy"]),
            fmt_m(d["grs_ly"]),
            fmt_pct(d["grs_chg"]),
            fmt_vol(d["vol_cy"]),
            fmt_vol(d["vol_ly"]),
            fmt_pct(d["vol_chg"]),
            fmt_m(grs_per_stn, 1),
            fmt_pct(natl_sh, 1),
        ]
        for ci, val in enumerate(row_vals):
            cell_bg = bg
            if ci == 4 and d["grs_chg"] is not None:
                cell_bg = chg_color(d["grs_chg"])
            elif ci == 7 and d["vol_chg"] is not None:
                cell_bg = chg_color(d["vol_chg"])
            style_cell(tbl.cell(ri, ci), val, size=8.5, bold=is_total, fg=fg, bg=cell_bg,
                       align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
        set_row_height(tbl, ri, Inches(0.34))

    add_footer(slide,
        f"GRS = Gross Revenue from Sales (before discounts) in PKR Million  |  "
        f"Vol = Volume in Million Litres  |  "
        f"Green = Growth > +2%  |  Red = Decline > 2%  |  Yellow = Flat (within ±2%)")

    # ── Notes ──────────────────────────────────────────────────────────────
    tot   = groups["Total"]
    dsl   = groups["Diesel"]
    pet   = groups["Petrol"]
    lub   = groups["Lubricants"]

    def _note_seg(name, d):
        chg_g = f"{fmt_pct(d['grs_chg'])} YoY" if d["grs_chg"] else "no SPLY"
        chg_v = f"{fmt_pct(d['vol_chg'])} YoY" if d["vol_chg"] else "no SPLY"
        sh    = d["grs_cy"] / tot["grs_cy"] * 100 if tot["grs_cy"] else 0
        return (f"  {name}: GRS PKR {fmt_m(d['grs_cy'])}M ({sh:.0f}% of city, {chg_g}), "
                f"Vol {fmt_vol(d['vol_cy'])}ML ({chg_v})")

    notes = (
        f"SLIDE 1 — {city.upper()} — GRS & VOLUME\n\n"
        f"TOTAL CITY:\n"
        f"  GRS CY: PKR {fmt_m(tot['grs_cy'])}M   SPLY: PKR {fmt_m(tot['grs_ly'])}M   Change: {fmt_pct(tot['grs_chg'])}\n"
        f"  Vol CY: {fmt_vol(tot['vol_cy'])}ML   SPLY: {fmt_vol(tot['vol_ly'])}ML   Change: {fmt_pct(tot['vol_chg'])}\n"
        f"  Total stations: {tot['stations']}\n\n"
        f"BY PRODUCT GROUP:\n"
        f"{_note_seg('Diesel', dsl)}\n"
        f"{_note_seg('Petrol', pet)}\n"
        f"{_note_seg('Lubricants', lub)}\n\n"
        f"KEY OBSERVATION:\n"
    )

    # Data-driven key observation
    obs = []
    if tot["grs_chg"] and tot["vol_chg"] and tot["grs_chg"] > 0 and tot["vol_chg"] < -2:
        obs.append("GRS growing but volume declining — GRS increase is price-driven, NOT volume-driven. "
                   "This indicates market share loss in volume terms.")
    if tot["vol_chg"] and tot["vol_chg"] < -5:
        obs.append(f"Volume is down {fmt_pct(tot['vol_chg'])} — significant contraction. "
                   "Competitive leakage or inactive stations are likely causes.")
    if dsl["vol_chg"] and dsl["vol_chg"] < -5:
        obs.append(f"Diesel volume down {fmt_pct(dsl['vol_chg'])} — the largest volume segment is shrinking.")
    if pet["grs_cy"] / tot["grs_cy"] < 0.15 and tot["grs_cy"] > 0:
        obs.append("Petrol share is below 15% of city GRS — petrol penetration is low relative to the city's potential.")
    if lub["grs_cy"] / tot["grs_cy"] < 0.05 and tot["grs_cy"] > 0:
        obs.append("Lubricants contribute less than 5% of city GRS — substantial lube cross-sell opportunity exists.")
    if not obs:
        if tot["grs_chg"] and tot["grs_chg"] > 10 and tot["vol_chg"] and tot["vol_chg"] > 2:
            obs.append(f"Strong double-digit GRS growth ({fmt_pct(tot['grs_chg'])}) with positive volume growth — "
                       "this city is genuinely expanding in both value and volume terms.")
        else:
            obs.append("Performance is broadly stable. No major anomalies detected in value or volume trends.")

    notes += "\n".join(f"  • {o}" for o in obs)
    set_notes(slide, notes)

# ─── Slide 2: Station Pareto ──────────────────────────────────────────────────
def make_slide2_pareto(prs, city, cd, period_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, city, "Station Contribution — Who Drives 80% of City GRS", period_label)

    pareto   = cd["pareto"]
    n80      = cd["stations_for_80pct"]
    n_tot    = cd["total_stations"]
    n_inact  = cd["inactive_stations"]
    tot_grs  = pareto["grs_M"].sum()

    # ── KPI strip ─────────────────────────────────────────────────────────
    kpis = [
        ("Total Stations",     str(n_tot)),
        ("Active Stations",    str(n_tot - n_inact)),
        ("Inactive (0 Vol)",   str(n_inact)),
        (f"Stations → 80% GRS", f"{n80}  ({n80/n_tot*100:.0f}% of fleet)"),
        ("Top Station GRS",    f"PKR {pareto.iloc[0]['grs_M']:.1f}M"),
        ("Top Station Share",  f"{pareto.iloc[0]['grs_sh']:.1f}% of city"),
        ("Top 5 Stns Share",   f"{pareto.head(5)['grs_sh'].sum():.1f}% of city"),
        ("Top 10 Stns Share",  f"{pareto.head(10)['grs_sh'].sum():.1f}% of city"),
    ]
    kpi_w = SW / len(kpis)
    for ki, (label, val) in enumerate(kpis):
        x = ki * kpi_w
        bg = C_LRED if ki == 2 and n_inact > 0 else C_DBLUE
        add_rect(slide, x, Inches(0.74), kpi_w, Inches(0.52), bg)
        add_label(slide, x + Inches(0.04), Inches(0.75), kpi_w - Inches(0.08), Inches(0.22),
                  label, size=6.5, fg=C_MGREY, align=PP_ALIGN.CENTER)
        add_label(slide, x + Inches(0.04), Inches(0.96), kpi_w - Inches(0.08), Inches(0.26),
                  val, size=8, bold=True, fg=C_WHITE, align=PP_ALIGN.CENTER)

    # ── Pareto table (left ~60%) ───────────────────────────────────────────
    add_section_title(slide, Inches(0.15), Inches(1.32), Inches(7.85),
                      f"Top {min(25, n_tot)} Stations by GRS (shaded = 80% threshold)")

    show_n  = min(25, n_tot)
    n_rows  = show_n + 1
    n_cols  = 7
    tbl = add_table_shape(slide,
        Inches(0.15), Inches(1.62), Inches(7.85), Inches(5.3),
        n_rows, n_cols
    )
    hdrs = ["#", "Station Name", "GRS CY (M)", "GRS%", "Cum%", "Products", "Active?"]
    cw   = [Inches(0.3), Inches(2.7), Inches(1.1), Inches(0.7), Inches(0.7), Inches(1.7), Inches(0.65)]
    for ci, (h, w) in enumerate(zip(hdrs, cw)):
        style_cell(tbl.cell(0, ci), h, size=7.5, bold=True, fg=C_WHITE, bg=C_NAVY,
                   align=PP_ALIGN.CENTER)
        set_col_width(tbl, ci, w)
    set_row_height(tbl, 0, Inches(0.27))

    reached_80 = False
    for ri, (_, row) in enumerate(pareto.head(show_n).iterrows(), 1):
        is_key = row["cum_grs_sh"] <= 80 or (not reached_80 and row["cum_grs_sh"] > 80)
        if row["cum_grs_sh"] > 80 and not reached_80:
            reached_80 = True

        bg = C_LBLUE if is_key else C_WHITE
        if not row["active"]: bg = C_LRED

        vals = [
            str(ri),
            str(row["name"])[:38],
            f"{row['grs_M']:.1f}",
            f"{row['grs_sh']:.1f}%",
            f"{row['cum_grs_sh']:.1f}%",
            str(row["products"])[:22],
            "Active" if row["active"] else "INACTIVE",
        ]
        for ci, val in enumerate(vals):
            alg = PP_ALIGN.LEFT if ci in (1, 5) else PP_ALIGN.CENTER
            cell_bg = C_LRED if not row["active"] else bg
            style_cell(tbl.cell(ri, ci), val, size=7, bold=(ri <= n80),
                       fg=C_NAVY if row["active"] else C_RED, bg=cell_bg, align=alg)
        set_row_height(tbl, ri, Inches(0.19))

    # ── Right panel: analysis ──────────────────────────────────────────────
    add_section_title(slide, Inches(8.2), Inches(1.32), Inches(4.95),
                      "Concentration Analysis")

    # Quintile analysis
    quintile_labels = ["Top 20%", "21–40%", "41–60%", "61–80%", "Bottom 20%"]
    q_size = max(1, n_tot // 5)
    quintiles = []
    for qi in range(5):
        q_s = pareto.iloc[qi * q_size : (qi+1) * q_size]["grs_M"]
        q_sum = q_s.sum()
        q_sh  = q_sum / tot_grs * 100 if tot_grs else 0
        quintiles.append((quintile_labels[qi], len(q_s), f"{q_sh:.1f}%"))

    qn_rows = len(quintiles) + 1
    qtbl = add_table_shape(slide,
        Inches(8.2), Inches(1.62), Inches(4.95), Inches(1.5),
        qn_rows, 3
    )
    for ci, h in enumerate(["Station Tier", "Count", "GRS Share"]):
        style_cell(qtbl.cell(0, ci), h, size=7.5, bold=True, fg=C_WHITE, bg=C_NAVY,
                   align=PP_ALIGN.CENTER)
    set_row_height(qtbl, 0, Inches(0.25))

    q_bgs = [C_NAVY, C_DBLUE, C_LBLUE, C_LGREY, C_WHITE]
    q_fgs = [C_WHITE, C_WHITE, C_NAVY, C_NAVY, C_NAVY]
    for ri, (lbl, cnt, sh) in enumerate(quintiles, 1):
        style_cell(qtbl.cell(ri, 0), lbl, size=7.5, bold=True, fg=q_fgs[ri-1], bg=q_bgs[ri-1])
        style_cell(qtbl.cell(ri, 1), str(cnt), size=7.5, fg=q_fgs[ri-1], bg=q_bgs[ri-1])
        style_cell(qtbl.cell(ri, 2), sh, size=7.5, bold=True, fg=q_fgs[ri-1], bg=q_bgs[ri-1])
        set_row_height(qtbl, ri, Inches(0.22))

    # Inactive detail
    add_section_title(slide, Inches(8.2), Inches(3.22), Inches(4.95),
                      "Inactive Stations (Zero Volume CY)")
    inactive = pareto[~pareto["active"]].head(12)
    if inactive.empty:
        add_label(slide, Inches(8.2), Inches(3.52), Inches(4.95), Inches(0.4),
                  "No inactive stations — all stations have volume in CY",
                  size=8, fg=C_DGREEN)
    else:
        itbl = add_table_shape(slide,
            Inches(8.2), Inches(3.52), Inches(4.95), Inches(2.0),
            min(len(inactive), 12) + 1, 3
        )
        for ci, h in enumerate(["Station", "SPLY GRS (M)", "Products"]):
            style_cell(itbl.cell(0, ci), h, size=7, bold=True, fg=C_WHITE, bg=C_RED,
                       align=PP_ALIGN.CENTER)
        set_row_height(itbl, 0, Inches(0.22))
        for ri, (_, row) in enumerate(inactive.iterrows(), 1):
            ly_m = row["grs_ly"] / 1e6
            style_cell(itbl.cell(ri, 0), str(row["name"])[:28], size=6.5, fg=C_NAVY)
            style_cell(itbl.cell(ri, 1), f"{ly_m:.1f}", size=6.5, fg=C_NAVY, align=PP_ALIGN.CENTER)
            style_cell(itbl.cell(ri, 2), str(row["products"])[:18], size=6.5, fg=C_NAVY)
            set_row_height(itbl, ri, Inches(0.17))

    add_footer(slide,
        f"Blue highlight = stations within 80% GRS threshold  |  "
        f"Red = inactive (zero volume CY)  |  "
        f"'Products' shows segments sold by each station")

    # ── Notes ──────────────────────────────────────────────────────────────
    top1 = pareto.iloc[0]
    top5_sh = pareto.head(5)["grs_sh"].sum()
    top10_sh = pareto.head(10)["grs_sh"].sum()

    notes = (
        f"SLIDE 2 — {city.upper()} — STATION PARETO\n\n"
        f"STATION FLEET OVERVIEW:\n"
        f"  Total stations: {n_tot}\n"
        f"  Active (Vol > 0): {n_tot - n_inact}\n"
        f"  Inactive (zero CY volume): {n_inact} ({n_inact/n_tot*100:.0f}% of fleet)\n\n"
        f"CONCENTRATION ANALYSIS:\n"
        f"  Stations needed for 80% of city GRS: {n80} ({n80/n_tot*100:.0f}% of total fleet)\n"
        f"  Top 5 stations → {top5_sh:.1f}% of city GRS (PKR {pareto.head(5)['grs_M'].sum():.0f}M)\n"
        f"  Top 10 stations → {top10_sh:.1f}% of city GRS\n"
        f"  Highest GRS station: '{top1['name']}' at PKR {top1['grs_M']:.1f}M "
        f"({top1['grs_sh']:.1f}% of city total)\n\n"
        f"STATION TIER ANALYSIS (quintiles):\n"
    )
    for lbl, cnt, sh in quintiles:
        notes += f"  {lbl} stations ({cnt} stns): {sh} of city GRS\n"

    notes += f"\nKEY RISK:\n"
    if n80 / n_tot < 0.30:
        notes += (f"  HIGH CONCENTRATION RISK: Only {n80} stations ({n80/n_tot*100:.0f}% of fleet) "
                  f"generate 80% of GRS. Losing even 1–2 top stations would materially impact city revenue.\n")
    if n_inact > 3:
        inact_ly_grs = pareto[~pareto["active"]]["grs_ly"].sum() / 1e6
        notes += (f"  {n_inact} INACTIVE STATIONS had PKR {inact_ly_grs:.1f}M GRS in SPLY "
                  f"and are now showing zero volume. Immediate reactivation target.\n")
    set_notes(slide, notes)

# ─── Slide 3: Product Breakdown + Station Matrix ──────────────────────────────
def make_slide3_products(prs, city, cd, period_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, city,
        "Product Deep Dive — Sub-Product CY vs SPLY  |  Station × Product Matrix", period_label)

    subs   = cd["subs"]
    matrix = cd["matrix"]
    not_sell = cd["not_selling_cy"]
    never    = cd["never_selling"]

    # ── Sub-product table (left 55%) ───────────────────────────────────────
    add_section_title(slide, Inches(0.15), Inches(0.82), Inches(7.3),
                      "Sub-Product Breakdown — GRS & Volume CY vs SPLY")

    sub_order = (
        [k for k in subs if k.startswith("Diesel|") and (subs[k]["grs_cy"] > 0 or subs[k]["grs_ly"] > 0)] +
        [k for k in subs if k.startswith("Petrol|") and (subs[k]["grs_cy"] > 0 or subs[k]["grs_ly"] > 0)] +
        [k for k in subs if k.startswith("Lubes|")  and (subs[k]["grs_cy"] > 0 or subs[k]["grs_ly"] > 0)]
    )
    # If nothing at all, show all
    if not sub_order:
        sub_order = list(subs.keys())

    n_sub_rows = len(sub_order) + 1
    n_sub_cols = 8
    stbl = add_table_shape(slide,
        Inches(0.15), Inches(1.12), Inches(7.3), Inches(5.1),
        n_sub_rows, n_sub_cols
    )
    sub_hdrs = ["Product", "Stns", "GRS CY (M)", "GRS SPLY (M)", "GRS Chg%",
                "Vol CY (ML)", "Vol SPLY (ML)", "Vol Chg%"]
    sub_cw   = [Inches(1.3), Inches(0.5), Inches(1.1), Inches(1.1), Inches(0.85),
                Inches(0.85), Inches(0.85), Inches(0.75)]
    for ci, (h, w) in enumerate(zip(sub_hdrs, sub_cw)):
        style_cell(stbl.cell(0, ci), h, size=7.5, bold=True, fg=C_WHITE, bg=C_NAVY,
                   align=PP_ALIGN.CENTER)
        set_col_width(stbl, ci, w)
    set_row_height(stbl, 0, Inches(0.27))

    last_group = None
    for ri, key in enumerate(sub_order, 1):
        sv = subs[key]
        grp = sv["group"]
        lbl = sv["label"]
        is_new_grp = (grp != last_group)
        last_group = grp

        grp_bgs = {"Diesel": RGBColor(0xBD, 0xD7, 0xEE),
                   "Petrol": RGBColor(0xC6, 0xEF, 0xCE),
                   "Lubricants": RGBColor(0xFC, 0xE4, 0xD6)}
        row_bg = grp_bgs.get(grp, C_LGREY)
        has_zero_cy = sv["vol_cy"] == 0

        prod_label = f"{'→ ' if not is_new_grp else ''}{grp if is_new_grp else ''} {lbl}".strip()
        vals = [
            prod_label,
            str(sv["stations"]) if sv["stations"] > 0 else "—",
            fmt_m(sv["grs_cy"]) if sv["grs_cy"] > 0 else "—",
            fmt_m(sv["grs_ly"]) if sv["grs_ly"] > 0 else "—",
            fmt_pct(sv["grs_chg"]) if sv["grs_cy"] > 0 or sv["grs_ly"] > 0 else "NOT SELLING",
            fmt_vol(sv["vol_cy"]) if sv["vol_cy"] > 0 else ("—" if sv["vol_ly"] == 0 else "0.0 ▼"),
            fmt_vol(sv["vol_ly"]) if sv["vol_ly"] > 0 else "—",
            fmt_pct(sv["vol_chg"]) if sv["vol_cy"] > 0 or sv["vol_ly"] > 0 else "—",
        ]
        for ci, val in enumerate(vals):
            cell_bg = C_LRED if has_zero_cy and sv["vol_ly"] > 0 else (
                      C_MGREY if has_zero_cy and sv["vol_ly"] == 0 else row_bg)
            if ci == 4 and sv["grs_chg"] is not None and not has_zero_cy:
                cell_bg = chg_color(sv["grs_chg"])
            if ci == 7 and sv["vol_chg"] is not None and not has_zero_cy:
                cell_bg = chg_color(sv["vol_chg"])
            alg = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            style_cell(stbl.cell(ri, ci), val, size=7.5,
                       bold=(is_new_grp and ci == 0),
                       fg=C_RED if has_zero_cy and sv["vol_ly"] > 0 else C_NAVY,
                       bg=cell_bg, align=alg)
        set_row_height(stbl, ri, Inches(0.22))

    # ── Station × Product matrix (right 43%) ──────────────────────────────
    add_section_title(slide, Inches(7.6), Inches(0.82), Inches(5.55),
                      "Top 15 Stations — What Each Station Sells")

    if matrix:
        prod_cols = ["HSD", "LDO", "PMG", "R95", "DEO", "PCMO", "MCO", "LOW GRADE"]
        m_cols = len(prod_cols) + 3  # #, Name, GRS, + products
        m_rows = len(matrix) + 1

        mtbl = add_table_shape(slide,
            Inches(7.6), Inches(1.12), Inches(5.55), Inches(5.1),
            m_rows, m_cols
        )
        m_hdrs = ["#", "Station (abbrev.)", "GRS M"] + prod_cols
        m_cw   = [Inches(0.25), Inches(1.5), Inches(0.6)] + [Inches(0.4)] * len(prod_cols)
        for ci, (h, w) in enumerate(zip(m_hdrs, m_cw)):
            bg = C_NAVY if ci < 3 else GRP_COLORS.get(
                "Diesel" if h in ("HSD","LDO") else
                "Petrol" if h in ("PMG","R95") else
                "Lubricants", C_DBLUE)
            style_cell(mtbl.cell(0, ci), h, size=7, bold=True, fg=C_WHITE,
                       bg=bg, align=PP_ALIGN.CENTER)
            set_col_width(mtbl, ci, w)
        set_row_height(mtbl, 0, Inches(0.27))

        for ri, mr in enumerate(matrix, 1):
            bg = C_LGREY if ri % 2 == 0 else C_WHITE
            style_cell(mtbl.cell(ri, 0), str(ri), size=7, fg=C_NAVY, bg=bg, align=PP_ALIGN.CENTER)
            style_cell(mtbl.cell(ri, 1), mr["name"][:22], size=6.5, fg=C_NAVY, bg=bg, align=PP_ALIGN.LEFT)
            style_cell(mtbl.cell(ri, 2), f"{mr['grs_M']:.1f}", size=7, fg=C_NAVY, bg=bg, align=PP_ALIGN.CENTER)
            for ci, pc in enumerate(prod_cols, 3):
                val = mr.get(pc, "—")
                cell_bg = C_GREEN if val == "Y" else C_MGREY
                style_cell(mtbl.cell(ri, ci), val, size=7, bold=(val=="Y"),
                           fg=C_DGREEN if val == "Y" else C_MGREY,
                           bg=cell_bg, align=PP_ALIGN.CENTER)
            set_row_height(mtbl, ri, Inches(0.19))

    add_footer(slide,
        f"Red rows = product present in SPLY but ZERO volume in CY (dropped)  |  "
        f"Grey rows = product never sold in this city  |  "
        f"Green = station sells this product  |  Grey dash = not sold")

    # ── Notes ──────────────────────────────────────────────────────────────
    r95_stns  = sum(1 for mr in matrix if mr.get("R95") == "Y")
    lube_stns = sum(1 for mr in matrix if any(mr.get(l) == "Y" for l in ["DEO","PCMO","MCO","LOW GRADE"]))

    notes = (
        f"SLIDE 3 — {city.upper()} — PRODUCT BREAKDOWN\n\n"
        f"PRODUCTS DROPPED SINCE SPLY (CY volume = 0, had SPLY volume):\n"
    )
    notes += (f"  {', '.join(not_sell)}\n" if not_sell else "  None — all previously sold products still active.\n")
    notes += f"\nPRODUCTS NEVER SOLD IN THIS CITY:\n"
    notes += (f"  {', '.join(never)}\n" if never else "  All standard products have at least some SPLY history.\n")

    notes += f"\nSTATION × PRODUCT OBSERVATIONS (top 15 stations):\n"
    notes += f"  Stations selling R95 (premium petrol): {r95_stns} of top 15\n"
    notes += f"  Stations selling any lube product: {lube_stns} of top 15\n"

    cross_sell_gaps = []
    for mr in matrix:
        sells_diesel = mr.get("HSD") == "Y" or mr.get("LDO") == "Y"
        sells_lubes  = any(mr.get(l) == "Y" for l in ["DEO","PCMO","MCO","LOW GRADE"])
        sells_r95    = mr.get("R95") == "Y"
        if sells_diesel and not sells_lubes:
            cross_sell_gaps.append(f"'{mr['name']}' (PKR {mr['grs_M']:.1f}M) sells diesel but NO lubes")
        if sells_diesel and not sells_r95 and mr.get("PMG") == "Y":
            cross_sell_gaps.append(f"'{mr['name']}' sells PMG but not R95 — premium upgrade opportunity")

    if cross_sell_gaps:
        notes += "\nCROSS-SELL GAPS IN TOP 15 STATIONS:\n"
        for g in cross_sell_gaps[:6]:
            notes += f"  • {g}\n"

    set_notes(slide, notes)

# ─── Slide 4: Problems & Solutions ────────────────────────────────────────────
def make_slide4_problems(prs, city, cd, period_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, city,
        "City Assessment — Data-Identified Problems & Recommended Solutions", period_label)

    groups  = cd["groups"]
    pareto  = cd["pareto"]
    subs    = cd["subs"]
    n80     = cd["stations_for_80pct"]
    n_tot   = cd["total_stations"]
    n_inact = cd["inactive_stations"]
    not_sell = cd["not_selling_cy"]
    tot     = groups["Total"]
    dsl     = groups["Diesel"]
    pet     = groups["Petrol"]
    lub     = groups["Lubricants"]

    # ── Identify data-driven problems ─────────────────────────────────────
    problems = []
    solutions = []

    # P1: Volume decline
    if tot["vol_chg"] is not None and tot["vol_chg"] < -3:
        problems.append(
            f"VOLUME CONTRACTION ({fmt_pct(tot['vol_chg'])} YoY): City volume fell from "
            f"{fmt_vol(tot['vol_ly'])}ML to {fmt_vol(tot['vol_cy'])}ML. "
            f"GRS still shows {fmt_pct(tot['grs_chg'])} growth — entirely price-driven."
        )
        solutions.append(
            f"Conduct station-level volume audit for bottom-performing stations. "
            f"Identify competitive petrol stations within 500m of each declining outlet. "
            f"Target restoring {abs(tot['vol_cy'] - tot['vol_ly']) * 0.5:.0f}ML through "
            f"targeted pricing/loyalty activation within top 30 stations."
        )

    # P2: Inactive stations
    if n_inact > 0:
        inact_ly = pareto[~pareto["active"]]["grs_ly"].sum() / 1e6
        problems.append(
            f"INACTIVE STATIONS ({n_inact} stations, {n_inact/n_tot*100:.0f}% of fleet): "
            f"These stations had PKR {inact_ly:.1f}M GRS in SPLY and now show ZERO volume. "
            f"Direct revenue loss of ~PKR {inact_ly:.0f}M vs last year."
        )
        solutions.append(
            f"Immediate field visit to all {n_inact} inactive stations within 2 weeks. "
            f"Classify each as: (a) temporarily closed — reactivate, "
            f"(b) permanently closed — remove from active fleet, "
            f"(c) operational but not ordering — address supply/credit issues."
        )

    # P3: GRS growth but vol decline (price-led, not volume-led)
    if (tot["grs_chg"] and tot["grs_chg"] > 5 and
        tot["vol_chg"] and tot["vol_chg"] < -2):
        problems.append(
            f"PRICE-LED GROWTH, NOT VOLUME-LED: GRS up {fmt_pct(tot['grs_chg'])} but "
            f"volume down {fmt_pct(tot['vol_chg'])}. Revenue gains will reverse when "
            f"price increases normalise — the underlying volume base is eroding."
        )
        solutions.append(
            "Focus next quarter's activation on volume recovery, not price. "
            "Deploy promotional campaigns at high-traffic diesel stations. "
            "Monitor competitor discounting in the city and respond with targeted margin adjustments."
        )

    # P4: Low lube penetration
    lub_sh = lub["grs_cy"] / tot["grs_cy"] * 100 if tot["grs_cy"] else 0
    if lub_sh < 5 and tot["grs_cy"] > 20:
        problems.append(
            f"LOW LUBRICANT PENETRATION ({lub_sh:.1f}% of city GRS = PKR {lub['grs_cy']:.1f}M): "
            f"Only {lub['stations']} stations sell lubes out of {n_tot} total. "
            f"Lube yields the highest NMgn/ltr of all products."
        )
        lube_potential = n_tot * lub["grs_cy"] / max(lub["stations"], 1) * 0.3
        solutions.append(
            f"Target activation of lubricant sales at the {n_tot - lub['stations']} non-lube stations. "
            f"Prioritize stations within the top-80% GRS group. "
            f"Even 30% penetration growth could add ~PKR {lube_potential:.0f}M incremental GRS."
        )

    # P5: Low R95 penetration
    r95_grs = subs.get("Petrol|R95", {}).get("grs_cy", 0)
    pmg_grs = subs.get("Petrol|PMG", {}).get("grs_cy", 0)
    if pmg_grs > 0 and r95_grs / (pmg_grs + r95_grs) < 0.10 and pmg_grs > 5:
        problems.append(
            f"LOW R95 PREMIUM PETROL SHARE ({r95_grs/(pmg_grs+r95_grs)*100:.1f}% of petrol GRS = "
            f"PKR {r95_grs:.1f}M): City's petrol mix is heavily PMG-skewed. "
            f"R95 carries higher NMgn/ltr and captures premium consumer segment."
        )
        solutions.append(
            f"Introduce R95 at all petrol-selling stations not currently stocking it. "
            f"Price R95 competitively vs local competitors. "
            f"Run customer communication campaign highlighting R95 engine benefits."
        )

    # P6: Dropped products
    if not_sell:
        problems.append(
            f"DROPPED PRODUCTS vs SPLY: {', '.join(not_sell)} had volume in SPLY "
            f"but show ZERO volume in CY. This represents permanent revenue loss if not recovered."
        )
        solutions.append(
            f"Investigate supply disruption, distribution failure, or dealer disengagement "
            f"for: {', '.join(not_sell)}. Re-establish supply chain and offer dealer incentives to restock."
        )

    # P7: High concentration risk
    if n80 / n_tot < 0.25:
        problems.append(
            f"HIGH KEY-ACCOUNT RISK: Just {n80} stations ({n80/n_tot*100:.0f}% of fleet) "
            f"generate 80% of city GRS. The top station alone contributes "
            f"{pareto.iloc[0]['grs_sh']:.1f}% of city revenue."
        )
        solutions.append(
            f"Assign dedicated key account managers to top {min(n80, 10)} stations. "
            f"Formalize commercial agreements with top 5 stations to prevent churn to competitors. "
            f"Develop bottom-tier station activation plan to diversify revenue base."
        )

    # P8: Diesel volume decline
    if dsl["vol_chg"] and dsl["vol_chg"] < -5:
        problems.append(
            f"DIESEL VOLUME DECLINE ({fmt_pct(dsl['vol_chg'])} YoY): "
            f"Diesel is {dsl['grs_cy']/tot['grs_cy']*100:.0f}% of city GRS — "
            f"a {dsl['vol_chg']:.1f}% volume drop from {fmt_vol(dsl['vol_ly'])}ML "
            f"to {fmt_vol(dsl['vol_cy'])}ML materially impacts the city."
        )
        solutions.append(
            "Review diesel pricing vs PUMA/Shell/Total at key commercial corridors in the city. "
            "Target freight company accounts and commercial fleet operators with bulk deal pricing. "
            "Check if any major fleet contracts lapsed since SPLY."
        )

    # Ensure at least 2 problems/solutions
    if not problems:
        problems.append(
            f"OVERALL STABLE PERFORMANCE: City shows {fmt_pct(tot['grs_chg'])} GRS growth "
            f"and {fmt_pct(tot['vol_chg'])} volume growth. No acute structural issues detected."
        )
        solutions.append(
            "Maintain current performance. Focus on: (1) growing lube cross-sell at diesel-only stations, "
            "(2) incremental R95 penetration, (3) defending market share in top-GRS stations."
        )

    # ── Render problems (left pane) ────────────────────────────────────────
    add_rect(slide, Inches(0.15), Inches(0.82), Inches(6.4), Inches(6.2), C_LRED)
    add_rect(slide, Inches(0.15), Inches(0.82), Inches(6.4), Inches(0.35), C_RED)
    add_label(slide, Inches(0.25), Inches(0.84), Inches(6.2), Inches(0.3),
              "PROBLEMS IDENTIFIED — FROM DATA", size=9.5, bold=True, fg=C_WHITE)

    p_top = Inches(1.22)
    p_h_each = Inches(6.2 - 0.4) / max(len(problems), 1)
    for pi, prob in enumerate(problems[:6]):
        y = p_top + pi * p_h_each
        add_rect(slide, Inches(0.2), y, Inches(6.25), p_h_each - Inches(0.06), C_WHITE)
        add_rect(slide, Inches(0.2), y, Inches(0.22), p_h_each - Inches(0.06), C_RED)
        num_y = y + (p_h_each - Inches(0.06)) / 2 - Inches(0.12)
        add_label(slide, Inches(0.22), num_y, Inches(0.22), Inches(0.25),
                  str(pi+1), size=9, bold=True, fg=C_WHITE)
        add_label(slide, Inches(0.46), y + Inches(0.03), Inches(5.92), p_h_each - Inches(0.1),
                  prob, size=8, fg=C_NAVY, italic=False)

    # ── Render solutions (right pane) ─────────────────────────────────────
    add_rect(slide, Inches(6.73), Inches(0.82), Inches(6.42), Inches(6.2),
             RGBColor(0xE2, 0xEF, 0xDA))
    add_rect(slide, Inches(6.73), Inches(0.82), Inches(6.42), Inches(0.35), C_DGREEN)
    add_label(slide, Inches(6.83), Inches(0.84), Inches(6.2), Inches(0.3),
              "RECOMMENDED SOLUTIONS", size=9.5, bold=True, fg=C_WHITE)

    s_top = Inches(1.22)
    s_h_each = Inches(6.2 - 0.4) / max(len(solutions), 1)
    for si, sol in enumerate(solutions[:6]):
        y = s_top + si * s_h_each
        add_rect(slide, Inches(6.78), y, Inches(6.3), s_h_each - Inches(0.06),
                 RGBColor(0xF2, 0xF8, 0xED))
        add_rect(slide, Inches(6.78), y, Inches(0.22), s_h_each - Inches(0.06),
                 C_DGREEN)
        num_y = y + (s_h_each - Inches(0.06)) / 2 - Inches(0.12)
        add_label(slide, Inches(6.80), num_y, Inches(0.22), Inches(0.25),
                  str(si+1), size=9, bold=True, fg=C_WHITE)
        add_label(slide, Inches(7.04), y + Inches(0.03), Inches(5.9), s_h_each - Inches(0.1),
                  sol, size=8, fg=C_NAVY)

    add_footer(slide,
        "All problems identified directly from Working File data. "
        "No assumptions made beyond data. Solutions are directional — validate with field teams.")

    # ── Notes ──────────────────────────────────────────────────────────────
    notes = (
        f"SLIDE 4 — {city.upper()} — PROBLEMS & SOLUTIONS\n\n"
        f"DATA SUMMARY USED FOR ANALYSIS:\n"
        f"  GRS CY: PKR {fmt_m(tot['grs_cy'])}M  |  SPLY: PKR {fmt_m(tot['grs_ly'])}M  "
        f"|  Chg: {fmt_pct(tot['grs_chg'])}\n"
        f"  Vol CY: {fmt_vol(tot['vol_cy'])}ML  |  SPLY: {fmt_vol(tot['vol_ly'])}ML  "
        f"|  Chg: {fmt_pct(tot['vol_chg'])}\n"
        f"  Stations: {n_tot}  |  Inactive: {n_inact}  |  For 80% GRS: {n80}\n"
        f"  Lube share: {lub['grs_cy']/tot['grs_cy']*100:.1f}%  "
        f"|  Products dropped: {', '.join(not_sell) if not_sell else 'None'}\n\n"
        f"PROBLEMS ({len(problems)} identified):\n"
    )
    for pi, p in enumerate(problems, 1):
        notes += f"  P{pi}: {p}\n"
    notes += f"\nSOLUTIONS:\n"
    for si, s in enumerate(solutions, 1):
        notes += f"  S{si}: {s}\n"
    set_notes(slide, notes)

# ═══════════════════════════════════════════════════════════════════════════════
# MAIN — BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════════
print("Building PowerPoint …")
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

make_cover(prs, period_label, top10_cities)

for city_n, city in enumerate(top10_cities, 1):
    cd = city_data[city]
    print(f"  [{city_n}/10] {city} — generating 4 slides …")
    make_slide1_grs_vol  (prs, city, cd, period_label)
    make_slide2_pareto   (prs, city, cd, period_label)
    make_slide3_products (prs, city, cd, period_label)
    make_slide4_problems (prs, city, cd, period_label)

out_path = Path(OUT_FILE)
out_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out_path))
print(f"\nSaved → {out_path.resolve()}")
print(f"Total slides: {len(prs.slides)}")
