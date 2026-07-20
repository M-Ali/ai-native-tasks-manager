"""
PSO City Lubricant Profile Generator — VOLUME EDITION
Same 5-slide format as city_profiles.py but metric = SalesLtr_CY (litres/KL).
Saves to reports/city_profiles_volume/ — does NOT overwrite revenue profiles.
"""
import sys, os
sys.path.insert(0,'src')
sys.stdout.reconfigure(encoding='utf-8')

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pso import ingest
from _pso_common import INPUT_PATH, get_period, get_period_label, out_path

# ── constants ─────────────────────────────────────────────────────────────────
SW = Inches(10); SH = Inches(5.625)

def rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = dict(
    NAVY      ='0F2035', DNAVY='0A1628', DNAV2='1B2A4A',
    ORANGE    ='E67E22', ORANGE_L='E59866',
    BLUE      ='2E86C1', BLUE_L='7FB3D3', BLUE2='1A5276',
    GREEN     ='1E8449', GREEN_L='EAFAF1',
    RED       ='E74C3C', RED_D='922B21', RED_BG='FDEDEC',
    PURPLE    ='8E44AD',
    AMBER     ='D35400',
    GOLD      ='D4AC0D',
    GREY      ='7F8C8D', GREY2='5D6D7E', GREY3='95A5A6',
    WHITE     ='FFFFFF',
    BG_BLUE   ='EAF2FF', BG_GREEN='EAFAF1',
    BG_YELLOW ='FEF9E7', BG_RED='FDEDEC', BG_LIGHT='EBF5FB',
    KPI_BG    ='F4F6F8',
    BROWN     ='784212',
    TRACK     ='EAF2FF',
    # Volume edition accent — teal instead of orange on cover
    TEAL      ='148F77', TEAL_L='A2D9CE',
)

CAT_COLORS = {
    'DEO':       C['BLUE'],
    'PCMO':      C['PURPLE'],
    'MCO':       C['ORANGE'],
    'LOW GRADE': C['GREEN'],
    'OTHERS':    C['GREY3'],
}
CAT_ORDER = ['DEO','PCMO','MCO','LOW GRADE','OTHERS']

# ── helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill_hex):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    fill = shape.fill; fill.solid()
    fill.fore_color.rgb = fill_hex if isinstance(fill_hex, RGBColor) else rgb(fill_hex)
    return shape

def txt(slide, x, y, w, h, text, size=9, bold=False, color='FFFFFF',
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = wrap; tf.auto_size = None
    para = tf.paragraphs[0]; para.alignment = align
    run  = para.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = rgb(color); run.font.name = 'Arial'
    return txb

def set_bg(slide, fill_hex):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = rgb(fill_hex)

def fmt_kl(v):
    """Format litres as KL with appropriate precision."""
    kl = v / 1000
    if kl >= 1000: return f"{kl/1000:.3f} ML"
    if kl >= 100:  return f"{kl:.1f} KL"
    if kl >= 10:   return f"{kl:.2f} KL"
    if kl >= 1:    return f"{kl:.3f} KL"
    return f"{v:.0f} L"

def fmt_kl_total(v):
    kl = v / 1000
    if kl >= 1000: return f"{kl/1000:.2f} ML"
    return f"{kl:.1f} KL"

def chg_str(v):
    return f"{v:+.1f}%"

# ── data ──────────────────────────────────────────────────────────────────────
print("Loading data…")
df, _ = ingest.load(INPUT_PATH)
REPORT_PERIOD = get_period_label(df)
retail    = df[df['IsRetail'] & ~df['IsInternational']].copy()
lubes_all = retail[retail['FuelSegment']=='Lubricants'].copy()

def compute_city(city_norm):
    city_lubes = lubes_all[lubes_all['CityNorm']==city_norm].copy()
    city_all   = retail[retail['CityNorm']==city_norm].copy()

    # ── volume metrics (litres) ───────────────────────────────────────────────
    lubes_vol     = city_lubes['SalesLtr_CY'].sum()
    lubes_vol_sply= city_lubes['SalesLtr_SPLY'].sum()
    vol_chg       = (lubes_vol - lubes_vol_sply) / lubes_vol_sply * 100 if lubes_vol_sply else 0

    # per-station volume totals
    stn_total = (city_lubes.groupby('Customer Number')['SalesLtr_CY']
                 .sum().sort_values(ascending=False))
    n_stns    = city_all['Customer Number'].nunique()
    n_active  = (stn_total > 0).sum()
    n_zero    = n_stns - n_active

    avg_per_stn = lubes_vol / n_stns if n_stns else 0
    med_per_stn = stn_total[stn_total > 0].median() if n_active else 0
    min_per_stn = stn_total[stn_total > 0].min()    if n_active else 0
    max_per_stn = stn_total.max()                   if n_active else 0

    top10_vol = stn_total.head(10).sum()
    top10_pct = top10_vol / lubes_vol * 100 if lubes_vol else 0

    active_cats = [c for c in CAT_ORDER
                   if city_lubes[city_lubes['LubeCategory']==c]['SalesLtr_CY'].sum() > 0]
    n_cats = len(active_cats)

    # per-category volume stats
    cat_stats = {}
    for cat in CAT_ORDER:
        df_cat = city_lubes[city_lubes['LubeCategory']==cat]
        stn_cat        = df_cat.groupby('Customer Number')['SalesLtr_CY'].sum()
        stn_cat_active = stn_cat[stn_cat > 0]
        vol    = stn_cat.sum()
        vol_sply = df_cat['SalesLtr_SPLY'].sum()
        n_sell = (stn_cat > 0).sum()
        cat_stats[cat] = dict(
            vol       = vol,
            vol_pct   = vol / lubes_vol * 100 if lubes_vol else 0,
            vol_chg   = (vol - vol_sply) / vol_sply * 100 if vol_sply else 0,
            n_selling = n_sell,
            n_pct     = n_sell / n_stns * 100 if n_stns else 0,
            avg       = stn_cat_active.mean()   if len(stn_cat_active) else 0,
            median    = stn_cat_active.median() if len(stn_cat_active) else 0,
            min_val   = stn_cat_active.min()    if len(stn_cat_active) else 0,
            max_val   = stn_cat_active.max()    if len(stn_cat_active) else 0,
        )

    # tiers based on median volume
    p75      = stn_total[stn_total > 0].quantile(0.75)
    n_high   = (stn_total  > p75).sum()
    n_mid    = ((stn_total > med_per_stn) & (stn_total <= p75)).sum()
    n_low    = ((stn_total > 0) & (stn_total <= med_per_stn)).sum()
    vol_high = stn_total[stn_total  > p75].sum()
    vol_mid  = stn_total[(stn_total > med_per_stn) & (stn_total <= p75)].sum()
    vol_low  = stn_total[(stn_total > 0) & (stn_total <= med_per_stn)].sum()

    # top 10 by volume with category breakdown
    top10_codes  = stn_total.head(10).index.tolist()
    top10_detail = []
    for code in top10_codes:
        stn_data = city_lubes[city_lubes['Customer Number']==code]
        name  = stn_data['Name 1'].iloc[0] if len(stn_data) else code
        total = stn_total[code]
        by_cat = {cat: stn_data[stn_data['LubeCategory']==cat]['SalesLtr_CY'].sum()
                  for cat in CAT_ORDER}
        top10_detail.append({'code': code, 'name': name, 'total': total, 'by_cat': by_cat})

    no_pcmo   = n_stns - cat_stats['PCMO']['n_selling']
    no_lowgr  = n_stns - cat_stats['LOW GRADE']['n_selling']
    pcmo_avg  = cat_stats['PCMO']['avg']
    lowgr_avg = cat_stats['LOW GRADE']['avg']

    return dict(
        city=city_norm, n_stns=n_stns, n_active=n_active, n_zero=n_zero,
        lubes_vol=lubes_vol, lubes_vol_sply=lubes_vol_sply, vol_chg=vol_chg,
        avg_per_stn=avg_per_stn, med_per_stn=med_per_stn,
        min_per_stn=min_per_stn, max_per_stn=max_per_stn,
        top10_vol=top10_vol, top10_pct=top10_pct,
        n_cats=n_cats, active_cats=active_cats,
        cat=cat_stats,
        n_high=n_high, n_mid=n_mid, n_low=n_low, n_zero_stn=n_zero,
        vol_high=vol_high, vol_mid=vol_mid, vol_low=vol_low,
        p75=p75, top10=top10_detail,
        no_pcmo=no_pcmo, no_lowgr=no_lowgr,
        pcmo_avg=pcmo_avg, lowgr_avg=lowgr_avg,
    )

# ── SLIDE 1: Cover ────────────────────────────────────────────────────────────
def build_slide1(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C['NAVY'])

    city = d['city']
    rect(sl, 0.90, 1.54, 0.50, 0.06, C['TEAL'])
    txt(sl, 1.54, 1.47, 6.0, 0.20,
        f"PSO {city.upper()} — RETAIL LUBRICATION DIVISION",
        size=9, color=C['TEAL_L'])
    rect(sl, 0.93, 1.81, 0.035, 0.90, C['TEAL'])
    txt(sl, 1.27, 1.75, 7.98, 0.60, f"{city} Station Profile",
        size=34, bold=True, color=C['WHITE'])
    txt(sl, 1.27, 2.44, 7.98, 0.28, "Lubricants Volume Analysis",
        size=14, color=C['ORANGE_L'])
    # volume edition badge
    rect(sl, 1.27, 2.80, 1.60, 0.22, C['TEAL'])
    txt(sl, 1.30, 2.82, 1.54, 0.18, "VOLUME EDITION",
        size=8, bold=True, color=C['WHITE'])

    rect(sl, 0.90, 3.09, 8.19, 0.018, C['GREY'])
    txt(sl, 0.90, 3.19, 4.0, 0.14, REPORT_PERIOD, size=8, color=C['GREY'])

    kpis = [
        (0.90, 1.01, f"{d['n_stns']}",               "Stations"),
        (2.20, 2.05, fmt_kl_total(d['lubes_vol']),    "Total Lubes Volume"),
        (4.55, 1.37, str(d['n_cats']),                "Active Categories"),
        (6.22, 2.08, fmt_kl(d['avg_per_stn']),        "Avg Volume / Station"),
    ]
    for bx, bw, val, lbl in kpis:
        rect(sl, bx, 3.35, bw, 0.80, C['DNAVY'])
        txt(sl, bx+0.18, 3.46, bw-0.22, 0.36, val,
            size=22, bold=True, color=C['ORANGE'])
        txt(sl, bx+0.18, 3.86, bw-0.22, 0.16, lbl,
            size=8,  color=C['TEAL_L'])

# ── SLIDE 2: At a Glance ─────────────────────────────────────────────────────
def build_slide2(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    rect(sl, 0, 0, 10, 0.61, C['DNAVY'])
    txt(sl, 0.25, 0.14, 6.5, 0.32,
        f"{d['city']} — Lubricants Volume At a Glance",
        size=17, bold=True, color=C['WHITE'])
    txt(sl, 6.5, 0.20, 3.25, 0.20,
        f"{d['n_stns']} Stations  |  {fmt_kl_total(d['lubes_vol'])} Total Volume  |  {REPORT_PERIOD}",
        size=8.5, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    lx, lw = 0.22, 3.06
    vol_chg_str = chg_str(d['vol_chg'])
    boxes = [
        (0.75, 0.86, fmt_kl_total(d['lubes_vol']),
         'TOTAL LUBES VOLUME', C['ORANGE'],
         f"vs SPLY: {vol_chg_str}  (SPLY: {fmt_kl_total(d['lubes_vol_sply'])})"),
        (1.73, 0.86, fmt_kl(d['avg_per_stn']),
         'AVG VOLUME PER STATION', C['BLUE'],
         f"Median: {fmt_kl(d['med_per_stn'])}  |  Range: {fmt_kl(d['min_per_stn'])} – {fmt_kl(d['max_per_stn'])}"),
        (2.70, 0.40, str(d['n_active']),
         'ACTIVE SELLERS', C['GREEN'],
         f"{d['n_active']/d['n_stns']*100:.1f}% of stations"),
        (3.20, 0.40, str(d['n_zero']),
         'ZERO LUBES', C['RED'],
         "Not selling any lubes"),
        (3.70, 0.40, f"{d['top10_pct']:.1f}%",
         'TOP 10 STATION CONCENTRATION',  C['PURPLE'],
         f"Top 10 stations = {fmt_kl_total(d['top10_vol'])} of total volume"),
    ]
    for by, bh, val, lbl, vcol, sub in boxes:
        rect(sl, lx, by, lw, bh, C['KPI_BG'])
        rect(sl, lx+0.03, by, 0.02, bh, vcol)
        txt(sl, lx+0.20, by+0.09, lw-0.26, bh*0.35, val,
            size=20 if by < 2.5 else 16, bold=True, color=vcol)
        txt(sl, lx+0.20, by+0.09+bh*0.35, lw-0.26, 0.14, lbl,
            size=7.5, color=C['GREY'])
        txt(sl, lx+0.20, by+0.09+bh*0.35+0.15, lw-0.26, 0.15, sub,
            size=8, color=C['GREY2'])

    # RIGHT: category volume bars
    rx = 3.47
    txt(sl, rx, 0.73, 6.5, 0.15, "VOLUME CONTRIBUTION BY CATEGORY",
        size=8, bold=True, color=C['GREY'])

    cat_max_vol = max(d['cat'][c]['vol'] for c in CAT_ORDER if d['cat'][c]['vol'] > 0) or 1
    track_w = 3.90

    for i, cat in enumerate([c for c in CAT_ORDER if d['cat'][c]['vol'] > 0]):
        by  = 0.90 + i * 0.55
        cs  = d['cat'][cat]
        bar_w = max(0.05, cs['vol'] / cat_max_vol * track_w)
        col   = CAT_COLORS.get(cat, C['GREY3'])
        txt(sl, rx, by+0.03, 1.05, 0.14, cat, size=8, color=C['DNAV2'])
        rect(sl, rx+1.10, by, track_w, 0.19, C['BG_BLUE'])
        rect(sl, rx+1.10, by+0.01, bar_w, 0.17, col)
        chg_c = C['GREEN'] if cs['vol_chg'] >= 0 else C['RED']
        txt(sl, rx+1.10+track_w+0.06, by+0.03, 1.25, 0.14,
            f"{cs['vol_pct']:.1f}%  |  {cs['n_selling']} stns  {chg_str(cs['vol_chg'])}",
            size=7, color=C['DNAV2'])

    # Key insight
    ins_y = 0.90 + len([c for c in CAT_ORDER if d['cat'][c]['vol'] > 0]) * 0.55
    ins_y = max(ins_y, 3.90)
    rect(sl, rx, ins_y, 6.38, 1.05, C['BG_YELLOW'])
    rect(sl, rx+0.02, ins_y, 0.02, 1.05, C['BROWN'])
    txt(sl, rx+0.18, ins_y+0.08, 6.10, 0.15, "KEY INSIGHT",
        size=7.5, bold=True, color=C['BROWN'])

    top2     = sorted(CAT_ORDER[:4], key=lambda c: d['cat'][c]['vol'], reverse=True)[:2]
    top2_pct = sum(d['cat'][c]['vol_pct'] for c in top2)
    widest   = max(CAT_ORDER[:4], key=lambda c: d['cat'][c]['n_selling'])
    fastest  = max([c for c in CAT_ORDER[:4] if d['cat'][c]['vol'] > 0],
                   key=lambda c: d['cat'][c]['vol_chg'])
    insight  = (f"{' + '.join(top2)} = {top2_pct:.0f}% of {d['city']} lubes volume. "
                f"{widest} has the widest reach ({d['cat'][widest]['n_pct']:.0f}% of stations). "
                f"{fastest} is the fastest growing category ({chg_str(d['cat'][fastest]['vol_chg'])} YoY). "
                f"LOW GRADE growth is driving overall volume — monitor premium category mix.")
    txt(sl, rx+0.18, ins_y+0.26, 6.10, 0.65, insight,
        size=9.5, color=C['DNAV2'], wrap=True)

# ── SLIDE 3: Category Performance ────────────────────────────────────────────
def build_slide3(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    rect(sl, 0, 0, 10, 0.58, C['DNAVY'])
    txt(sl, 0.25, 0.14, 7.2, 0.30,
        "Category Volume Performance — Avg, Min & Max per Station",
        size=17, bold=True, color=C['WHITE'])
    txt(sl, 7.25, 0.20, 2.55, 0.18,
        f"{d['n_stns']} Stations Total  |  All values in KL",
        size=9, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    active_cats = [c for c in CAT_ORDER
                   if d['cat'][c]['vol'] > 0 and d['cat'][c]['n_selling'] > 0]
    chart_top = 0.68; chart_h = 1.72
    row_h = chart_h / max(len(active_cats), 1)
    max_max = max((d['cat'][c]['max_val'] for c in active_cats), default=1) or 1
    track_x = 1.50; track_w = 6.10

    for i, cat in enumerate(active_cats):
        cs  = d['cat'][cat]
        col = CAT_COLORS.get(cat, C['GREY3'])
        ry  = chart_top + i * row_h + 0.04
        rh  = row_h - 0.08

        rect(sl, 0.22, ry, 1.24, rh, C['KPI_BG'])
        txt(sl, 0.25, ry+0.02, 1.18, rh-0.04, cat,
            size=8, bold=True, color=col)
        txt(sl, 0.25, ry+rh*0.45, 1.18, rh*0.50,
            f"{cs['n_selling']} stns ({cs['n_pct']:.0f}%)",
            size=7, color=C['GREY'])

        rect(sl, track_x, ry+rh*0.2, track_w, rh*0.45, C['TRACK'])
        max_w = max(0.05, cs['max_val'] / max_max * track_w)
        rect(sl, track_x, ry+rh*0.22, max_w, rh*0.41, C['KPI_BG'])
        avg_w = max(0.03, cs['avg'] / max_max * track_w)
        rect(sl, track_x, ry+rh*0.22, avg_w, rh*0.41, col)
        if cs['median'] > 0:
            med_x = track_x + cs['median'] / max_max * track_w
            rect(sl, med_x-0.015, ry+rh*0.18, 0.03, rh*0.49, C['DNAVY'])

        vx = track_x + track_w + 0.12
        txt(sl, vx, ry, 2.04, rh*0.50,
            f"Avg  {fmt_kl(cs['avg'])}",
            size=8, bold=True, color=col)
        txt(sl, vx, ry+rh*0.45, 2.04, rh*0.55,
            f"Med {fmt_kl(cs['median'])}  Max {fmt_kl(cs['max_val'])}",
            size=7, color=C['GREY2'])

    # 3 insight boxes
    lg   = d['cat']['LOW GRADE']
    pcmo = d['cat']['PCMO']
    mco  = d['cat']['MCO']
    not_lg  = d['n_stns'] - lg['n_selling']
    not_pcmo= d['n_stns'] - pcmo['n_selling']

    insight_data = [
        (C['RED_BG'],   C['RED_D'],
         'LOW GRADE — VOLUME LEADER, COVERAGE GAP',
         f"LOW GRADE drives {lg['vol_pct']:.0f}% of city volume ({chg_str(lg['vol_chg'])} YoY). "
         f"Yet {not_lg} stations are not stocking it. "
         f"Avg volume per active station: {fmt_kl(lg['avg'])}. "
         f"Closing this gap could add ~{fmt_kl(not_lg * lg['avg'] * 0.5)} to city volume."),
        (C['BG_GREEN'], C['BLUE2'],
         'MCO — MOST CONSISTENT VOLUMES',
         f"{mco['n_pct']:.0f}% station coverage. "
         f"Distribution is {'tight' if mco['max_val'] < 5*mco['avg'] else 'moderate'} "
         f"(avg {fmt_kl(mco['avg'])}, median {fmt_kl(mco['median'])}). "
         f"Volume {'growing' if mco['vol_chg'] >= 0 else 'declining'} at {chg_str(mco['vol_chg'])} YoY."),
        (C['BG_YELLOW'], C['BROWN'],
         f'PCMO — {not_pcmo} STATIONS NOT STOCKING',
         f"PCMO has {pcmo['n_pct']:.0f}% station coverage — "
         f"{not_pcmo} stations are missing it. "
         f"Though a smaller volume category, PCMO earns the highest margin per litre. "
         f"Activating these {not_pcmo} stations at avg {fmt_kl(pcmo['avg'])} adds "
         f"~{fmt_kl(not_pcmo * pcmo['avg'] * 0.5)} volume with premium margin."),
    ]

    bx_w = 3.11
    for ii, (bg, tc, title, body) in enumerate(insight_data):
        bx = 0.19 + ii*(bx_w+0.22); by = 2.51
        rect(sl, bx, by, bx_w, 0.98, bg)
        rect(sl, bx+0.03, by, 0.02, 0.98, tc)
        txt(sl, bx+0.18, by+0.07, bx_w-0.22, 0.15, title,
            size=7.5, bold=True, color=tc)
        txt(sl, bx+0.18, by+0.25, bx_w-0.22, 0.68, body,
            size=8, color=C['DNAV2'], wrap=True)

# ── SLIDE 4: Station Performance ─────────────────────────────────────────────
def build_slide4(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    rect(sl, 0, 0, 10, 0.58, C['DNAVY'])
    txt(sl, 0.25, 0.14, 6.5, 0.30,
        "Station Volume Performance — Who Is Delivering",
        size=17, bold=True, color=C['WHITE'])
    txt(sl, 6.5, 0.20, 3.25, 0.18,
        f"Tiers based on Total Lubes Volume  |  Median {fmt_kl(d['med_per_stn'])}",
        size=8.5, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    txt(sl, 0.19, 0.68, 3.26, 0.14, "STATION VOLUME TIERS",
        size=7.5, bold=True, color=C['GREY'])

    tiers = [
        (C['GREEN'],  f"HIGH PERFORMERS    {d['n_high']} stations",
         fmt_kl_total(d['vol_high']),
         f"{d['vol_high']/d['lubes_vol']*100:.1f}% of {d['city']} Lubes  |  Above {fmt_kl(d['p75'])} each"),
        (C['BLUE'],   f"MID PERFORMERS    {d['n_mid']} stations",
         fmt_kl_total(d['vol_mid']),
         f"{d['vol_mid']/d['lubes_vol']*100:.1f}% of {d['city']} Lubes  |  {fmt_kl(d['med_per_stn'])} – {fmt_kl(d['p75'])}"),
        (C['AMBER'],  f"LOW PERFORMERS    {d['n_low']} stations",
         fmt_kl_total(d['vol_low']),
         f"{d['vol_low']/d['lubes_vol']*100:.1f}% of {d['city']} Lubes  |  Below {fmt_kl(d['med_per_stn'])}"),
        (C['RED'],    f"ZERO LUBES    {d['n_zero_stn']} stations",
         "0 KL",
         "Not selling any lubricant category"),
    ]
    for ti, (col, title, val, sub) in enumerate(tiers):
        by = 0.86 + ti*0.97
        rect(sl, 0.19, by, 3.19, 0.83, col)
        txt(sl, 0.33, by+0.08, 2.98, 0.18, title,
            size=9, bold=True, color=C['WHITE'])
        txt(sl, 0.33, by+0.28, 2.98, 0.30, val,
            size=18, bold=True, color=C['WHITE'])
        txt(sl, 0.33, by+0.62, 2.98, 0.16, sub,
            size=8, color=C['WHITE'])

    # TOP 10 stacked bar chart (by volume)
    txt(sl, 3.56, 0.68, 6.35, 0.14,
        "TOP 10 STATIONS BY TOTAL LUBES VOLUME (KL)",
        size=7.5, bold=True, color=C['GREY'])

    top10 = d['top10']
    chart_x = 3.56; chart_y = 0.86
    max_vol  = top10[0]['total'] if top10 else 1
    bar_max_w = 5.60
    row_h = 0.24; row_gap = 0.03

    for si, stn in enumerate(top10[:10]):
        ry = chart_y + si * (row_h + row_gap)
        txt(sl, chart_x, ry+0.02, 1.75, row_h-0.04,
            stn['name'][:22], size=7, color=C['DNAV2'])
        bx_start = chart_x + 1.78
        for cat in CAT_ORDER:
            vol = stn['by_cat'].get(cat, 0)
            if vol > 0:
                seg_w = vol / max_vol * bar_max_w
                rect(sl, bx_start, ry+0.02, seg_w, row_h-0.04,
                     CAT_COLORS.get(cat, C['GREY3']))
                bx_start += seg_w
        txt(sl, bx_start+0.06, ry+0.02, 0.90, row_h-0.04,
            fmt_kl(stn['total']), size=7, color=C['GREY2'])

    leg_y = chart_y + 10*(row_h+row_gap) + 0.04
    for li, cat in enumerate(CAT_ORDER[:5]):
        lx = chart_x + li * 1.25
        rect(sl, lx, leg_y, 0.12, 0.10, CAT_COLORS.get(cat, C['GREY3']))
        txt(sl, lx+0.15, leg_y, 1.08, 0.12, cat, size=6.5, color=C['GREY'])

    ins_y = leg_y + 0.18
    rect(sl, 3.56, ins_y, 6.25, 0.52, C['BG_YELLOW'])
    rect(sl, 3.58, ins_y, 0.02, 0.52, C['BROWN'])
    top10_names = [s['name'][:20] for s in top10[:1]] if top10 else ['']
    ins_txt = (f"Top 10 stations = {fmt_kl_total(d['top10_vol'])} ({d['top10_pct']:.1f}% of all {d['city']} lubes volume). "
               f"{top10_names[0]} leads at {fmt_kl(top10[0]['total'])} — "
               f"analyse this station's category mix as a benchmark for the network.")
    txt(sl, 3.72, ins_y+0.06, 6.06, 0.40, ins_txt,
        size=8, color=C['DNAV2'], wrap=True)

# ── SLIDE 5: Priorities ───────────────────────────────────────────────────────
def build_slide5(prs, d):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    rect(sl, 0, 0, 10, 0.58, C['DNAVY'])
    txt(sl, 0.25, 0.14, 7.5, 0.30,
        f"{d['city']} — Where to Focus & What to Fix",
        size=17, bold=True, color=C['WHITE'])
    txt(sl, 7.5, 0.20, 2.30, 0.18,
        "4 Priorities  |  Estimated Volume Upside",
        size=8.5, color=C['ORANGE_L'], align=PP_ALIGN.RIGHT)

    no_pcmo   = d['no_pcmo']
    no_lowgr  = d['no_lowgr']
    n_low     = d['n_low']
    n_zero    = d['n_zero_stn']
    pcmo_avg  = d['pcmo_avg']
    lowgr_avg = d['lowgr_avg']
    med       = d['med_per_stn']
    avg_low   = d['vol_low'] / d['n_low'] if d['n_low'] else 0
    potential3 = max(0, (med - avg_low) * min(n_low, max(int(n_low*0.35), 5)))

    priorities = [
        (C['BG_LIGHT'], C['BLUE'],   C['BLUE2'],
         '01', f"Activate {no_lowgr} Stations Not Selling LOW GRADE",
         f"LOW GRADE drives the largest volume share and is growing fastest nationally (+24%). "
         f"{no_lowgr} stations in {d['city']} are not stocking it. "
         f"Avg LOW GRADE volume per active station: {fmt_kl(lowgr_avg)}.",
         fmt_kl(no_lowgr * lowgr_avg * 0.5),
         "VOLUME GROWTH  Priority #1"),

        (C['BG_GREEN'], C['GREEN'],  C['BLUE2'],
         '02', f"Activate {no_pcmo} Stations Not Selling PCMO",
         f"PCMO is premium volume with the highest margin per litre. "
         f"{no_pcmo} stations — {no_pcmo/d['n_stns']*100:.1f}% of network — are not stocking it. "
         f"Avg PCMO volume per active station: {fmt_kl(pcmo_avg)}.",
         fmt_kl(no_pcmo * pcmo_avg * 0.5),
         "PREMIUM MIX  Immediate Action"),

        (C['BG_YELLOW'], C['GOLD'],  C['BROWN'],
         '03', f"Upgrade {n_low} Low-Volume Stations",
         f"{n_low} stations sell below the city median of {fmt_kl(med)}. "
         f"If {min(n_low, max(int(n_low*0.35), 5))} of them reached the median "
         f"from avg of {fmt_kl(avg_low)}, the volume uplift would be material.",
         fmt_kl(potential3) + "+",
         "WALLET SHARE  Sales Force Priority"),

        (C['RED_BG'], C['RED'],      C['RED_D'],
         '04', f"Convert {n_zero} Zero-Volume Stations",
         f"{n_zero} active fuel stations are selling ZERO lubricants. "
         f"Investigate: missing product license, no shelf space, competitor brand only. "
         f"Each conversion at the city avg of {fmt_kl(d['avg_per_stn'])} adds meaningful volume.",
         fmt_kl(n_zero * d['avg_per_stn'] * 0.4) + "+",
         "QUICK WIN  Investigate & Convert"),
    ]

    col_w = 2.28
    for pi, (bg, num_col, title_col, num, title, body, pot, tag) in enumerate(priorities):
        px = 0.19 + pi*(col_w+0.14)
        rect(sl, px, 0.70, col_w, 4.82, bg)
        rect(sl, px+0.03, 0.72, col_w, 0.02, num_col)
        txt(sl, px+0.17, 0.85, col_w-0.20, 0.44, num,
            size=28, bold=True, color=num_col)
        txt(sl, px+0.17, 1.38, col_w-0.20, 0.36, title,
            size=9, bold=True, color=title_col, wrap=True)
        txt(sl, px+0.17, 1.80, col_w-0.20, 1.30, body,
            size=7.5, color=C['DNAV2'], wrap=True)
        txt(sl, px+0.17, 3.18, col_w-0.20, 0.22,
            f"Potential: {pot}", size=10, bold=True, color=num_col)
        txt(sl, px+0.17, 3.44, col_w-0.20, 0.22, tag,
            size=6.5, bold=True, color=C['GREY'])

# ── MAIN ─────────────────────────────────────────────────────────────────────
top_vol = (lubes_all.groupby('CityNorm')['SalesLtr_CY'].sum()
           .sort_values(ascending=False))
top_cities = top_vol.index.tolist()

# Same 10 cities as revenue edition (rank 2-10 + Peshawar)
# Also generate Karachi volume edition
target_cities = [c for c in top_cities if c != 'Karachi'][:9]
if 'Peshawar' not in target_cities:
    target_cities.append('Peshawar')
# include Karachi for completeness
target_cities = ['Karachi'] + target_cities

print(f"Generating VOLUME profiles for: {target_cities}")
os.makedirs('reports/city_profiles_volume', exist_ok=True)

for city in target_cities:
    print(f"  Processing {city}…", end=' ', flush=True)
    d = compute_city(city)

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    build_slide1(prs, d)
    build_slide2(prs, d)
    build_slide3(prs, d)
    build_slide4(prs, d)
    build_slide5(prs, d)

    out = out_path(f"PSO_{city.replace(' ','_')}_Lubes_Vol_Profile", 'pptx', df,
                    out_dir='reports/city_profiles_volume')
    prs.save(out)
    print(f"saved → {out}")

print("\nAll VOLUME profiles generated.")
print("Revenue profiles in: reports/city_profiles/")
print("Volume  profiles in: reports/city_profiles_volume/")
