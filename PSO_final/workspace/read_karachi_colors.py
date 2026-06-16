"""Read fill colors and font sizes from karachi_profile.pptx shapes."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.dml.color import RGBColor
from lxml import etree

prs = Presentation('karachi_profile.pptx')

for si, slide in enumerate(prs.slides):
    print(f"\n{'='*60} SLIDE {si+1}")
    for shape in slide.shapes:
        # fill color
        fill_hex = None
        try:
            fill = shape.fill
            if fill.type is not None:
                fg = fill.fore_color
                if fg.type is not None:
                    fill_hex = str(fg.rgb)
        except: pass

        # font info from first text run
        font_sz, font_bold, font_hex, font_name = None, None, None, None
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_sz = run.font.size.pt
                        font_bold = run.font.bold
                        try:
                            font_hex = str(run.font.color.rgb)
                        except: pass
                        font_name = run.font.name
                        break
                    if font_sz: break
        except: pass

        txt_preview = ''
        if hasattr(shape,'text') and shape.text.strip():
            txt_preview = shape.text.strip()[:50]

        if fill_hex or txt_preview or font_sz:
            print(f"  '{shape.name}' | fill={fill_hex} | font={font_sz}pt bold={font_bold} color={font_hex} face={font_name}")
            if txt_preview:
                print(f"    text: {txt_preview}")
